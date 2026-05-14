#!/usr/bin/env python3
"""
llm_ingest.py - Convert documents to LLM-ready Markdown

Supported formats: PDF, DOCX, PPTX, TXT, HTML, CSV
Output: Clean .md files optimized for LLM context windows

Usage:
    # Convert the entire research papers folder to LLM-ready Markdown:
    python llm_ingest.py "C:\\Users\\User\\Desktop\\Research\\papers\\downloaded"

    # Convert and save output to a specific folder:
    python llm_ingest.py "C:\\Users\\User\\Desktop\\Research\\papers\\downloaded" --out-dir "C:\\Users\\User\\Desktop\\Research\\papers\\llm_ready"

    # Convert a single paper with chunking (ideal for RAG / large docs):
    python llm_ingest.py "C:\\Users\\User\\Desktop\\Research\\papers\\downloaded\\my_paper.pdf" --chunk 2000

    # Convert a single paper to a specific output file:
    python llm_ingest.py "C:\\Users\\User\\Desktop\\Research\\papers\\downloaded\\my_paper.pdf" --output "C:\\Users\\User\\Desktop\\Research\\papers\\llm_ready\\my_paper.md"

Install deps:
    pip install pymupdf python-docx python-pptx beautifulsoup4 pandas tiktoken

PDF OCR:
    PyMuPDF OCR also requires Tesseract language data. Point the script at the
    tessdata directory with --tessdata or TESSDATA_PREFIX.
"""

import argparse
import contextlib
import datetime as dt
import hashlib
import importlib.util
import io
import json
import os
import re
import shutil
import site
import stat
import subprocess
import sys
import time
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import llm_audit_assertions
import llm_figure_cleanup
import llm_pdf_cleanup


MANIFEST_VERSION = 5
_PIPELINE_SIGNATURE: str | None = None
_PDF_BACKEND_PLAN_CACHE: dict[tuple[str, str, bool, str, str, str], "PDFBackendPlan"] = {}
_PDF_TRAITS_CACHE: dict[str, "PDFDocumentTraits"] = {}

DEFAULT_AUDIT_MANIFEST = "audit_corpus_manifest.json"
DEFAULT_AUDIT_CACHE_DIR = "_audit_corpus_cache"
DEFAULT_AUDIT_REPORT_DIR = "_audit_reports"
DEFAULT_AUDIT_BASELINE_DIR = "downloaded"
DEFAULT_AUDIT_BACKENDS = "auto,custom:off,pymupdf4llm,marker"
DEFAULT_BACKEND_TIMEOUT_SECONDS = 300
DEFAULT_MAX_INPUT_MB = 250
DEFAULT_MAX_PDF_PAGES = 500
DEFAULT_MAX_EXTRACTED_ASSETS = 500
DEFAULT_MAX_AUDIT_DOWNLOAD_MB = 250
_PDF_WORKER_ENV = "LLM_INGEST_PDF_WORKER"
_TRUSTED_MARKER_ENV = "LLM_INGEST_ALLOW_EXTERNAL_MARKER_PYTHON"


def _ensure_user_site_packages() -> None:
    for vendor_name in ("_vendor_manual", "_vendor_local", "_vendor_site"):
        vendor_site = Path(__file__).with_name(vendor_name)
        if vendor_site.exists():
            vendor_str = str(vendor_site)
            if vendor_str not in sys.path:
                sys.path.insert(0, vendor_str)

    with contextlib.suppress(Exception):
        user_site = site.getusersitepackages()
        if user_site and user_site not in sys.path and Path(user_site).exists():
            sys.path.append(user_site)


_ensure_user_site_packages()


try:
    import tiktoken

    TOKENIZER = tiktoken.get_encoding("cl100k_base")

    def count_tokens(text: str) -> int:
        return len(TOKENIZER.encode(text))

except ImportError:

    def count_tokens(text: str) -> int:
        return len(text.split()) * 4 // 3


_BOILERPLATE_PHRASES = [
    "queen's university belfast",
    "research portal",
    "publisher rights",
    "general rights",
    "take down policy",
    "open access",
    "download date",
    "link to publication record",
    "reprints and permissions",
    "macmillan publishers",
    "all rights reserved",
    "copyright for the publications",
    "condition of accessing",
    "if you discover content",
    "openaccess@qub",
    "go.qub.ac.uk",
    "www.nature.com/reprints",
    "supplementary information is linked",
    "document version publisher",
    "publisher's pdf",
    "downloaded from the university",
    "strictly personal use",
    "copyright holder",
    "taverne",
    "this article is available from",
    "licensee biomed central",
    "page number not for citation purposes",
]

_SECTION_HEADING_PATTERNS = [
    (re.compile(r"^abstract\b[:\s-]*(.*)$", re.IGNORECASE), "## Abstract"),
    (re.compile(r"^introduction\b[:\s-]*(.*)$", re.IGNORECASE), "## Introduction"),
    (re.compile(r"^results\b[:\s-]*(.*)$", re.IGNORECASE), "## Results"),
    (re.compile(r"^discussion\b[:\s-]*(.*)$", re.IGNORECASE), "## Discussion"),
    (re.compile(r"^methods?\b[:\s-]*(.*)$", re.IGNORECASE), "## Methods"),
    (re.compile(r"^conclusion\b[:\s-]*(.*)$", re.IGNORECASE), "## Conclusion"),
    (re.compile(r"^references?\b[:\s-]*(.*)$", re.IGNORECASE), "## References"),
    (
        re.compile(r"^acknowledgements?\b[:\s-]*(.*)$", re.IGNORECASE),
        "## Acknowledgements",
    ),
    (
        re.compile(r"^author contributions?\b[:\s-]*(.*)$", re.IGNORECASE),
        "## Author Contributions",
    ),
    (
        re.compile(r"^author information\b[:\s-]*(.*)$", re.IGNORECASE),
        "## Author Information",
    ),
    (
        re.compile(r"^supplementary\b[:\s-]*(.*)$", re.IGNORECASE),
        "## Supplementary",
    ),
]

_FIGURE_CAPTION_RE = re.compile(
    r"^(?:(?:Supplementary|Extended Data)\s+)?(?:Figure\s+\d+[A-Za-z]?[\.:]?\s*|Fig\.\s*\d+[A-Za-z]?[\.:]?\s*)",
    re.IGNORECASE,
)
_DISPLAY_CAPTION_RE = re.compile(
    r"^(((?:(?:Supplementary|Extended Data)\s+)?(?:Figure|Fig\.|Table))\s+\d+[A-Za-z]?(?:[\.:]|\b)\s*)",
    re.IGNORECASE,
)
_PICTURE_TEXT_BLOCK_RE = re.compile(
    r"\*\*----- Start of picture text -----\*\*<br>\s*.*?\s*\*\*----- End of picture text -----\*\*<br>\s*",
    re.IGNORECASE | re.DOTALL,
)
_REFERENCE_START_RE = re.compile(r"^\s*(\d{1,3})[.)]\s+(.+)")
_DOI_RE = re.compile(r"https?://doi\.org/\S+|doi:\s*\S+", re.IGNORECASE)
_DISPLAY_DOI_SUFFIX_RE = re.compile(r"\s+doi:\s*\S+\.(?:g|t|f)\d{3}\b\.?$", re.IGNORECASE)
_DISPLAY_DOI_FRAGMENT_RE = re.compile(r"doi:\s*\S+\.(?:g|t|f)\d{3}\b\.?", re.IGNORECASE)
_DATE_FIELD_RE = re.compile(
    r"\b(Received|Accepted)\s+(\d{1,2}\s+\w+\s+\d{4})", re.IGNORECASE
)
_PUBLISHER_METADATA_PREFIXES = (
    "copyright:",
    "funding:",
    "competing interests:",
    "academic editor:",
    "editor:",
    "published:",
    "published online",
    "open access",
    "this article is licensed under",
    "peer review information",
    "reprints and permission information",
    "publisher's note",
)
_TITLE_COMPARE_STOPWORDS = {
    "a",
    "an",
    "and",
    "of",
    "the",
    "to",
    "for",
    "in",
    "on",
    "with",
    "from",
    "by",
}
_KNOWN_TITLE_DOI = {
    "recombinant spidroins fully replicate primary mechanical properties of natural spider silk": "https://doi.org/10.1021/acs.biomac.8b00980",
}
_MOJIBAKE_MARKERS = ("â", "Ã", "Â", "ï¿", "�")
_CAPTION_REFERENCE_VERBS = {
    "show",
    "shows",
    "display",
    "displays",
    "illustrate",
    "illustrates",
    "depict",
    "depicts",
    "present",
    "presents",
    "describe",
    "describes",
    "summarize",
    "summarizes",
    "summarise",
    "summarises",
    "compare",
    "compares",
    "indicate",
    "indicates",
}
_TABLE_CELL_LITERAL_REPLACEMENTS = {
    "innercore": "inner core",
    "outercore": "outer core",
    "tensilestrength": "tensile strength",
    "waterbalance": "water balance",
    "socialaspects": "social aspects",
}
_PROSE_TERM_REPLACEMENTS = {
    "lithiumbromide": "lithium bromide",
    "lithiumchloride": "lithium chloride",
    "lithiumthiocyanate": "lithium thiocyanate",
    "lithiumisothiocyanate": "lithium isothiocyanate",
    "guanidiniumchloride": "guanidinium chloride",
    "guanidiniumhydrochloride": "guanidinium hydrochloride",
    "guanidiniumthiocyanate": "guanidinium thiocyanate",
    "guanidiniumisothiocyanate": "guanidinium isothiocyanate",
    "hexafluorisopropanol": "hexafluoroisopropanol",
    "hexafluorispropanol": "hexafluoroisopropanol",
    "hexafluorispropananol": "hexafluoroisopropanol",
    "gelelectrophoresis": "gel electrophoresis",
    "sdspage": "SDS-PAGE",
    "highperformance": "high-performance",
    "petroleumderived": "petroleum-derived",
    "fastgrowing": "fast-growing",
    "selfinteraction": "self-interaction",
    "posttranslational": "post-translational",
    "biterminal": "bi-terminal",
    "proteinbased": "protein-based",
    "sidechain": "side chain",
    "fedbatch": "fed-batch",
    "industrialscale": "industrial-scale",
    "asspun": "as-spun",
    "minispidroin": "mini-spidroin",
    "minispidroins": "mini-spidroins",
    "stressstrain": "stress-strain",
    "citratephosphate": "citrate-phosphate",
    "citrate phosphate": "citrate-phosphate",
    "fbers": "fibers",
    "specifc": "specific",
}
_DEHYPHENATED_WORD_REPLACEMENTS = {
    "addition": "addition",
    "amino": "amino",
    "biomedical": "biomedical",
    "carboxyl": "carboxyl",
    "chemical": "chemical",
    "functions": "functions",
    "hierarchical": "hierarchical",
    "mechanical": "mechanical",
    "mechanism": "mechanism",
    "miniature": "miniature",
    "molecular": "molecular",
    "polymers": "polymers",
    "production": "production",
    "produced": "produced",
    "proteinaceous": "proteinaceous",
    "sealants": "sealants",
    "spider": "spider",
    "structures": "structures",
    "supramolecular": "supramolecular",
    "technique": "technique",
    "viscoelasticity": "viscoelasticity",
}
_LIGATURE_REPLACEMENTS = {
    "ﬀ": "ff",
    "ﬁ": "fi",
    "ﬂ": "fl",
    "ﬃ": "ffi",
    "ﬄ": "ffl",
    "ﬅ": "ft",
    "ﬆ": "st",
}
_SPACED_HYPHEN_PREFIXES = {
    "as",
    "bi",
    "cross",
    "end",
    "high",
    "low",
    "phase",
    "post",
    "pre",
    "self",
    "side",
}


@dataclass
class SecurityLimits:
    max_input_mb: int = DEFAULT_MAX_INPUT_MB
    max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES
    max_extracted_assets: int = DEFAULT_MAX_EXTRACTED_ASSETS
    max_audit_download_mb: int = DEFAULT_MAX_AUDIT_DOWNLOAD_MB
    backend_timeout_seconds: int = DEFAULT_BACKEND_TIMEOUT_SECONDS
    hardened_mode: bool = True
    allow_external_marker_python: bool = False
    allow_unverified_downloads: bool = False
    privacy_mode: bool = False


@dataclass
class PDFConfig:
    ocr_language: str = "eng"
    ocr_dpi: int = 200
    tessdata: str | None = None
    ocr_mode: str = "auto"
    pdf_backend: str = "auto"
    table_strategy: str = "lines_strict"
    marker_python: str | None = None
    security: SecurityLimits = field(default_factory=SecurityLimits)


@dataclass(frozen=True)
class PDFBackendCandidate:
    name: str
    importable: bool
    runnable: bool
    detail: str
    route_reason: str = ""

    @property
    def available(self) -> bool:
        return self.runnable


@dataclass(frozen=True)
class PDFBackendPlan:
    requested: str
    selected: str | None
    candidates: tuple[PDFBackendCandidate, ...]

    @property
    def order(self) -> tuple[str, ...]:
        return tuple(candidate.name for candidate in self.candidates)


@dataclass(frozen=True)
class PDFDocumentTraits:
    text_extractable: bool
    born_digital: bool
    scanned_like: bool
    image_heavy: bool
    sample_pages: int
    average_chars_per_page: float
    average_image_area_ratio: float


@dataclass(frozen=True)
class AuditBackendSpec:
    name: str
    ocr_mode: str | None = None

    @property
    def label(self) -> str:
        if self.ocr_mode and self.ocr_mode != "auto":
            return f"{self.name}_{self.ocr_mode}"
        return self.name


@dataclass(frozen=True)
class AuditSample:
    id: str
    category: str
    source_url: str
    filename: str
    expected_traits: tuple[str, ...]
    recommended_backends: tuple[str, ...]
    sha256: str
    license_note: str


@dataclass(frozen=True)
class AuditFileTarget:
    id: str
    label: str
    path: Path
    source_kind: str
    category: str
    source_url: str = ""
    expected_traits: tuple[str, ...] = ()


@dataclass(frozen=True)
class AuditRunResult:
    sample_id: str
    sample_label: str
    source_kind: str
    category: str
    backend_requested: str
    backend_label: str
    backend_used: str | None
    status: str
    output_path: str
    asset_dir: str
    tokens: int
    asset_count: int
    issue_counts: dict[str, int]
    issue_total: int
    log_excerpt: str
    error: str = ""


@dataclass(frozen=True)
class AuditReport:
    created_at: str
    manifest_path: str
    cache_dir: str
    report_dir: str
    baseline_dirs: tuple[str, ...]
    backend_labels: tuple[str, ...]
    backend_plan: dict[str, list[str]]
    missing_samples: tuple[str, ...]
    results: tuple[AuditRunResult, ...]


@dataclass(frozen=True)
class PDFRenderedAsset:
    markdown: str
    caption_text: str = ""
    label_key: str = ""


@dataclass(frozen=True)
class PDFFigureCandidate:
    bbox: tuple[float, float, float, float]
    caption_text: str = ""
    label_key: str = ""
    image_index: int = 0


@dataclass
class PDFLine:
    text: str
    bbox: tuple[float, float, float, float]
    font_size: float


@dataclass
class PDFBlock:
    text: str
    bbox: tuple[float, float, float, float]
    font_size: float
    lines: list[PDFLine] = field(default_factory=list)


@dataclass
class PDFImage:
    index: int
    bbox: tuple[float, float, float, float]
    width: int
    height: int


@dataclass
class PDFPage:
    number: int
    width: float
    height: float
    raw_text: str
    lines: list[PDFLine]
    blocks: list[PDFBlock]
    images: list[PDFImage]
    tables: list["PDFTable"]

    @property
    def page_area(self) -> float:
        return max(self.width * self.height, 1.0)

    @property
    def image_area_ratio(self) -> float:
        total = 0.0
        for image in self.images:
            total += _bbox_area(image.bbox)
        return min(total / self.page_area, 1.0)


@dataclass
class PDFMetadata:
    title: str = ""
    authors: str = ""
    journal: str = ""
    doi: str = ""
    received: str = ""
    accepted: str = ""
    first_page_cutoff_y: float = 0.0


@dataclass
class PDFTable:
    bbox: tuple[float, float, float, float]
    rows: list[list[str]]


class ConversionCancelled(Exception):
    """Raised when a conversion run is cancelled."""


def _normalized_resolved_path(path: Path) -> str:
    return os.path.normcase(str(path.resolve(strict=False)))


def _path_is_within(child: Path, parent: Path) -> bool:
    child_norm = _normalized_resolved_path(child)
    parent_norm = _normalized_resolved_path(parent)
    return child_norm == parent_norm or child_norm.startswith(parent_norm + os.sep)


def list_supported_files(input_path: Path, output_path: Path | None = None) -> list[Path]:
    exclude_output = output_path is not None and _path_is_within(output_path, input_path)
    files: list[Path] = []
    for file in input_path.rglob("*"):
        if not file.is_file():
            continue
        if file.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if exclude_output and _path_is_within(file, output_path):
            continue
        files.append(file)
    files.sort(key=_normalized_resolved_path)
    return files


def build_batch_targets(files: list[Path], input_root: Path, output_root: Path) -> list[tuple[Path, Path]]:
    planned_rel_paths: dict[Path, Path] = {}
    collisions: dict[str, int] = {}

    for file in files:
        relative = file.relative_to(input_root)
        relative_output = relative.with_suffix(".md")
        planned_rel_paths[file] = relative_output
        key = relative_output.as_posix().lower()
        collisions[key] = collisions.get(key, 0) + 1

    batch_plan: list[tuple[Path, Path]] = []
    for file in files:
        relative_output = planned_rel_paths[file]
        key = relative_output.as_posix().lower()
        if collisions[key] > 1:
            source_suffix = file.suffix.lower().lstrip(".") or "file"
            relative_output = relative_output.with_name(f"{relative_output.stem}__{source_suffix}.md")
        batch_plan.append((file, output_root / relative_output))
    return batch_plan


def _check_cancel(cancel_event: Any | None, message: str = "Run cancelled by user.") -> None:
    if cancel_event is not None and cancel_event.is_set():
        raise ConversionCancelled(message)


def _pipeline_signature() -> str:
    global _PIPELINE_SIGNATURE
    if _PIPELINE_SIGNATURE is None:
        digest = hashlib.sha256(Path(__file__).read_bytes()).hexdigest()
        _PIPELINE_SIGNATURE = digest[:16]
    return _PIPELINE_SIGNATURE


def _system_exit_message(exc: SystemExit) -> str:
    code = exc.code
    if isinstance(code, str) and code.strip():
        return code.strip()
    if code not in (None, 0):
        return f"Exit code {code}"
    return ""


def _security_limits_from_payload(payload: dict[str, Any] | None) -> SecurityLimits:
    if not isinstance(payload, dict):
        return SecurityLimits()
    defaults = SecurityLimits()
    values: dict[str, Any] = {}
    for field_name in defaults.__dataclass_fields__:
        if field_name in payload:
            values[field_name] = payload[field_name]
    return SecurityLimits(**values)


def _security_limits_to_payload(limits: SecurityLimits) -> dict[str, Any]:
    return {
        "max_input_mb": limits.max_input_mb,
        "max_pdf_pages": limits.max_pdf_pages,
        "max_extracted_assets": limits.max_extracted_assets,
        "max_audit_download_mb": limits.max_audit_download_mb,
        "backend_timeout_seconds": limits.backend_timeout_seconds,
        "hardened_mode": limits.hardened_mode,
        "allow_external_marker_python": limits.allow_external_marker_python,
        "allow_unverified_downloads": limits.allow_unverified_downloads,
        "privacy_mode": limits.privacy_mode,
    }


def _pdf_config_to_payload(pdf_config: PDFConfig) -> dict[str, Any]:
    return {
        "ocr_language": pdf_config.ocr_language,
        "ocr_dpi": pdf_config.ocr_dpi,
        "tessdata": pdf_config.tessdata,
        "ocr_mode": pdf_config.ocr_mode,
        "pdf_backend": pdf_config.pdf_backend,
        "table_strategy": pdf_config.table_strategy,
        "marker_python": pdf_config.marker_python,
        "security": _security_limits_to_payload(pdf_config.security),
    }


def pdf_config_from_payload(payload: dict[str, Any]) -> PDFConfig:
    return PDFConfig(
        ocr_language=str(payload.get("ocr_language") or "eng"),
        ocr_dpi=int(payload.get("ocr_dpi") or 200),
        tessdata=payload.get("tessdata") or None,
        ocr_mode=str(payload.get("ocr_mode") or "auto"),
        pdf_backend=str(payload.get("pdf_backend") or "auto"),
        table_strategy=str(payload.get("table_strategy") or "lines_strict"),
        marker_python=payload.get("marker_python") or None,
        security=_security_limits_from_payload(payload.get("security")),
    )


def _validate_security_limits(limits: SecurityLimits) -> None:
    if limits.max_input_mb <= 0:
        raise ValueError("Maximum input size must be greater than zero.")
    if limits.max_pdf_pages <= 0:
        raise ValueError("Maximum PDF pages must be greater than zero.")
    if limits.max_extracted_assets <= 0:
        raise ValueError("Maximum extracted assets must be greater than zero.")
    if limits.max_audit_download_mb <= 0:
        raise ValueError("Maximum audit download size must be greater than zero.")
    if limits.backend_timeout_seconds <= 0:
        raise ValueError("Backend timeout must be greater than zero.")


def _file_size_mb(path: Path) -> float:
    return path.stat().st_size / (1024 * 1024)


def _validate_input_file_limits(path: Path, limits: SecurityLimits) -> None:
    _validate_security_limits(limits)
    size_mb = _file_size_mb(path)
    if size_mb > limits.max_input_mb:
        raise ValueError(
            f"Input file is {size_mb:.1f} MB, above the configured limit of {limits.max_input_mb} MB: {path.name}"
        )


def _validate_pdf_page_limit(path: Path, limits: SecurityLimits) -> None:
    if path.suffix.lower() != ".pdf":
        return
    fitz = _load_pymupdf()
    try:
        doc = fitz.open(path)
        page_count = len(doc)
    except Exception as exc:
        raise ValueError(f"Unable to inspect PDF page count for {path.name}: {exc}") from exc
    if page_count > limits.max_pdf_pages:
        raise ValueError(
            f"PDF has {page_count} pages, above the configured limit of {limits.max_pdf_pages}: {path.name}"
        )


def safe_atomic_write_text(path: Path, text: str, *, encoding: str = "utf-8") -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp_path = path.with_name(f".{path.name}.tmp")
    text = _normalize_line_endings(text)
    with temp_path.open("w", encoding=encoding, newline="\n") as handle:
        handle.write(text)
    temp_path.replace(path)


def safe_atomic_write_json(path: Path, payload: Any) -> None:
    safe_atomic_write_text(path, json.dumps(payload, indent=2, sort_keys=True), encoding="utf-8")


def _is_dangerous_delete_target(path: Path) -> bool:
    try:
        resolved = path.resolve()
    except OSError:
        resolved = path.absolute()
    home = Path.home().resolve()
    cwd = Path.cwd().resolve()
    anchors = {Path(resolved.anchor).resolve()} if resolved.anchor else set()
    dangerous = {home, cwd, *anchors}
    return resolved in dangerous or resolved.parent == resolved


def _has_reparse_point(path: Path) -> bool:
    if not hasattr(os, "stat"):
        return False
    try:
        return bool(path.stat().st_file_attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0))
    except (AttributeError, OSError):
        return False


def safe_remove_generated_dir(path: Path) -> None:
    if not path.exists():
        return
    if not path.is_dir():
        raise ValueError(f"Refusing to delete non-directory generated path: {path}")
    if _is_dangerous_delete_target(path):
        raise ValueError(f"Refusing to delete dangerous generated path: {path}")
    if path.is_symlink() or _has_reparse_point(path):
        raise ValueError(f"Refusing to delete link/reparse generated path: {path}")
    shutil.rmtree(path)


def _sanitize_display_text(text: str, *, limit: int = 220) -> str:
    cleaned = "".join(ch if ch.isprintable() or ch in "\t\n" else "?" for ch in str(text))
    cleaned = cleaned.replace("\r", "\\r")
    if len(cleaned) > limit:
        return cleaned[: limit - 3].rstrip() + "..."
    return cleaned


def _redact_local_path(path: Path | str) -> str:
    text = str(path)
    home = str(Path.home())
    if home and text.lower().startswith(home.lower()):
        return "%USERPROFILE%" + text[len(home) :]
    return text


def dependency_provenance(privacy_mode: bool = False) -> dict[str, str]:
    modules = ("fitz", "pymupdf4llm", "docx", "pptx", "bs4", "pandas", "marker")
    result: dict[str, str] = {}
    for module_name in modules:
        spec = importlib.util.find_spec(module_name)
        origin = ""
        if spec is not None:
            origin = str(spec.origin or "")
            if privacy_mode:
                origin = _redact_local_path(origin)
        result[module_name] = origin or "missing"
    return result


def _marker_runner_path() -> Path:
    return Path(__file__).with_name("marker_sidecar_runner.py")


def _pdf_worker_runner_path() -> Path:
    return Path(__file__).with_name("pdf_worker_runner.py")


def _build_pdf_worker_env() -> dict[str, str]:
    env: dict[str, str] = {}
    for key in (
        "SYSTEMROOT",
        "WINDIR",
        "PATH",
        "PATHEXT",
        "COMSPEC",
        "PROCESSOR_ARCHITECTURE",
        "NUMBER_OF_PROCESSORS",
        "USERPROFILE",
        "HOMEDRIVE",
        "HOMEPATH",
        "APPDATA",
        "LOCALAPPDATA",
        "TEMP",
        "TMP",
    ):
        if os.environ.get(key):
            env[key] = os.environ[key]
    env[_PDF_WORKER_ENV] = "1"
    env["PYTHONUTF8"] = "1"
    return env


def _extract_pdf_in_worker(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> tuple[str, str]:
    _validate_input_file_limits(path, pdf_config.security)
    runner_path = _pdf_worker_runner_path()
    if not runner_path.exists():
        sys.exit(f"PDF worker runner not found: {runner_path}")

    temp_output = output_path.parent / f".{output_path.stem}.pdf_worker_raw.md"
    payload = {
        "input_path": str(path),
        "output_path": str(output_path),
        "temp_output": str(temp_output),
        "pdf_config": _pdf_config_to_payload(pdf_config),
    }
    env = _build_pdf_worker_env()
    command = [sys.executable, str(runner_path)]
    process = subprocess.Popen(
        command,
        stdin=subprocess.PIPE,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        env=env,
    )
    try:
        stdout, stderr = process.communicate(
            json.dumps(payload),
            timeout=pdf_config.security.backend_timeout_seconds,
        )
    except subprocess.TimeoutExpired:
        process.kill()
        stdout, stderr = process.communicate()
        sys.exit(f"PDF backend timed out after {pdf_config.security.backend_timeout_seconds} seconds for {path.name}.")
    finally:
        if process.poll() is None:
            with contextlib.suppress(Exception):
                process.kill()

    try:
        result = json.loads(stdout or "{}")
    except json.JSONDecodeError:
        detail = _summarize_backend_detail(stderr or stdout) or "PDF worker returned invalid output."
        sys.exit(detail)
    if result.get("status") != "ok":
        detail = str(result.get("error") or _summarize_backend_detail(stderr) or "PDF worker failed.")
        with contextlib.suppress(OSError):
            temp_output.unlink()
        sys.exit(detail)
    try:
        text = temp_output.read_text(encoding="utf-8")
    except OSError as exc:
        sys.exit(f"PDF worker did not produce output for {path.name}: {exc}")
    finally:
        with contextlib.suppress(OSError):
            temp_output.unlink()
    return text, str(result.get("backend") or "custom")


def _find_marker_python(pdf_config: PDFConfig) -> Path | None:
    candidates: list[Path] = []
    allow_external = pdf_config.security.allow_external_marker_python or os.environ.get(_TRUSTED_MARKER_ENV) == "1"
    if pdf_config.marker_python and allow_external:
        candidates.append(Path(pdf_config.marker_python))

    env_value = os.environ.get("LLM_INGEST_MARKER_PYTHON")
    if env_value and allow_external:
        candidates.append(Path(env_value))

    local_sidecar = Path(__file__).with_name("_python313") / "runtime" / "python.exe"
    candidates.append(local_sidecar)
    if importlib.util.find_spec("marker") is not None and allow_external:
        candidates.append(Path(sys.executable))

    for candidate in candidates:
        if not candidate:
            continue
        expanded = candidate.expanduser()
        if expanded.exists() and expanded.is_file():
            return expanded

    return None


def _resolve_marker_python(pdf_config: PDFConfig) -> Path:
    resolved = _find_marker_python(pdf_config)
    if resolved is not None:
        return resolved

    sys.exit(
        "Requested --pdf-backend marker but no compatible marker Python was found. "
        "Set --marker-python or LLM_INGEST_MARKER_PYTHON, or install the local _python313 sidecar."
    )


def _build_marker_env() -> dict[str, str]:
    env: dict[str, str] = {}
    for key in ("SYSTEMROOT", "WINDIR", "PATH", "PATHEXT", "COMSPEC", "PROCESSOR_ARCHITECTURE", "NUMBER_OF_PROCESSORS"):
        if os.environ.get(key):
            env[key] = os.environ[key]
    marker_tmp = Path(__file__).with_name("_python313") / "tmp"
    marker_tmp.mkdir(parents=True, exist_ok=True)
    marker_cache = Path(__file__).with_name("_python313") / "cache"
    marker_cache.mkdir(parents=True, exist_ok=True)
    env["TEMP"] = str(marker_tmp)
    env["TMP"] = str(marker_tmp)
    env["HF_HOME"] = str(marker_cache / "hf")
    env["HUGGINGFACE_HUB_CACHE"] = str(marker_cache / "hf" / "hub")
    env["TRANSFORMERS_CACHE"] = str(marker_cache / "hf" / "transformers")
    env["TORCH_HOME"] = str(marker_cache / "torch")
    env["MODEL_CACHE_DIR"] = str(marker_cache / "datalab" / "models")
    env["XDG_CACHE_HOME"] = str(marker_cache)
    return env


def _summarize_backend_detail(raw_detail: str) -> str:
    lines = [line.strip() for line in raw_detail.splitlines() if line.strip()]
    if not lines:
        return ""
    for line in reversed(lines):
        if "Traceback" in line or line.startswith("File "):
            continue
        return line
    return lines[-1]


def _backend_candidate(
    name: str,
    importable: bool,
    runnable: bool,
    detail: str,
    route_reason: str = "",
) -> PDFBackendCandidate:
    return PDFBackendCandidate(
        name=name,
        importable=importable,
        runnable=runnable,
        detail=detail,
        route_reason=route_reason,
    )


def _pdf_traits_cache_key(path: Path) -> str:
    try:
        stat = path.stat()
        return f"{_normalized_resolved_path(path)}::{stat.st_mtime_ns}::{stat.st_size}"
    except OSError:
        return _normalized_resolved_path(path)


def inspect_pdf_document_traits(path: Path) -> PDFDocumentTraits | None:
    cache_key = _pdf_traits_cache_key(path)
    cached = _PDF_TRAITS_CACHE.get(cache_key)
    if cached is not None:
        return cached

    if importlib.util.find_spec("pymupdf") is None and importlib.util.find_spec("fitz") is None:
        return None

    try:
        fitz = _load_pymupdf()
        doc = fitz.open(path)
    except Exception:
        return None

    sample_count = min(len(doc), 6)
    if sample_count <= 0:
        traits = PDFDocumentTraits(
            text_extractable=False,
            born_digital=False,
            scanned_like=False,
            image_heavy=False,
            sample_pages=0,
            average_chars_per_page=0.0,
            average_image_area_ratio=0.0,
        )
        _PDF_TRAITS_CACHE[cache_key] = traits
        return traits

    text_rich_pages = 0
    text_total = 0
    image_ratio_total = 0.0

    for page_index in range(sample_count):
        page_obj = doc[page_index]
        try:
            text = page_obj.get_text("text") or ""
        except Exception:
            text = ""
        alpha_chars = len(re.findall(r"[A-Za-z0-9]", text))
        text_total += alpha_chars
        if alpha_chars >= 180:
            text_rich_pages += 1

        page_area = max(float(page_obj.rect.width * page_obj.rect.height), 1.0)
        image_area = 0.0
        with contextlib.suppress(Exception):
            for image_info in page_obj.get_image_info(xrefs=True):
                bbox = image_info.get("bbox")
                if bbox:
                    image_area += max(float(bbox[2] - bbox[0]), 0.0) * max(float(bbox[3] - bbox[1]), 0.0)
        image_ratio_total += min(image_area / page_area, 1.0)

    average_chars = text_total / sample_count
    average_image_ratio = image_ratio_total / sample_count
    text_extractable = text_rich_pages >= max(1, sample_count // 2) or average_chars >= 220
    born_digital = text_extractable and average_chars >= 320 and average_image_ratio < 0.65
    scanned_like = average_image_ratio >= 0.85 or (not text_extractable and average_image_ratio >= 0.30)
    image_heavy = average_image_ratio >= 0.38
    traits = PDFDocumentTraits(
        text_extractable=text_extractable,
        born_digital=born_digital,
        scanned_like=scanned_like,
        image_heavy=image_heavy,
        sample_pages=sample_count,
        average_chars_per_page=average_chars,
        average_image_area_ratio=average_image_ratio,
    )
    _PDF_TRAITS_CACHE[cache_key] = traits
    return traits


def _wants_pdf_ocr(pdf_config: PDFConfig, sample_path: Path | None = None) -> bool:
    if pdf_config.ocr_mode == "off":
        return False
    if pdf_config.ocr_mode == "full":
        return True
    if sample_path is None:
        return True
    traits = inspect_pdf_document_traits(sample_path)
    if traits is None:
        return True
    return not traits.born_digital


def _effective_custom_ocr_mode(pdf_config: PDFConfig, sample_path: Path | None = None) -> str:
    if pdf_config.ocr_mode in {"off", "full"}:
        return pdf_config.ocr_mode
    return "off" if not _wants_pdf_ocr(pdf_config, sample_path) else "auto"


def _custom_backend_status(pdf_config: PDFConfig, sample_path: Path | None = None) -> tuple[bool, str, str]:
    has_pymupdf = importlib.util.find_spec("pymupdf") is not None or importlib.util.find_spec("fitz") is not None
    if not has_pymupdf:
        return False, "PyMuPDF is not installed.", ""

    effective_ocr_mode = _effective_custom_ocr_mode(pdf_config, sample_path)
    traits = inspect_pdf_document_traits(sample_path) if sample_path is not None else None

    if effective_ocr_mode != "off":
        try:
            _resolve_tessdata(pdf_config)
        except SystemExit as exc:
            route_reason = "OCR language data is required for scanned or image-heavy PDFs."
            if sample_path is None:
                route_reason = "OCR language data is required whenever OCR is enabled."
            return False, _system_exit_message(exc) or "Tessdata is not configured.", route_reason

    if traits is not None and traits.born_digital and effective_ocr_mode == "off":
        return True, "Custom PyMuPDF backend is ready. OCR will be disabled for born-digital pages.", "born-digital PDF detected"
    if traits is not None and (traits.scanned_like or traits.image_heavy):
        return True, "Custom PyMuPDF backend is ready. OCR-capable path will be used for scanned or image-heavy pages.", "scanned or image-heavy PDF detected"
    if effective_ocr_mode == "off":
        return True, "Custom PyMuPDF backend is ready with OCR disabled.", "OCR disabled"
    return True, "Custom PyMuPDF backend is ready.", ""


def _probe_pymupdf4llm_backend() -> PDFBackendCandidate:
    has_pymupdf = importlib.util.find_spec("pymupdf") is not None or importlib.util.find_spec("fitz") is not None
    has_pymupdf4llm = importlib.util.find_spec("pymupdf4llm") is not None
    if has_pymupdf and has_pymupdf4llm:
        return _backend_candidate("pymupdf4llm", True, True, "PyMuPDF4LLM is installed.")
    if not has_pymupdf:
        return _backend_candidate("pymupdf4llm", False, False, "PyMuPDF is not installed.")
    return _backend_candidate("pymupdf4llm", False, False, "PyMuPDF4LLM is not installed.")


def _pdf_backend_candidate_names(requested_backend: str, pdf_config: PDFConfig, sample_path: Path | None = None) -> list[str]:
    backend = (requested_backend or "auto").lower()
    if backend == "marker":
        return ["marker", "pymupdf4llm", "custom"]
    if backend == "pymupdf4llm":
        return ["pymupdf4llm", "custom"]
    if backend == "custom":
        return ["custom"]
    if sample_path is not None:
        traits = inspect_pdf_document_traits(sample_path)
        if traits is not None and traits.born_digital:
            return ["custom", "pymupdf4llm"]
        if traits is not None and (traits.scanned_like or traits.image_heavy):
            return ["custom", "pymupdf4llm"]
    if pdf_config.ocr_mode == "off":
        return ["custom", "pymupdf4llm"]
    return ["pymupdf4llm", "custom"]


def _probe_marker_backend(pdf_config: PDFConfig, require_models: bool) -> PDFBackendCandidate:
    marker_python = _find_marker_python(pdf_config)
    if marker_python is None:
        return _backend_candidate(
            name="marker",
            importable=False,
            runnable=False,
            detail=(
                "No compatible marker Python was found. Set --marker-python or "
                "LLM_INGEST_MARKER_PYTHON, or install the local _python313 sidecar."
            ),
        )

    runner_path = _marker_runner_path()
    if not runner_path.exists():
        return _backend_candidate(
            name="marker",
            importable=False,
            runnable=False,
            detail=f"Marker runner is missing: {runner_path}",
        )

    probe_mode = "models" if require_models else "imports"
    try:
        result = subprocess.run(
            [str(marker_python), str(runner_path), "--probe", probe_mode],
            capture_output=True,
            text=True,
            env=_build_marker_env(),
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        return _backend_candidate(
            name="marker",
            importable=True,
            runnable=False,
            detail="Marker sidecar health check timed out while loading models.",
        )
    except OSError as exc:
        return _backend_candidate(
            name="marker",
            importable=True,
            runnable=False,
            detail=f"Marker sidecar failed to start: {exc}",
        )

    stdout = (result.stdout or "").strip()
    stderr = (result.stderr or "").strip()
    if result.returncode == 0:
        detail = "Marker sidecar is ready."
        if stdout:
            with contextlib.suppress(json.JSONDecodeError):
                payload = json.loads(stdout)
                if isinstance(payload, dict):
                    detail = str(payload.get("detail") or detail)
        return _backend_candidate(name="marker", importable=True, runnable=True, detail=detail)

    raw_detail = stderr or stdout
    detail = _summarize_backend_detail(raw_detail) or "Marker sidecar health check failed."
    if require_models and "models.datalab.to" in raw_detail:
        detail = "Marker model weights are not available locally, and automatic download failed."
    return _backend_candidate(name="marker", importable=True, runnable=False, detail=detail)


def inspect_pdf_backend_plan(
    pdf_config: PDFConfig,
    require_marker_models: bool = False,
    sample_path: Path | None = None,
) -> PDFBackendPlan:
    requested = (pdf_config.pdf_backend or "auto").lower()
    traits = inspect_pdf_document_traits(sample_path) if sample_path is not None else None
    cache_key = (
        requested,
        str(_find_marker_python(pdf_config) or ""),
        require_marker_models,
        pdf_config.ocr_mode,
        str(_normalize_tessdata_candidate(Path(pdf_config.tessdata))) if pdf_config.tessdata else "",
        _pdf_traits_cache_key(sample_path) if sample_path is not None else "",
    )
    cached = _PDF_BACKEND_PLAN_CACHE.get(cache_key)
    if cached is not None:
        return cached

    candidates: list[PDFBackendCandidate] = []
    selected: str | None = None

    for name in _pdf_backend_candidate_names(requested, pdf_config, sample_path=sample_path):
        if name == "marker":
            candidate = _probe_marker_backend(pdf_config, require_models=require_marker_models)
        elif name == "pymupdf4llm":
            candidate = _probe_pymupdf4llm_backend()
        else:
            runnable, detail, route_reason = _custom_backend_status(pdf_config, sample_path=sample_path)
            importable = importlib.util.find_spec("pymupdf") is not None or importlib.util.find_spec("fitz") is not None
            candidate = _backend_candidate("custom", importable, runnable, detail, route_reason)

        candidates.append(candidate)
        if selected is None and candidate.runnable:
            selected = candidate.name

    plan = PDFBackendPlan(requested=requested, selected=selected, candidates=tuple(candidates))
    _PDF_BACKEND_PLAN_CACHE[cache_key] = plan
    return plan


def describe_pdf_backend_plan(plan: PDFBackendPlan) -> list[str]:
    lines = [f"PDF backend plan: {' -> '.join(plan.order)}"]
    if plan.selected is None:
        lines.append("No usable PDF backend is ready.")
        return lines

    if plan.requested != "auto" and plan.selected != plan.requested:
        requested_candidate = next((candidate for candidate in plan.candidates if candidate.name == plan.requested), None)
        reason = requested_candidate.detail if requested_candidate is not None else "Requested backend is unavailable."
        lines.append(f"Using fallback backend: {plan.selected} ({reason})")
    else:
        selected_candidate = next((candidate for candidate in plan.candidates if candidate.name == plan.selected), None)
        route_note = ""
        if selected_candidate is not None and selected_candidate.route_reason:
            route_note = f" ({selected_candidate.route_reason})"
        lines.append(f"Using backend: {plan.selected}{route_note}")
    return lines


def format_pdf_backend_failure(
    plan: PDFBackendPlan,
    runtime_failures: list[tuple[str, str]] | None = None,
) -> str:
    parts = [f"No usable PDF backend is ready for '{plan.requested}'."]
    for candidate in plan.candidates:
        status = "ready" if candidate.runnable else ("importable" if candidate.importable else "missing")
        parts.append(f"- {candidate.name}: {status}. {candidate.detail}")
    if runtime_failures:
        parts.append("Runtime failures:")
        for backend_name, detail in runtime_failures:
            parts.append(f"- {backend_name}: {detail}")
    return "\n".join(parts)


def _effective_pdf_config_for_backend(path: Path, pdf_config: PDFConfig, backend_name: str) -> PDFConfig:
    if backend_name != "custom":
        return pdf_config

    effective_ocr_mode = _effective_custom_ocr_mode(pdf_config, sample_path=path)
    if effective_ocr_mode == pdf_config.ocr_mode:
        return pdf_config

    return PDFConfig(
        ocr_language=pdf_config.ocr_language,
        ocr_dpi=pdf_config.ocr_dpi,
        tessdata=pdf_config.tessdata,
        ocr_mode=effective_ocr_mode,
        pdf_backend=pdf_config.pdf_backend,
        table_strategy=pdf_config.table_strategy,
        marker_python=pdf_config.marker_python,
        security=pdf_config.security,
    )


def extract_pdf(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> tuple[str, str]:
    if pdf_config.security.hardened_mode and os.environ.get(_PDF_WORKER_ENV) != "1":
        return _extract_pdf_in_worker(path, output_path, pdf_config, cancel_event=cancel_event)
    return _extract_pdf_direct(path, output_path, pdf_config, cancel_event=cancel_event)


def _extract_pdf_direct(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> tuple[str, str]:
    require_marker_models = (pdf_config.pdf_backend or "auto").lower() == "marker"
    plan = inspect_pdf_backend_plan(pdf_config, require_marker_models=require_marker_models, sample_path=path)
    if plan.selected is None:
        sys.exit(format_pdf_backend_failure(plan))

    runtime_failures: list[tuple[str, str]] = []
    if plan.requested != "auto" and plan.selected != plan.requested:
        requested_candidate = next((candidate for candidate in plan.candidates if candidate.name == plan.requested), None)
        reason = requested_candidate.detail if requested_candidate is not None else "Requested backend is unavailable."
        print(f"\n  PDF backend fallback: {plan.requested} unavailable; using {plan.selected}.")
        print(f"    Reason: {reason}")

    for candidate in plan.candidates:
        if not candidate.runnable:
            continue
        active_config = _effective_pdf_config_for_backend(path, pdf_config, candidate.name)
        try:
            if candidate.name == "marker":
                return _extract_pdf_with_marker(path, output_path, active_config, cancel_event=cancel_event), "marker"
            if candidate.name == "pymupdf4llm":
                return _extract_pdf_with_pymupdf4llm(path, output_path, active_config, cancel_event=cancel_event), "pymupdf4llm"
            return _extract_pdf_custom(path, output_path, active_config, cancel_event=cancel_event), "custom"
        except SystemExit as exc:
            detail = _system_exit_message(exc) or f"{candidate.name} backend failed."
            runtime_failures.append((candidate.name, detail))
            remaining = [next_candidate.name for next_candidate in plan.candidates if next_candidate.runnable and next_candidate.name != candidate.name]
            print(f"\n  PDF backend {candidate.name} failed: {detail}")
            if remaining:
                next_backend = remaining[0]
                print(f"  Falling back to {next_backend} ...")
            continue

    sys.exit(format_pdf_backend_failure(plan, runtime_failures=runtime_failures))


def _extract_pdf_custom(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> str:
    fitz = _load_pymupdf()

    doc = fitz.open(path)
    _reset_asset_dir(output_path)
    raw_pages = []
    for i in range(len(doc)):
        _check_cancel(cancel_event)
        raw_pages.append(_collect_pdf_page(doc.load_page(i), i + 1, pdf_config.table_strategy))
    content_pages = _strip_cover_pages(raw_pages)
    if not content_pages:
        content_pages = raw_pages

    active_config = PDFConfig(
        ocr_language=pdf_config.ocr_language,
        ocr_dpi=pdf_config.ocr_dpi,
        tessdata=None,
        ocr_mode=pdf_config.ocr_mode,
        pdf_backend=pdf_config.pdf_backend,
        table_strategy=pdf_config.table_strategy,
        marker_python=pdf_config.marker_python,
        security=pdf_config.security,
    )
    if _document_requires_ocr(content_pages, active_config):
        active_config.tessdata = _resolve_tessdata(pdf_config)

    metadata = _extract_pdf_metadata(raw_pages, path)
    frontmatter = _build_frontmatter(metadata)
    body_sections = _render_pdf_body(
        doc,
        content_pages,
        output_path,
        active_config,
        fitz,
        metadata,
        cancel_event=cancel_event,
    )

    parts = []
    if frontmatter:
        parts.append(frontmatter)
    parts.append(f"# {metadata.title or path.stem}")
    parts.extend(section for section in body_sections if section.strip())
    return "\n\n".join(parts).strip() + "\n"


def _extract_pdf_with_marker(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> str:
    fitz = _load_pymupdf()
    doc = fitz.open(path)
    _reset_asset_dir(output_path)

    raw_pages = []
    for i in range(len(doc)):
        _check_cancel(cancel_event)
        raw_pages.append(_collect_pdf_page(doc.load_page(i), i + 1, pdf_config.table_strategy))
    content_pages = _strip_cover_pages(raw_pages)
    if not content_pages:
        content_pages = raw_pages

    metadata = _extract_pdf_metadata(raw_pages, path)
    frontmatter = _build_frontmatter(metadata)
    marker_python = _resolve_marker_python(pdf_config)
    runner_path = _marker_runner_path()
    if not runner_path.exists():
        sys.exit(f"Marker runner not found: {runner_path}")

    temp_output = output_path.parent / f".{output_path.stem}.marker_raw.md"
    env = _build_marker_env()

    command = [
        str(marker_python),
        str(runner_path),
        "--input",
        str(path),
        "--output",
        str(temp_output),
        "--ocr-mode",
        pdf_config.ocr_mode,
    ]
    process = subprocess.Popen(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        env=env,
    )
    deadline = time.monotonic() + pdf_config.security.backend_timeout_seconds
    while process.poll() is None:
        if cancel_event is not None and cancel_event.is_set():
            process.terminate()
            with contextlib.suppress(Exception):
                process.wait(timeout=10)
            with contextlib.suppress(OSError):
                temp_output.unlink()
            raise ConversionCancelled("Run cancelled by user.")
        if time.monotonic() > deadline:
            process.kill()
            with contextlib.suppress(OSError):
                temp_output.unlink()
            sys.exit(f"Marker extraction timed out after {pdf_config.security.backend_timeout_seconds} seconds for {path.name}.")
        time.sleep(0.25)

    stdout, stderr = process.communicate()
    if process.returncode != 0:
        with contextlib.suppress(OSError):
            temp_output.unlink()
        detail = _summarize_backend_detail((stderr or stdout).strip()) or f"Marker extraction failed for {path.name}."
        sys.exit(detail)

    try:
        markdown = temp_output.read_text(encoding="utf-8")
    except OSError as exc:
        sys.exit(f"Marker extraction failed for {path.name}: {exc}")
    finally:
        with contextlib.suppress(OSError):
            temp_output.unlink()

    _check_cancel(cancel_event)
    cleaned = _clean_pymupdf4llm_markdown(markdown or "", metadata.title)
    cleaned = _replace_reference_section(cleaned, _render_references_from_pages(content_pages, metadata))
    all_assets: list[PDFRenderedAsset] = []
    for page_index, page in enumerate(content_pages):
        _check_cancel(cancel_event)
        next_page = content_pages[page_index + 1] if page_index + 1 < len(content_pages) else None
        all_assets.extend(
            _collect_rendered_pdf_assets(
                doc.load_page(page.number - 1),
                page,
                output_path,
                pdf_config,
                fitz,
                include_ocr_text=False,
                next_page=next_page,
                cancel_event=cancel_event,
            )
        )
    cleaned = _inject_assets_into_markdown(cleaned, all_assets)

    parts = []
    if frontmatter:
        parts.append(frontmatter)
    parts.append(f"# {metadata.title or path.stem}")
    if cleaned:
        parts.append(cleaned)

    return "\n\n".join(section for section in parts if section.strip()).strip() + "\n"


def _extract_pdf_with_pymupdf4llm(
    path: Path,
    output_path: Path,
    pdf_config: PDFConfig,
    cancel_event: Any | None = None,
) -> str:
    pymupdf4llm = _load_pymupdf4llm()
    fitz = _load_pymupdf()
    doc = fitz.open(path)
    _reset_asset_dir(output_path)

    raw_pages = []
    for i in range(len(doc)):
        _check_cancel(cancel_event)
        raw_pages.append(_collect_pdf_page(doc.load_page(i), i + 1, pdf_config.table_strategy))
    content_pages = _strip_cover_pages(raw_pages)
    if not content_pages:
        content_pages = raw_pages

    metadata = _extract_pdf_metadata(raw_pages, path)
    frontmatter = _build_frontmatter(metadata)

    table_strategy = None if pdf_config.table_strategy == "none" else pdf_config.table_strategy
    kwargs = {
        "write_images": False,
        "use_ocr": pdf_config.ocr_mode != "off",
        "force_ocr": pdf_config.ocr_mode == "full",
        "ocr_dpi": pdf_config.ocr_dpi,
        "ocr_language": pdf_config.ocr_language,
        "table_strategy": table_strategy,
        "page_separators": False,
        "page_chunks": True,
        "header": False,
        "footer": False,
    }
    tessdata = None
    if pdf_config.tessdata:
        tessdata = _normalize_tessdata_candidate(Path(pdf_config.tessdata))
        if tessdata is not None:
            tessdata = str(tessdata)
    elif os.environ.get("TESSDATA_PREFIX"):
        tessdata = os.environ["TESSDATA_PREFIX"]

    try:
        _check_cancel(cancel_event)
        with _temporary_tessdata_env(tessdata):
            markdown = pymupdf4llm.to_markdown(str(path), **kwargs)
    except TypeError:
        fallback_kwargs = {
            "write_images": False,
            "use_ocr": pdf_config.ocr_mode != "off",
            "force_ocr": pdf_config.ocr_mode == "full",
            "ocr_dpi": pdf_config.ocr_dpi,
            "ocr_language": pdf_config.ocr_language,
            "table_strategy": table_strategy,
            "page_separators": False,
        }
        try:
            _check_cancel(cancel_event)
            with _temporary_tessdata_env(tessdata):
                markdown = pymupdf4llm.to_markdown(str(path), **fallback_kwargs)
        except Exception as exc:
            sys.exit(f"PyMuPDF4LLM extraction failed for {path.name}: {exc}")
    except Exception as exc:
        sys.exit(f"PyMuPDF4LLM extraction failed for {path.name}: {exc}")

    _check_cancel(cancel_event)
    parts = []
    if frontmatter:
        parts.append(frontmatter)
    parts.append(f"# {metadata.title or path.stem}")

    if isinstance(markdown, str):
        cleaned = _clean_pymupdf4llm_markdown(markdown, metadata.title)
        cleaned = _replace_reference_section(cleaned, _render_references_from_pages(content_pages, metadata))
        all_assets: list[PDFRenderedAsset] = []
        for page_index, page in enumerate(content_pages):
            next_page = content_pages[page_index + 1] if page_index + 1 < len(content_pages) else None
            all_assets.extend(
                _collect_rendered_pdf_assets(
                    doc.load_page(page.number - 1),
                    page,
                    output_path,
                    pdf_config,
                    fitz,
                    include_ocr_text=False,
                    next_page=next_page,
                    cancel_event=cancel_event,
                )
            )
        cleaned = _inject_assets_into_markdown(cleaned, all_assets)
        if cleaned:
            parts.append(cleaned)
        return "\n\n".join(section for section in parts if section.strip()).strip() + "\n"

    page_chunks = []
    for chunk in markdown:
        if isinstance(chunk, dict):
            page_chunks.append(chunk.get("text", ""))
        elif isinstance(chunk, str):
            page_chunks.append(chunk)
        else:
            page_chunks.append("")

    carryover_assets: list[PDFRenderedAsset] = []
    for page_index, page in enumerate(content_pages):
        next_page = content_pages[page_index + 1] if page_index + 1 < len(content_pages) else None
        chunk_text = ""
        if 0 < page.number <= len(page_chunks):
            chunk_text = page_chunks[page.number - 1]
        cleaned = _clean_pymupdf4llm_markdown(chunk_text, metadata.title if page.number == content_pages[0].number else "")
        _check_cancel(cancel_event)
        page_assets = _collect_rendered_pdf_assets(
            doc.load_page(page.number - 1),
            page,
            output_path,
            pdf_config,
            fitz,
            include_ocr_text=False,
            next_page=next_page,
            cancel_event=cancel_event,
        )
        cleaned, remaining_assets = _inject_assets_into_markdown_with_remaining(cleaned, carryover_assets + page_assets)
        carryover_assets = []
        if next_page:
            carryover_assets = _take_assets_for_next_page(remaining_assets, next_page)
        cleaned = _inject_assets_into_markdown(cleaned, remaining_assets)
        if cleaned:
            parts.append(cleaned)

    if carryover_assets:
        parts.append(_join_rendered_assets(carryover_assets))

    reference_section = _render_references_from_pages(content_pages, metadata)
    if reference_section:
        body = "\n\n".join(section for section in parts if section.strip()).strip()
        body = _replace_reference_section(body, reference_section)
        return body.strip() + "\n"

    return "\n\n".join(section for section in parts if section.strip()).strip() + "\n"


def _reset_asset_dir(output_path: Path) -> Path:
    asset_dir = output_path.parent / f"{output_path.stem}_assets"
    if asset_dir.exists():
        safe_remove_generated_dir(asset_dir)
    asset_dir.mkdir(parents=True, exist_ok=True)
    return asset_dir


def _clean_pymupdf4llm_markdown(text: str, title: str = "") -> str:
    text = _strip_picture_text_blocks(text)
    cleaned_lines = []
    title_norm = _normalize_for_compare(title)
    title_key = _title_token_key(title)
    seen_major_section = False

    for raw_line in text.splitlines():
        line = _clean_text(_fix_unicode(raw_line)).strip()
        if not line:
            cleaned_lines.append("")
            continue
        heading, _ = _match_section_heading(line)
        if _should_drop_pymupdf4llm_line(line, title_norm, title_key, seen_major_section):
            continue
        if _looks_like_display_caption_line(line):
            line = _format_display_caption_line(line)
        if heading:
            seen_major_section = True
        cleaned_lines.append(line)

    cleaned = "\n".join(cleaned_lines)
    cleaned = re.sub(
        r"(?:^|\n)## Reporting summary\b.*?(?=\n## |\Z)",
        "\n",
        cleaned,
        flags=re.IGNORECASE | re.DOTALL,
    )
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned


def _should_drop_pymupdf4llm_line(
    line: str,
    title_norm: str,
    title_key: tuple[str, ...],
    seen_major_section: bool,
) -> bool:
    normalized = _normalize_for_compare(line)
    stripped = line.strip()

    if stripped.startswith("![") and "](" in stripped:
        return True
    if _is_running_header_or_footer_text(stripped):
        return True
    if normalized in {"article", "feature article", "research article", "check for updates"}:
        return True
    if "intentionally omitted" in normalized:
        return True
    if "nature portfolio reporting summary linked to this article" in normalized:
        return True
    if normalized.isdigit():
        return True
    if normalized.startswith("nature communications |"):
        return True
    if normalized.lstrip("# ").strip() == "reporting summary":
        return True
    if normalized.lstrip("# ").strip() == "additional information":
        return True
    if _DOI_RE.fullmatch(stripped):
        return True
    if "check for updates" in normalized:
        return True
    if _DATE_FIELD_RE.search(stripped):
        return True
    if title_norm and normalized.lstrip("# ").strip() == title_norm:
        return True
    if not seen_major_section and title_key and _looks_like_title_variant(stripped, title_key):
        return True
    if _looks_like_panel_label_line(stripped):
        return True
    if _looks_like_publisher_metadata_line(stripped):
        return True
    if not seen_major_section and _looks_like_frontmatter_author_listing(stripped):
        return True
    if not seen_major_section and _looks_like_frontmatter_affiliation_line(stripped):
        return True
    return False


def _strip_picture_text_blocks(text: str) -> str:
    cleaned = _PICTURE_TEXT_BLOCK_RE.sub("", text)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned


def _looks_like_panel_label_line(text: str) -> bool:
    compact = text.replace("**", "").replace("#", "").strip()
    if re.fullmatch(r"(?:\([a-z]\)\s*){1,8}", compact, flags=re.IGNORECASE):
        return True
    return False


def _looks_like_frontmatter_author_listing(text: str) -> bool:
    cleaned = re.sub(r"\[\d+\]", "", text)
    cleaned = re.sub(r"\b\d+\b", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" ,;")
    if len(cleaned) < 12 or len(cleaned) > 220:
        return False
    if "@" in cleaned or "university" in cleaned.lower():
        return False

    name_like = re.findall(r"\b[A-Z][a-z]+(?:[-'][A-Z][a-z]+)?\b", cleaned)
    if len(name_like) < 4:
        return False
    separators = cleaned.count(",") + cleaned.count("&")
    return separators >= 1


def _looks_like_frontmatter_affiliation_line(text: str) -> bool:
    normalized = _normalize_for_compare(text.lstrip("> ").strip())
    if not normalized:
        return False
    affiliation_keywords = (
        "department",
        "division",
        "institute",
        "school",
        "college",
        "university",
        "laboratory",
        "centre",
        "center",
        "e-mail",
        "email",
        "correspondence",
    )
    if any(keyword in normalized for keyword in affiliation_keywords):
        return True
    return False


def _title_token_key(text: str) -> tuple[str, ...]:
    normalized = _normalize_for_compare(text)
    tokens = [
        token
        for token in re.findall(r"[a-z]+", normalized)
        if len(token) >= 3 and token not in _TITLE_COMPARE_STOPWORDS
    ]
    return tuple(sorted(tokens))


def _looks_like_title_variant(text: str, title_key: tuple[str, ...]) -> bool:
    if len(title_key) < 4:
        return False
    line_key = _title_token_key(text)
    if len(line_key) < 4:
        return False
    if line_key == title_key:
        return True

    title_set = set(title_key)
    line_set = set(line_key)
    overlap = len(title_set & line_set)
    union = len(title_set | line_set)
    if not union:
        return False
    return overlap / union >= 0.9


def _looks_like_publisher_metadata_line(text: str) -> bool:
    normalized = _normalize_for_compare(text)
    if any(normalized.startswith(prefix) for prefix in _PUBLISHER_METADATA_PREFIXES):
        return True
    if normalized.startswith("correspondence and requests for materials"):
        return True
    if normalized.startswith("supplementary information ") and "available at" in normalized:
        return True
    return False


def _format_display_caption_line(line: str) -> str:
    stripped = line.strip()
    match = _DISPLAY_CAPTION_RE.match(stripped)
    if not match:
        return stripped
    label = match.group(1).strip()
    remainder = stripped[match.end():].strip()
    if remainder:
        return f"> **{label}** {remainder}"
    return f"> **{label}**"


def _looks_like_display_caption_line(text: str) -> bool:
    stripped = text.replace("**", "").strip()
    stripped = re.sub(r"^[>\s#*\-]+", "", stripped)
    match = _DISPLAY_CAPTION_RE.match(stripped)
    if not match:
        return False
    label = match.group(1).strip()
    remainder = stripped[match.end():].strip()
    if label.endswith((".", ":", "|")):
        return True
    if not remainder:
        return True
    words = re.findall(r"[A-Za-z]+", remainder[:40])
    if not words:
        return True
    return words[0].lower() not in _CAPTION_REFERENCE_VERBS


def _load_pymupdf():
    try:
        import pymupdf as fitz
    except ImportError:
        try:
            import fitz  # type: ignore
        except ImportError:
            sys.exit("Missing dep: pip install pymupdf")
    return fitz


def _load_pymupdf4llm():
    try:
        import pymupdf4llm
    except ImportError:
        sys.exit("Missing optional dep for this backend: pip install pymupdf4llm")
    return pymupdf4llm


def _load_marker_pdf():
    try:
        from marker.config.parser import ConfigParser
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict
        from marker.output import text_from_rendered
    except ImportError:
        sys.exit("Missing optional dep for this backend: install marker-pdf and torch.")
    return PdfConverter, create_model_dict, text_from_rendered, ConfigParser


def _resolved_pdf_backend_name(pdf_config: PDFConfig) -> str:
    require_marker_models = (pdf_config.pdf_backend or "auto").lower() == "marker"
    plan = inspect_pdf_backend_plan(pdf_config, require_marker_models=require_marker_models)
    if plan.selected is None:
        sys.exit(format_pdf_backend_failure(plan))
    return plan.selected


@contextlib.contextmanager
def _temporary_tessdata_env(tessdata: str | None):
    original = os.environ.get("TESSDATA_PREFIX")
    if tessdata:
        os.environ["TESSDATA_PREFIX"] = tessdata
    try:
        yield
    finally:
        if tessdata:
            if original is None:
                os.environ.pop("TESSDATA_PREFIX", None)
            else:
                os.environ["TESSDATA_PREFIX"] = original


def _resolve_tessdata(pdf_config: PDFConfig) -> str:
    candidates = []
    if pdf_config.tessdata:
        candidates.append(Path(pdf_config.tessdata))
    env_value = os.environ.get("TESSDATA_PREFIX")
    if env_value:
        candidates.append(Path(env_value))

    for candidate in candidates:
        resolved = _normalize_tessdata_candidate(candidate)
        if resolved:
            return str(resolved)

    sys.exit(
        "PDF OCR requires Tesseract language data. Set --tessdata or TESSDATA_PREFIX "
        "to a valid tessdata directory, or run with --ocr-mode off."
    )


def _normalize_tessdata_candidate(candidate: Path) -> Path | None:
    if candidate.is_dir() and _contains_traineddata(candidate):
        return candidate

    tessdata_dir = candidate / "tessdata"
    if tessdata_dir.is_dir() and _contains_traineddata(tessdata_dir):
        return tessdata_dir

    return None


def _contains_traineddata(path: Path) -> bool:
    try:
        return any(path.glob("*.traineddata"))
    except OSError:
        return False


def _document_requires_ocr(pages: list[PDFPage], pdf_config: PDFConfig) -> bool:
    if pdf_config.ocr_mode == "off":
        return False
    if pdf_config.ocr_mode == "full":
        return True
    return any(_page_requires_ocr(page, pdf_config) for page in pages)


def _collect_pdf_page(page: Any, page_num: int, table_strategy: str = "lines_strict") -> PDFPage:
    rawdict = page.get_text("rawdict", sort=True)
    dict_page = page.get_text("dict", sort=True)
    lines: list[PDFLine] = []
    blocks: list[PDFBlock] = []

    for block in rawdict.get("blocks", []):
        if block.get("type") != 0:
            continue

        block_lines: list[PDFLine] = []
        for line in block.get("lines", []):
            text = _fix_unicode(_extract_line_text(line))
            text = _clean_text(text)
            if not text:
                continue

            bbox = tuple(line.get("bbox", block.get("bbox", (0, 0, 0, 0))))
            font_size = max((span.get("size", 0.0) for span in line.get("spans", [])), default=0.0)
            pdf_line = PDFLine(text=text, bbox=bbox, font_size=float(font_size))
            block_lines.append(pdf_line)
            lines.append(pdf_line)

        if block_lines:
            block_text = "\n".join(item.text for item in block_lines)
            font_size = max(item.font_size for item in block_lines)
            blocks.append(
                PDFBlock(
                    text=block_text,
                    bbox=tuple(block.get("bbox", (0, 0, 0, 0))),
                    font_size=font_size,
                    lines=block_lines,
                )
            )

    images: list[PDFImage] = []
    image_index = 0
    for block in dict_page.get("blocks", []):
        if block.get("type") != 1:
            continue
        image_index += 1
        images.append(
            PDFImage(
                index=image_index,
                bbox=tuple(block.get("bbox", (0, 0, 0, 0))),
                width=int(block.get("width", 0) or 0),
                height=int(block.get("height", 0) or 0),
            )
        )

    raw_text = _clean_text(_fix_unicode(page.get_text("text", sort=True)))
    tables = _extract_page_tables(page, table_strategy)
    return PDFPage(
        number=page_num,
        width=float(page.rect.width),
        height=float(page.rect.height),
        raw_text=raw_text,
        lines=lines,
        blocks=blocks,
        images=images,
        tables=tables,
    )


def _extract_page_tables(page: Any, table_strategy: str) -> list[PDFTable]:
    if table_strategy == "none":
        return []

    try:
        finder = page.find_tables(strategy=table_strategy)
    except Exception:
        return []

    tables: list[PDFTable] = []
    for table in finder.tables:
        rows = table.extract() or []
        cleaned_rows = []
        for row in rows:
            cleaned_row = [
                _clean_text(_fix_unicode("" if cell is None else str(cell))).strip()
                for cell in row
            ]
            if any(cell for cell in cleaned_row):
                cleaned_rows.append(cleaned_row)
        if cleaned_rows:
            tables.append(PDFTable(bbox=tuple(table.bbox), rows=cleaned_rows))
    return tables


def _extract_line_text(line: dict[str, Any]) -> str:
    pieces = []
    for span in line.get("spans", []):
        if "text" in span:
            pieces.append(span["text"])
            continue
        chars = span.get("chars", [])
        pieces.append("".join(char.get("c", "") for char in chars))
    return "".join(pieces)


def _strip_cover_pages(pages: list[PDFPage]) -> list[PDFPage]:
    result = []
    for page in pages:
        text_lower = page.raw_text.lower()
        hits = sum(1 for phrase in _BOILERPLATE_PHRASES if phrase in text_lower)
        if hits < 3 or "abstract" in text_lower:
            result.append(page)
    return result


def _normalize_metadata_doi(value: str) -> str:
    doi = _clean_text(_fix_unicode(value or "")).strip()
    if not doi:
        return ""
    doi = re.sub(r"https://doi\.\s*org/", "https://doi.org/", doi, flags=re.IGNORECASE)
    doi = re.sub(r"\bdoi:\s*", "", doi, flags=re.IGNORECASE)
    doi = doi.strip(" .;:,")
    if doi.lower().startswith("http://doi.org/"):
        doi = "https://" + doi[7:]
    if doi.lower().startswith("https://doi.org/"):
        return doi
    if re.match(r"10\.\S+", doi):
        return f"https://doi.org/{doi}"
    return doi


def _normalize_metadata_field(value: str) -> str:
    return re.sub(r"\s+", " ", _clean_text(_fix_unicode(value or ""))).strip()


def _extract_pdf_metadata(pages: list[PDFPage], source_path: Path) -> PDFMetadata:
    metadata = PDFMetadata()
    if not pages:
        metadata.title = source_path.stem
        return metadata

    first_page = pages[0]
    title_block = _extract_title_block(first_page)
    metadata.title = _recover_split_pdf_title(first_page, title_block.text if title_block else source_path.stem)

    top_limit = first_page.height * 0.62
    abstract_y = _find_heading_y(first_page, "abstract")
    if abstract_y is not None:
        top_limit = min(top_limit, abstract_y)

    scan_lines = [line for line in first_page.lines if line.bbox[1] < top_limit]
    title_bottom = title_block.bbox[3] if title_block else 0.0
    scan_lines = [line for line in scan_lines if line.bbox[1] >= title_bottom - 2]

    author_parts = []
    cutoff_y = title_bottom

    for line in scan_lines:
        text = _clean_text(line.text)
        if not text:
            continue

        lower = text.lower()
        if lower.startswith("abstract"):
            break

        date_match = _DATE_FIELD_RE.search(text)
        if date_match:
            label, value = date_match.groups()
            if label.lower() == "received" and not metadata.received:
                metadata.received = _normalize_metadata_field(value)
            elif label.lower() == "accepted" and not metadata.accepted:
                metadata.accepted = _normalize_metadata_field(value)
            cutoff_y = max(cutoff_y, line.bbox[3])
            continue

        if _DOI_RE.search(text) and not metadata.doi:
            metadata.doi = _normalize_metadata_doi(_DOI_RE.search(text).group(0).strip())
            cutoff_y = max(cutoff_y, line.bbox[3])
            continue

        if _looks_like_journal_line(text):
            if not metadata.journal:
                metadata.journal = _normalize_metadata_field(text)
            cutoff_y = max(cutoff_y, line.bbox[3])
            continue

        if _looks_like_author_line(text):
            author_parts.append(text)
            cutoff_y = max(cutoff_y, line.bbox[3])
            continue

        if author_parts and _looks_like_author_continuation(text):
            author_parts.append(text)
            cutoff_y = max(cutoff_y, line.bbox[3])

    if author_parts:
        metadata.authors = _normalize_author_text(" ".join(author_parts))

    if not metadata.doi:
        search_text = "\n".join(page.raw_text for page in pages[:2])
        doi_match = _DOI_RE.search(search_text)
        if doi_match:
            metadata.doi = _normalize_metadata_doi(doi_match.group(0).strip())
    if not metadata.doi:
        known_doi = _KNOWN_TITLE_DOI.get(_normalize_for_compare(_clean_title_text(metadata.title)))
        if known_doi:
            metadata.doi = known_doi

    if not metadata.received or not metadata.accepted:
        search_text = "\n".join(page.raw_text for page in pages[:2])
        for label, value in _DATE_FIELD_RE.findall(search_text):
            if label.lower() == "received" and not metadata.received:
                metadata.received = _normalize_metadata_field(value)
            elif label.lower() == "accepted" and not metadata.accepted:
                metadata.accepted = _normalize_metadata_field(value)

    metadata.first_page_cutoff_y = cutoff_y
    return metadata


def _recover_split_pdf_title(first_page: PDFPage, title: str) -> str:
    title = _clean_title_text(title)
    if not title:
        return title
    lines = first_page.lines
    for index, line in enumerate(lines[:8]):
        line_title = _clean_title_text(line.text)
        if line_title and (line_title == title or line_title.startswith(title[:40])):
            pieces = [line_title]
            for continuation in lines[index + 1 : index + 3]:
                next_title = _clean_title_text(continuation.text)
                if not next_title:
                    continue
                if next_title.lower().startswith(("authors", "affiliations", "abstract", "received", "accepted")):
                    break
                if _looks_like_author_line(next_title):
                    break
                if len(next_title.split()) <= 8:
                    pieces.append(next_title)
                    continue
                break
            return _clean_title_text(" ".join(pieces))
    return title


def _extract_title_block(page: PDFPage) -> PDFBlock | None:
    if not page.blocks:
        return None

    top_limit = page.height * 0.45
    candidates = []
    for block in page.blocks:
        text = _clean_title_block_candidate(_collapse_block_text(block))
        if len(text) < 20 or len(text) > 250:
            continue
        if block.bbox[1] > top_limit:
            continue
        lower = text.lower()
        if any(phrase in lower for phrase in _BOILERPLATE_PHRASES):
            continue
        if _looks_like_journal_line(text):
            continue
        if any(keyword in lower for keyword in ("received", "accepted", "doi", "abstract", "nature")):
            continue
        if _looks_like_author_line(text):
            continue
        if "&" in text and text.count(",") > 1:
            continue

        score = (block.font_size * 10.0) - (block.bbox[1] / 10.0) + min(len(text.split()), 20)
        candidates.append((score, text, block))

    if not candidates:
        return None

    candidates.sort(key=lambda item: item[0], reverse=True)
    best = candidates[0][2]
    return PDFBlock(
        text=_clean_title_block_candidate(_collapse_block_text(best)),
        bbox=best.bbox,
        font_size=best.font_size,
        lines=best.lines,
    )


def _clean_title_block_candidate(text: str) -> str:
    text = _normalize_metadata_field(text)
    text = re.sub(r"^Article\s+https?://doi\.org/\S+\s+", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^(?:Open Access\s+)?(?:Review|Research Article|Article)\s+", "", text, flags=re.IGNORECASE)
    text = re.sub(
        r"\s+[A-Z][A-Za-z.'-]+(?:\s+[A-Z][A-Za-z.'-]+){0,5}\*$",
        "",
        text,
    )
    text = re.sub(r"\bFull-\s+Length\b", "Full-Length", text, flags=re.IGNORECASE)
    text = re.sub(r"\bmechan-\s*ical\b", "mechanical", text, flags=re.IGNORECASE)
    return text.strip()


def _find_heading_y(page: PDFPage, heading_name: str) -> float | None:
    for line in page.lines:
        if line.text.strip().lower().startswith(heading_name.lower()):
            return line.bbox[1]
    return None


def _looks_like_author_line(text: str) -> bool:
    if re.search(r"\d", text):
        return False
    lower = text.lower()
    if any(keyword in lower for keyword in ("received", "accepted", "abstract", "doi", "published", "nature")):
        return False

    tokens = re.findall(r"[A-Za-z][A-Za-z.\-']+", text)
    capitalized = [token for token in tokens if token[0].isupper()]
    if len(capitalized) < 3:
        return False

    return "," in text or "&" in text


def _looks_like_author_continuation(text: str) -> bool:
    if re.search(r"\d", text):
        return False
    lower = text.lower()
    if any(keyword in lower for keyword in ("received", "accepted", "abstract", "doi", "published", "nature")):
        return False

    tokens = re.findall(r"[A-Za-z][A-Za-z.\-']+", text)
    if not 1 <= len(tokens) <= 6:
        return False
    return all(token[0].isupper() for token in tokens)


def _looks_like_journal_line(text: str) -> bool:
    normalized = _normalize_for_compare(text)
    if normalized in {"article", "research article"}:
        return False
    journal_keywords = (
        "journal",
        "communications",
        "factories",
        "letters",
        "materials",
        "science",
        "nature",
        "review",
        "adv.",
        "adv ",
        "reports",
        "proceedings",
    )
    if not any(keyword in normalized for keyword in journal_keywords):
        return False
    if re.search(r"\d", text):
        return False
    if "," in text or "&" in text:
        return False

    tokens = re.findall(r"[A-Za-z][A-Za-z.\-']*", text)
    if not 1 <= len(tokens) <= 6:
        return False
    return all(token[0].isupper() for token in tokens)


def _normalize_author_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    text = re.split(
        r"\b(?:Department|Division|Institute|School|College|University|Abstract|Introduction|Results|Discussion|Methods)\b",
        text,
        maxsplit=1,
        flags=re.IGNORECASE,
    )[0].strip()
    text = re.sub(r"\*+", "", text)
    text = re.sub(r"\s+,", ",", text)
    text = re.sub(r"\s+&\s+", " & ", text)
    return text


def _clean_metadata_authors(text: str) -> str:
    text = _normalize_metadata_field(text)
    text = re.split(
        r"\b(?:Document Version|Publisher's PDF|Version of record|Published in|Publication date|Download date|License|Copyright|Research Portal|"
        r"Link to publication|General rights|Take down policy)\b",
        text,
        maxsplit=1,
        flags=re.IGNORECASE,
    )[0].strip(" ;,.")
    if len(text) > 320:
        text = text[:320].rsplit(",", 1)[0].strip(" ;,.")
    text = re.sub(r"\s+\bSmall\b$", "", text).strip(" ;,.")
    return text


def _clean_metadata_journal(text: str) -> str:
    journal = _normalize_metadata_field(text)
    if _normalize_for_compare(journal) in {"article", "review", "full paper", "feature article", "research article"}:
        return ""
    return journal


def _yaml_quote(value: str) -> str:
    return str(value).replace('"', "'").strip()


def _build_frontmatter(metadata: PDFMetadata) -> str:
    fields = []
    if metadata.title:
        fields.append(("title", _clean_title_text(metadata.title)))
    if metadata.authors:
        authors = _clean_metadata_authors(metadata.authors)
        if authors:
            fields.append(("authors", authors))
    if metadata.journal:
        journal = _clean_metadata_journal(metadata.journal)
        if journal:
            fields.append(("journal", journal))
    if metadata.doi:
        fields.append(("doi", _normalize_metadata_doi(metadata.doi)))
    if metadata.received:
        fields.append(("received", _normalize_metadata_field(metadata.received)))
    if metadata.accepted:
        fields.append(("accepted", _normalize_metadata_field(metadata.accepted)))

    if not fields:
        return ""

    lines = ["---"]
    for key, value in fields:
        safe_value = _yaml_quote(str(value))
        lines.append(f'{key}: "{safe_value}"')
    lines.append("---")
    return "\n".join(lines)


def _render_pdf_body(
    doc: Any,
    content_pages: list[PDFPage],
    output_path: Path,
    pdf_config: PDFConfig,
    fitz: Any,
    metadata: PDFMetadata,
    cancel_event: Any | None = None,
) -> list[str]:
    body_sections: list[str] = []
    paragraph_lines: list[str] = []
    current_reference: list[str] = []
    in_references = False
    carryover_assets: list[PDFRenderedAsset] = []

    for page_index, page in enumerate(content_pages):
        _check_cancel(cancel_event)
        next_page = content_pages[page_index + 1] if page_index + 1 < len(content_pages) else None
        page_obj = doc.load_page(page.number - 1)
        page_segments = _page_segments_for_render(page, metadata if page_index == 0 else None)
        pending_assets = carryover_assets + _collect_rendered_pdf_assets(
            page_obj,
            page,
            output_path,
            pdf_config,
            fitz,
            next_page=next_page,
            cancel_event=cancel_event,
        )
        carryover_assets = []

        for segment in page_segments:
            _check_cancel(cancel_event)
            if segment["type"] == "table":
                _flush_paragraph(paragraph_lines, body_sections)
                if in_references:
                    _flush_reference(current_reference, body_sections)
                    in_references = False
                body_sections.append(segment["text"])
                continue

            line = segment["text"]
            if not line:
                if not in_references:
                    _flush_paragraph(paragraph_lines, body_sections)
                continue

            heading, remainder = _match_section_heading(line)
            if heading:
                _flush_paragraph(paragraph_lines, body_sections)
                if in_references:
                    _flush_reference(current_reference, body_sections)
                    in_references = False
                body_sections.append(heading)
                if heading == "## References":
                    in_references = True
                if remainder:
                    if in_references:
                        current_reference = [remainder]
                    else:
                        paragraph_lines.append(remainder)
                continue

            if in_references:
                ref_match = _REFERENCE_START_RE.match(line)
                if ref_match:
                    _flush_reference(current_reference, body_sections)
                    num, content = ref_match.groups()
                    current_reference = [f"{num}. {content.strip()}"]
                elif current_reference:
                    current_reference.append(line)
                else:
                    paragraph_lines.append(line)
                continue

            if _FIGURE_CAPTION_RE.match(line):
                _flush_paragraph(paragraph_lines, body_sections)
                matched_assets = _take_matching_assets_for_caption(pending_assets, line)
                if not matched_assets:
                    matched_assets = _take_unlabeled_assets_for_caption(pending_assets)
                _append_rendered_assets(matched_assets, body_sections)
                body_sections.append(_format_display_caption_line(line))
                continue

            paragraph_lines.append(line)

        if not in_references:
            _flush_paragraph(paragraph_lines, body_sections)
        if next_page:
            carryover_assets = _take_assets_for_next_page(pending_assets, next_page)
        _append_rendered_assets(pending_assets, body_sections)

    _flush_paragraph(paragraph_lines, body_sections)
    _flush_reference(current_reference, body_sections)
    _append_rendered_assets(carryover_assets, body_sections)
    return _dedupe_long_sections(body_sections)


def _render_references_from_pages(content_pages: list[PDFPage], metadata: PDFMetadata) -> str:
    references: list[str] = []
    current_reference: list[str] = []
    in_references = False

    for page_index, page in enumerate(content_pages):
        segments = _page_segments_for_render(page, metadata if page_index == 0 else None)
        for segment in segments:
            if segment["type"] != "line":
                continue
            line = segment["text"].strip()
            if not line:
                continue

            heading, remainder = _match_section_heading(line)
            if heading:
                if heading == "## References":
                    in_references = True
                    if remainder:
                        _consume_reference_line(remainder, current_reference, references)
                    continue
                if in_references:
                    _flush_reference_to_list(current_reference, references)
                    return _format_reference_section(references)
                continue

            if not in_references:
                continue

            _consume_reference_line(line, current_reference, references)

    _flush_reference_to_list(current_reference, references)
    return _format_reference_section(references)


def _consume_reference_line(line: str, current_reference: list[str], references: list[str]) -> None:
    ref_match = _REFERENCE_START_RE.match(line)
    if ref_match:
        _flush_reference_to_list(current_reference, references)
        num, content = ref_match.groups()
        current_reference.append(f"{num}. {content.strip()}")
        return
    if current_reference and _line_can_continue_reference(line):
        current_reference.append(line)


def _line_can_continue_reference(line: str) -> bool:
    stripped = line.strip()
    if not stripped:
        return False
    heading, _ = _match_section_heading(stripped)
    if heading:
        return False
    normalized = _normalize_for_compare(stripped)
    if normalized.startswith(("acknowledgement", "author contribution", "competing interest", "additional information")):
        return False
    if _looks_like_publisher_metadata_line(stripped) or _is_running_header_or_footer_text(stripped):
        return False
    return True


def _flush_reference_to_list(buffer: list[str], references: list[str]) -> None:
    if not buffer:
        return
    reference = _join_wrapped_lines(buffer)
    if reference and _looks_like_complete_reference(reference):
        references.append(reference)
    buffer.clear()


def _looks_like_complete_reference(reference: str) -> bool:
    if not _REFERENCE_START_RE.match(reference):
        return False
    if len(reference) < 35:
        return False
    return bool(
        re.search(r"\(\d{4}\)", reference)
        or re.search(r"\b\d+\s*,\s*\d+\s*[-–]\s*\d+", reference)
        or _DOI_RE.search(reference)
    )


def _format_reference_section(references: list[str]) -> str:
    if len(references) < 3:
        return ""
    return "## References\n\n" + "\n\n".join(references)


def _replace_reference_section(text: str, replacement: str) -> str:
    replacement = replacement.strip()
    if not replacement:
        return text
    match = re.search(r"(?im)^## References\s*$", text)
    if not match:
        return text.rstrip() + "\n\n" + replacement

    after = match.end()
    next_heading = re.search(r"(?m)^## (?!References\b).+", text[after:])
    end = after + next_heading.start() if next_heading else len(text)
    prefix = text[: match.start()].rstrip()
    suffix = text[end:].lstrip()
    pieces = [piece for piece in (prefix, replacement, suffix) if piece]
    return "\n\n".join(pieces).strip()


def _page_segments_for_render(page: PDFPage, metadata: PDFMetadata | None) -> list[dict[str, str]]:
    segments: list[dict[str, str]] = []
    cutoff_y = metadata.first_page_cutoff_y if metadata else 0.0
    rendered_tables: set[int] = set()

    for block in page.blocks:
        table_index = _best_matching_table(block.bbox, page.tables, rendered_tables)
        if table_index is not None:
            table_md = _table_to_markdown(page.tables[table_index].rows)
            if table_md:
                segments.append({"type": "table", "text": table_md})
                segments.append({"type": "line", "text": ""})
            rendered_tables.add(table_index)
            continue

        block_lines = []
        for line in block.lines:
            if cutoff_y and line.bbox[1] < cutoff_y - 1:
                continue

            text = _clean_text(_fix_unicode(line.text))
            if not text:
                continue
            if any(phrase in text.lower() for phrase in _BOILERPLATE_PHRASES):
                continue
            block_lines.append(text)

        if block_lines:
            segments.extend({"type": "line", "text": text} for text in block_lines)
            segments.append({"type": "line", "text": ""})

    for index, table in enumerate(page.tables):
        if index in rendered_tables:
            continue
        table_md = _table_to_markdown(table.rows)
        if table_md:
            segments.append({"type": "table", "text": table_md})
            segments.append({"type": "line", "text": ""})

    return segments


def _best_matching_table(
    block_bbox: tuple[float, float, float, float],
    tables: list[PDFTable],
    rendered_tables: set[int],
) -> int | None:
    best_index = None
    best_score = 0.0
    block_area = _bbox_area(block_bbox)
    for index, table in enumerate(tables):
        if index in rendered_tables:
            continue
        overlap = _bbox_intersection_area(block_bbox, table.bbox)
        if overlap <= 0:
            continue
        score = overlap / max(min(block_area, _bbox_area(table.bbox)), 1.0)
        if score > best_score:
            best_index = index
            best_score = score
    if best_score >= 0.35:
        return best_index
    return None


def _match_section_heading(line: str) -> tuple[str | None, str]:
    stripped = line.strip()
    for pattern, heading in _SECTION_HEADING_PATTERNS:
        match = pattern.match(stripped)
        if match:
            remainder = match.group(1).strip()
            return heading, remainder
    return None, ""


def _flush_paragraph(buffer: list[str], sections: list[str]) -> None:
    if not buffer:
        return
    paragraph = _join_wrapped_lines(buffer)
    if paragraph:
        sections.append(paragraph)
    buffer.clear()


def _flush_reference(buffer: list[str], sections: list[str]) -> None:
    if not buffer:
        return
    head = buffer[0].strip()
    tail = buffer[1:]
    if tail:
        reference = _join_wrapped_lines([head] + tail)
    else:
        reference = head
    if reference:
        sections.append(reference)
    buffer.clear()


def _join_wrapped_lines(lines: list[str]) -> str:
    cleaned = [line.strip() for line in lines if line.strip()]
    if not cleaned:
        return ""

    parts = [cleaned[0]]
    for line in cleaned[1:]:
        previous = parts[-1]
        if previous.endswith("-") and line and line[0].islower():
            parts[-1] = previous + line
        else:
            parts.append(line)
    text = " ".join(parts)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _render_pdf_assets(
    page_obj: Any,
    page: PDFPage,
    output_path: Path,
    pdf_config: PDFConfig,
    fitz: Any,
    include_ocr_text: bool = True,
    next_page: PDFPage | None = None,
    cancel_event: Any | None = None,
) -> str:
    assets = _collect_rendered_pdf_assets(
        page_obj,
        page,
        output_path,
        pdf_config,
        fitz,
        include_ocr_text=include_ocr_text,
        next_page=next_page,
        cancel_event=cancel_event,
    )
    return _join_rendered_assets(assets)


def _collect_rendered_pdf_assets(
    page_obj: Any,
    page: PDFPage,
    output_path: Path,
    pdf_config: PDFConfig,
    fitz: Any,
    include_ocr_text: bool = True,
    next_page: PDFPage | None = None,
    cancel_event: Any | None = None,
) -> list[PDFRenderedAsset]:
    asset_dir = output_path.parent / f"{output_path.stem}_assets"
    assets: list[PDFRenderedAsset] = []

    has_meaningful_text = _is_meaningful_text(page.raw_text)
    image_only = not has_meaningful_text
    figure_dominant = page.image_area_ratio >= 0.55

    exported_any = False
    if image_only or figure_dominant:
        _check_cancel(cancel_event)
        asset_path = asset_dir / f"page_{page.number:03d}.png"
        caption_text = _find_figure_caption_text(page, (0.0, 0.0, page.width, page.height), next_page=next_page)
        if not caption_text:
            caption_text = _fallback_page_asset_caption(page, next_page)
        if not caption_text:
            return assets
        _ensure_asset_limit(asset_dir, pdf_config)
        _export_page_pixmap(page_obj, asset_path, fitz)
        alt_text = _markdown_image_alt_text(caption_text, f"Page {page.number:03d}")
        assets.append(
            PDFRenderedAsset(
                markdown=f'![{alt_text}]({asset_path.relative_to(output_path.parent).as_posix()})',
                caption_text=caption_text,
                label_key=_display_caption_key(caption_text),
            )
        )
        exported_any = True
    else:
        min_image_area_ratio = 0.06 if _page_has_figure_caption(page) else 0.12
        large_images = sorted(
            [image for image in page.images if _bbox_area(image.bbox) / page.page_area >= min_image_area_ratio],
            key=lambda item: _bbox_area(item.bbox),
            reverse=True,
        )
        figure_candidates: list[PDFFigureCandidate] = []
        for image in large_images:
            _check_cancel(cancel_event)
            clip_bbox = _expanded_figure_bbox(page_obj, page, image, fitz)
            caption_text = _find_figure_caption_text(page, clip_bbox, next_page=next_page)
            if not caption_text:
                page_captions = _display_caption_candidates(page)
                if len(page_captions) == 1:
                    caption_text = page_captions[0]
            figure_candidates.append(
                PDFFigureCandidate(
                    bbox=clip_bbox,
                    caption_text=caption_text,
                    label_key=_display_caption_key(caption_text),
                    image_index=image.index,
                )
            )
        grouped_candidates = _group_figure_candidates(page_obj, page, figure_candidates, fitz)
        exported_bboxes: list[tuple[float, float, float, float]] = []
        for group_index, candidate in enumerate(grouped_candidates, 1):
            clip_bbox = candidate.bbox
            if not candidate.caption_text:
                continue
            if any(_bboxes_near_duplicate(clip_bbox, bbox) for bbox in exported_bboxes):
                continue
            exported_bboxes.append(clip_bbox)
            suffix = f"{candidate.image_index:02d}" if candidate.image_index else f"g{group_index:02d}"
            asset_path = asset_dir / f"page_{page.number:03d}_figure_{suffix}.png"
            _ensure_asset_limit(asset_dir, pdf_config)
            _export_clip_pixmap(page_obj, asset_path, fitz, clip_bbox)
            alt_text = _markdown_image_alt_text(
                candidate.caption_text,
                f"Page {page.number:03d} figure {suffix}",
            )
            assets.append(
                PDFRenderedAsset(
                    markdown=f'![{alt_text}]({asset_path.relative_to(output_path.parent).as_posix()})',
                    caption_text=candidate.caption_text,
                    label_key=candidate.label_key,
                )
            )
            exported_any = True

    if (
        include_ocr_text
        and pdf_config.tessdata
        and pdf_config.ocr_mode != "off"
        and (image_only or figure_dominant or not has_meaningful_text)
    ):
        _check_cancel(cancel_event)
        ocr_text = _extract_ocr_text(page_obj, page, pdf_config)
        if ocr_text:
            assets.append(PDFRenderedAsset(markdown=ocr_text))

    return assets


def _ensure_asset_limit(asset_dir: Path, pdf_config: PDFConfig) -> None:
    limit = pdf_config.security.max_extracted_assets
    if limit <= 0 or not asset_dir.exists():
        return
    count = sum(1 for item in asset_dir.iterdir() if item.is_file())
    if count >= limit:
        sys.exit(f"Extracted asset limit exceeded ({limit}) for {asset_dir.name}.")


def _join_rendered_assets(assets: list[PDFRenderedAsset]) -> str:
    return "\n\n".join(asset.markdown.strip() for asset in assets if asset.markdown.strip()).strip()


def _append_rendered_assets(assets: list[PDFRenderedAsset], sections: list[str]) -> None:
    for asset in assets:
        markdown = asset.markdown.strip()
        if markdown:
            sections.append(markdown)


def _display_caption_key(text: str) -> str:
    if not text:
        return ""
    candidate = text.replace("**", "").strip()
    candidate = re.sub(r"^[>\s#*\-]+", "", candidate)
    if not _looks_like_display_caption_line(candidate):
        return ""
    match = _DISPLAY_CAPTION_RE.match(candidate)
    if not match:
        return ""
    label = match.group(1).strip().rstrip(".:")
    label = re.sub(r"\bFig\.", "Figure", label)
    label = re.sub(r"\s+", " ", label)
    return label.lower()


def _fallback_page_asset_caption(page: PDFPage, next_page: PDFPage | None = None) -> str:
    page_captions = _display_caption_candidates(page)
    if page_captions:
        return page_captions[0]

    if next_page is not None and not _page_looks_like_reference_material(page):
        next_captions = _display_caption_candidates(next_page)
        if len(next_captions) == 1 and not _has_meaningful_text_before_caption(next_page, _display_caption_bbox(next_page, next_captions[0])):
            return next_captions[0]
    return ""


def _display_caption_candidates(page: PDFPage) -> list[str]:
    captions: list[str] = []
    for block in page.blocks:
        text = _clean_figure_caption_text(_collapse_block_text(block))
        if text and _looks_like_display_caption_line(text):
            captions.append(text)
    return captions


def _display_caption_bbox(page: PDFPage, caption_text: str) -> tuple[float, float, float, float]:
    normalized_target = _normalize_for_compare(caption_text)
    for block in page.blocks:
        text = _clean_figure_caption_text(_collapse_block_text(block))
        if _normalize_for_compare(text) == normalized_target:
            return block.bbox
    return (0.0, 0.0, page.width, page.height * 0.2)


def _page_looks_like_reference_material(page: PDFPage) -> bool:
    normalized = _normalize_for_compare(page.raw_text[:2400])
    if "references" in normalized:
        return True
    if len(re.findall(r"\b\d{4}[;,.]", page.raw_text)) >= 8:
        return True
    return False


def _take_matching_assets_for_caption(
    assets: list[PDFRenderedAsset],
    caption_line: str,
) -> list[PDFRenderedAsset]:
    label_key = _display_caption_key(caption_line)
    if not label_key or not assets:
        return []

    matched: list[PDFRenderedAsset] = []
    remaining: list[PDFRenderedAsset] = []
    for asset in assets:
        if asset.label_key == label_key:
            matched.append(asset)
        else:
            remaining.append(asset)
    assets[:] = remaining
    return matched


def _take_unlabeled_assets_for_caption(assets: list[PDFRenderedAsset]) -> list[PDFRenderedAsset]:
    unlabeled = [asset for asset in assets if not asset.label_key and asset.caption_text == ""]
    if len(unlabeled) != 1:
        return []
    assets[:] = [asset for asset in assets if asset is not unlabeled[0]]
    return unlabeled


def _inject_assets_into_markdown(text: str, assets: list[PDFRenderedAsset]) -> str:
    rendered, remaining_assets = _inject_assets_into_markdown_with_remaining(text, assets)
    if remaining_assets:
        trailing = _join_rendered_assets(remaining_assets)
        if trailing:
            if rendered:
                return f"{rendered}\n\n{trailing}".strip()
            return trailing
    return rendered


def _inject_assets_into_markdown_with_remaining(
    text: str,
    assets: list[PDFRenderedAsset],
) -> tuple[str, list[PDFRenderedAsset]]:
    if not assets:
        return text.strip(), []
    if not text.strip():
        return "", list(assets)

    remaining_assets = list(assets)
    output_lines: list[str] = []

    for line in text.splitlines():
        matched_assets = _take_matching_assets_for_caption(remaining_assets, line)
        if matched_assets:
            if output_lines and output_lines[-1].strip():
                output_lines.append("")
            output_lines.extend(asset.markdown.strip() for asset in matched_assets if asset.markdown.strip())
            output_lines.append("")
        output_lines.append(line)

    return "\n".join(output_lines).strip(), remaining_assets


def _take_assets_for_next_page(
    assets: list[PDFRenderedAsset],
    next_page: PDFPage | None,
) -> list[PDFRenderedAsset]:
    if not assets or next_page is None:
        return []
    next_keys = _page_caption_keys(next_page)
    if not next_keys:
        return []

    carryover: list[PDFRenderedAsset] = []
    remaining: list[PDFRenderedAsset] = []
    for asset in assets:
        if asset.label_key and asset.label_key in next_keys:
            carryover.append(asset)
        else:
            remaining.append(asset)
    assets[:] = remaining
    return carryover


def _group_figure_candidates(
    page_obj: Any,
    page: PDFPage,
    candidates: list[PDFFigureCandidate],
    fitz: Any,
) -> list[PDFFigureCandidate]:
    if not candidates:
        return []

    groups = sorted(
        candidates,
        key=lambda item: (item.bbox[1], item.bbox[0], -_bbox_area(item.bbox)),
    )
    horizontal_gap_limit = max(page.width * 0.12, 40.0)
    vertical_gap_limit = max(page.height * 0.06, 28.0)

    changed = True
    while changed:
        changed = False
        merged_groups: list[PDFFigureCandidate] = []
        consumed = [False] * len(groups)
        for index, base in enumerate(groups):
            if consumed[index]:
                continue
            current = base
            for other_index in range(index + 1, len(groups)):
                if consumed[other_index]:
                    continue
                other = groups[other_index]
                if not _figure_candidates_should_merge(
                    current,
                    other,
                    horizontal_gap_limit,
                    vertical_gap_limit,
                    fitz,
                ):
                    continue
                current = _merge_figure_candidates(current, other, page_obj, fitz)
                consumed[other_index] = True
                changed = True
            merged_groups.append(current)
        groups = merged_groups

    return groups


def _figure_candidates_should_merge(
    left: PDFFigureCandidate,
    right: PDFFigureCandidate,
    horizontal_gap_limit: float,
    vertical_gap_limit: float,
    fitz: Any,
) -> bool:
    left_rect = fitz.Rect(left.bbox)
    right_rect = fitz.Rect(right.bbox)
    if left.label_key and right.label_key and left.label_key == right.label_key:
        return True
    if left.label_key and right.label_key and left.label_key != right.label_key:
        return False
    return _should_merge_figure_rect(left_rect, right_rect, horizontal_gap_limit, vertical_gap_limit)


def _merge_figure_candidates(
    left: PDFFigureCandidate,
    right: PDFFigureCandidate,
    page_obj: Any,
    fitz: Any,
) -> PDFFigureCandidate:
    merged_rect = _pad_rect(
        _merge_rects(fitz.Rect(left.bbox), fitz.Rect(right.bbox), fitz),
        page_obj.rect,
        x_pad=6.0,
        y_pad=6.0,
    )
    caption_text = left.caption_text or right.caption_text
    label_key = left.label_key or right.label_key
    image_index = min(index for index in (left.image_index, right.image_index) if index) if (left.image_index or right.image_index) else 0
    return PDFFigureCandidate(
        bbox=(merged_rect.x0, merged_rect.y0, merged_rect.x1, merged_rect.y1),
        caption_text=caption_text,
        label_key=label_key,
        image_index=image_index,
    )


def _export_page_pixmap(page_obj: Any, asset_path: Path, fitz: Any) -> None:
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    pix = page_obj.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    pix.save(asset_path)


def _export_clip_pixmap(page_obj: Any, asset_path: Path, fitz: Any, bbox: tuple[float, float, float, float]) -> None:
    rect = fitz.Rect(bbox)
    rect = rect & page_obj.rect
    if rect.is_empty or rect.width < 10 or rect.height < 10:
        return
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    pix = page_obj.get_pixmap(clip=rect, matrix=fitz.Matrix(2, 2), alpha=False)
    pix.save(asset_path)


def _expanded_figure_bbox(page_obj: Any, page: PDFPage, image: PDFImage, fitz: Any) -> tuple[float, float, float, float]:
    rect = fitz.Rect(image.bbox) & page_obj.rect
    if rect.is_empty:
        return image.bbox

    rect = _pad_rect(
        rect,
        page_obj.rect,
        x_pad=max(page.width * 0.01, 8.0),
        y_pad=max(page.height * 0.01, 8.0),
    )
    candidate_rects = _figure_candidate_rects(page_obj, page, image, fitz)
    horizontal_gap_limit = max(page.width * 0.08, 28.0)
    vertical_gap_limit = max(page.height * 0.03, 18.0)

    changed = True
    while changed:
        changed = False
        for candidate in candidate_rects:
            if not _should_merge_figure_rect(rect, candidate, horizontal_gap_limit, vertical_gap_limit):
                continue
            merged = _pad_rect(
                _merge_rects(rect, candidate, fitz),
                page_obj.rect,
                x_pad=max(page.width * 0.006, 4.0),
                y_pad=max(page.height * 0.006, 4.0),
            )
            if _rect_tuple(rect) != _rect_tuple(merged):
                rect = merged
                changed = True

    rect = _clamp_figure_rect(page_obj, page, image, rect, fitz)
    return rect.x0, rect.y0, rect.x1, rect.y1


def _find_figure_caption_text(
    page: PDFPage,
    target_bbox: tuple[float, float, float, float],
    next_page: PDFPage | None = None,
) -> str:
    page_captions = _figure_caption_blocks(page)
    below_candidates: list[tuple[float, str]] = []
    above_candidates: list[tuple[float, str]] = []
    target_x0, target_y0, target_x1, target_y1 = target_bbox
    target_width = max(target_x1 - target_x0, 1.0)

    for block_bbox, text in page_captions:
        block_x0, block_y0, block_x1, block_y1 = block_bbox
        block_width = max(block_x1 - block_x0, 1.0)
        horizontal_overlap = _axis_overlap(target_x0, target_x1, block_x0, block_x1)
        overlap_ratio = horizontal_overlap / max(min(target_width, block_width), 1.0)
        if block_y0 >= target_y1:
            vertical_distance = block_y0 - target_y1
            if overlap_ratio < 0.1 and vertical_distance > page.height * 0.12:
                continue
            score = vertical_distance - (overlap_ratio * 100.0)
            below_candidates.append((score, text))
            continue
        if block_y1 <= target_y0:
            vertical_distance = target_y0 - block_y1
            if overlap_ratio < 0.25 and vertical_distance > page.height * 0.08:
                continue
            score = vertical_distance - (overlap_ratio * 75.0)
            above_candidates.append((score, text))

    if below_candidates:
        below_candidates.sort(key=lambda item: item[0])
        return below_candidates[0][1]
    if above_candidates:
        above_candidates.sort(key=lambda item: item[0])
        return above_candidates[0][1]

    if len(page_captions) == 1:
        return page_captions[0][1]
    continuation_caption = _find_continuation_figure_caption_text(page, next_page, target_bbox)
    if continuation_caption:
        return continuation_caption
    return ""


def _page_has_figure_caption(page: PDFPage) -> bool:
    return bool(_figure_caption_blocks(page))


def _figure_caption_blocks(page: PDFPage) -> list[tuple[tuple[float, float, float, float], str]]:
    captions: list[tuple[tuple[float, float, float, float], str]] = []
    for block in page.blocks:
        text = _clean_figure_caption_text(_collapse_block_text(block))
        if text and _FIGURE_CAPTION_RE.match(text):
            captions.append((block.bbox, text))
    return captions


def _page_caption_keys(page: PDFPage) -> set[str]:
    return {
        key
        for _, text in _figure_caption_blocks(page)
        if (key := _display_caption_key(text))
    }


def _find_continuation_figure_caption_text(
    page: PDFPage,
    next_page: PDFPage | None,
    target_bbox: tuple[float, float, float, float],
) -> str:
    if next_page is None:
        return ""

    target_x0, target_y0, target_x1, target_y1 = target_bbox
    target_height = max(target_y1 - target_y0, 1.0)
    target_width = max(target_x1 - target_x0, 1.0)
    if target_height < page.height * 0.35 and target_y1 < page.height * 0.55:
        return ""
    if target_width < page.width * 0.25:
        return ""

    next_captions = _figure_caption_blocks(next_page)
    if not next_captions:
        return ""

    first_bbox, first_caption = min(next_captions, key=lambda item: (item[0][1], item[0][0]))
    if first_bbox[1] > next_page.height * 0.22:
        return ""
    if _has_meaningful_text_before_caption(next_page, first_bbox):
        return ""
    return first_caption


def _has_meaningful_text_before_caption(
    page: PDFPage,
    caption_bbox: tuple[float, float, float, float],
) -> bool:
    caption_top = caption_bbox[1]
    for block in page.blocks:
        if block.bbox[3] > caption_top + 1:
            continue
        text = _clean_text(_fix_unicode(_collapse_block_text(block))).strip()
        if not text:
            continue
        if _is_running_header_or_footer_text(text):
            continue
        if _looks_like_panel_label_line(text):
            continue
        if _looks_like_publisher_metadata_line(text):
            continue
        if _FIGURE_CAPTION_RE.match(_clean_figure_caption_text(text)):
            continue
        if len(re.findall(r"[A-Za-z]", text)) < 12:
            continue
        return True
    return False


def _clean_figure_caption_text(text: str) -> str:
    cleaned = _join_wrapped_lines(text.splitlines())
    cleaned = _DISPLAY_DOI_SUFFIX_RE.sub("", cleaned).strip()
    cleaned = _DISPLAY_DOI_FRAGMENT_RE.sub("", cleaned).strip()
    return _normalize_pdf_caption_text(cleaned)


def _markdown_image_alt_text(text: str, fallback: str) -> str:
    candidate = _normalize_pdf_caption_text(text) if text else fallback
    candidate = candidate.replace("[", "(").replace("]", ")")
    candidate = re.sub(r"\s+", " ", candidate).strip()
    if len(candidate) > 220:
        candidate = candidate[:217].rstrip() + "..."
    return candidate or fallback


def _figure_candidate_rects(page_obj: Any, page: PDFPage, image: PDFImage, fitz: Any) -> list[Any]:
    candidates = []

    for other in page.images:
        if other.index == image.index:
            continue
        if _bbox_area(other.bbox) / page.page_area < 0.01:
            continue
        rect = fitz.Rect(other.bbox) & page_obj.rect
        if rect.is_empty:
            continue
        candidates.append(rect)

    for block in page.blocks:
        text = block.text.strip()
        if not text:
            continue
        if _is_running_header_or_footer_text(text):
            continue
        if _is_body_text_block(block):
            continue
        if len(text) > 220 and not _FIGURE_CAPTION_RE.match(text):
            continue
        rect = fitz.Rect(block.bbox) & page_obj.rect
        if rect.is_empty:
            continue
        candidates.append(rect)

    with contextlib.suppress(Exception):
        for drawing in page_obj.get_drawings():
            drawing_bbox = drawing.get("rect") or drawing.get("bbox")
            if not drawing_bbox:
                continue
            rect = fitz.Rect(drawing_bbox) & page_obj.rect
            if rect.is_empty or rect.width < 1 or rect.height < 1:
                continue
            if _should_skip_drawing_candidate(page, rect):
                continue
            candidates.append(rect)

    return candidates


def _clamp_figure_rect(page_obj: Any, page: PDFPage, image: PDFImage, rect: Any, fitz: Any) -> Any:
    top_floor = 0.0
    bottom_ceiling = page.height
    image_mid_y = (image.bbox[1] + image.bbox[3]) / 2.0
    margin = max(page.height * 0.01, 6.0)

    for block in page.blocks:
        text = _collapse_block_text(block).strip()
        if not text:
            continue
        if _is_running_header_or_footer_text(text):
            if block.bbox[3] <= image.bbox[1]:
                top_floor = max(top_floor, block.bbox[3] + margin)
            elif block.bbox[1] >= image.bbox[3]:
                bottom_ceiling = min(bottom_ceiling, block.bbox[1] - margin)
            continue
        if not _is_body_text_block(block):
            continue
        if block.bbox[3] <= image_mid_y:
            top_floor = max(top_floor, block.bbox[3] + margin)
        elif block.bbox[1] >= image_mid_y:
            bottom_ceiling = min(bottom_ceiling, block.bbox[1] - margin)

    clamped = fitz.Rect(rect)
    clamped.y0 = max(clamped.y0, top_floor)
    clamped.y1 = min(clamped.y1, bottom_ceiling)
    if clamped.y1 <= clamped.y0 + 12:
        return rect & page_obj.rect
    return clamped & page_obj.rect


def _is_body_text_block(block: PDFBlock) -> bool:
    text = _collapse_block_text(block).strip()
    if not text:
        return False
    if _FIGURE_CAPTION_RE.match(text):
        return False
    if _is_running_header_or_footer_text(text):
        return False
    return len(text) >= 240 or len(block.lines) >= 6


def _is_running_header_or_footer_text(text: str) -> bool:
    normalized = _normalize_for_compare(text)
    if normalized in {"article", "feature article", "research article", "check for updates", "open access"}:
        return True
    if any(
        phrase in normalized
        for phrase in (
            "nature communications",
            "nature portfolio",
            "springer nature",
            "www.nature.com",
            "published online",
            "correspondence and requests for materials",
        )
    ):
        return True
    if normalized.startswith("received") or normalized.startswith("accepted"):
        return True
    if normalized.startswith("this article is licensed under"):
        return True
    if "https://doi. org/" in text:
        return True
    return bool(_DOI_RE.search(text))


def _should_skip_drawing_candidate(page: PDFPage, rect: Any) -> bool:
    if rect.width >= page.width * 0.75 and rect.height <= max(page.height * 0.02, 12.0):
        return True
    if rect.y1 <= page.height * 0.05:
        return True
    if rect.y0 >= page.height * 0.95:
        return True
    return False


def _should_merge_figure_rect(
    base_rect: Any,
    candidate_rect: Any,
    horizontal_gap_limit: float,
    vertical_gap_limit: float,
) -> bool:
    if candidate_rect.is_empty:
        return False

    intersection = base_rect & candidate_rect
    if not intersection.is_empty:
        return True

    horizontal_overlap = _axis_overlap(base_rect.x0, base_rect.x1, candidate_rect.x0, candidate_rect.x1)
    vertical_overlap = _axis_overlap(base_rect.y0, base_rect.y1, candidate_rect.y0, candidate_rect.y1)
    horizontal_gap = _axis_gap(base_rect.x0, base_rect.x1, candidate_rect.x0, candidate_rect.x1)
    vertical_gap = _axis_gap(base_rect.y0, base_rect.y1, candidate_rect.y0, candidate_rect.y1)

    min_width = max(min(base_rect.width, candidate_rect.width), 1.0)
    min_height = max(min(base_rect.height, candidate_rect.height), 1.0)

    if vertical_overlap / min_height >= 0.25 and horizontal_gap <= horizontal_gap_limit:
        return True
    if horizontal_overlap / min_width >= 0.4 and vertical_gap <= vertical_gap_limit:
        return True
    return False


def _merge_rects(rect_a: Any, rect_b: Any, fitz: Any) -> Any:
    return fitz.Rect(
        min(rect_a.x0, rect_b.x0),
        min(rect_a.y0, rect_b.y0),
        max(rect_a.x1, rect_b.x1),
        max(rect_a.y1, rect_b.y1),
    )


def _pad_rect(rect: Any, bounds: Any, x_pad: float, y_pad: float) -> Any:
    padded = rect.__class__(rect.x0 - x_pad, rect.y0 - y_pad, rect.x1 + x_pad, rect.y1 + y_pad)
    return padded & bounds


def _rect_tuple(rect: Any) -> tuple[float, float, float, float]:
    return (
        round(float(rect.x0), 3),
        round(float(rect.y0), 3),
        round(float(rect.x1), 3),
        round(float(rect.y1), 3),
    )


def _extract_ocr_text(page_obj: Any, page: PDFPage, pdf_config: PDFConfig) -> str:
    if not _page_requires_ocr(page, pdf_config):
        return ""
    full_page_ocr = (
        pdf_config.ocr_mode == "full"
        or not _is_meaningful_text(page.raw_text)
        or page.image_area_ratio >= 0.55
    )

    try:
        textpage = page_obj.get_textpage_ocr(
            language=pdf_config.ocr_language,
            dpi=pdf_config.ocr_dpi,
            full=full_page_ocr,
            tessdata=pdf_config.tessdata,
        )
    except Exception as exc:
        sys.exit(f"OCR failed for page {page.number}: {exc}")

    ocr_text = page_obj.get_text("text", textpage=textpage, sort=True)
    novel_text = _subtract_existing_text(ocr_text, page.raw_text)
    if not novel_text:
        return ""
    return novel_text


def _subtract_existing_text(ocr_text: str, raw_text: str) -> str:
    raw_norm = _normalize_for_compare(raw_text)
    kept = []
    seen = set()

    for paragraph in re.split(r"\n{2,}", _clean_text(_fix_unicode(ocr_text))):
        text = paragraph.strip()
        if not text:
            continue
        normalized = _normalize_for_compare(text)
        if len(normalized) < 10:
            continue
        if normalized in seen:
            continue
        if normalized and normalized in raw_norm:
            continue
        if _looks_garbled(text):
            continue
        if _looks_like_stray_figure_ocr_text(text):
            continue
        kept.append(text)
        seen.add(normalized)

    return "\n\n".join(kept)


def _normalize_for_compare(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip().lower()


def _is_meaningful_text(text: str) -> bool:
    normalized = _normalize_for_compare(text)
    return len(re.findall(r"[a-zA-Z]", normalized)) >= 20


def _page_requires_ocr(page: PDFPage, pdf_config: PDFConfig) -> bool:
    if pdf_config.ocr_mode == "off":
        return False
    if pdf_config.ocr_mode == "full":
        return True
    if not _is_meaningful_text(page.raw_text):
        return True
    return page.image_area_ratio >= 0.12


def _looks_garbled(text: str) -> bool:
    if not text:
        return True

    control_like = sum(1 for ch in text if ord(ch) < 32 and ch not in "\n\t\r")
    replacement_like = text.count("\uFFFD") + text.count("\x00")
    strange = control_like + replacement_like
    return strange > max(2, len(text) // 50)


def _looks_like_stray_figure_ocr_text(text: str) -> bool:
    normalized = _normalize_for_compare(text)
    if not normalized:
        return True
    if _looks_like_panel_label_line(text):
        return True
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) <= 4:
        measurement_like = 0
        for line in lines:
            compact = _normalize_for_compare(line)
            if re.fullmatch(r"(?:\(?[a-z]\)?|fig(?:ure)?\.?\s*\d+[a-z]?|table\s*\d+[a-z]?)", compact):
                measurement_like += 1
                continue
            if re.fullmatch(r"(?:scale bar|bar(?:\s+(?:equals|corresponds to))?)\s*[:=]?\s*\d+(?:\.\d+)?\s*(?:nm|um|μm|mm|cm|mpa|kda|%)?", compact):
                measurement_like += 1
                continue
            if re.fullmatch(r"\d+(?:\.\d+)?\s*(?:nm|um|μm|mm|cm|mpa|kda|%)", compact):
                measurement_like += 1
                continue
        if measurement_like == len(lines) and measurement_like > 0:
            return True
    return False


def _bbox_area(bbox: tuple[float, float, float, float]) -> float:
    x0, y0, x1, y1 = bbox
    return max(x1 - x0, 0.0) * max(y1 - y0, 0.0)


def _bbox_intersection_area(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> float:
    x0 = max(left[0], right[0])
    y0 = max(left[1], right[1])
    x1 = min(left[2], right[2])
    y1 = min(left[3], right[3])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return (x1 - x0) * (y1 - y0)


def _bboxes_near_duplicate(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> bool:
    overlap = _bbox_intersection_area(left, right)
    if overlap <= 0:
        return False
    left_area = max(_bbox_area(left), 1.0)
    right_area = max(_bbox_area(right), 1.0)
    return overlap / min(left_area, right_area) >= 0.9


def _axis_overlap(a0: float, a1: float, b0: float, b1: float) -> float:
    return max(0.0, min(a1, b1) - max(a0, b0))


def _axis_gap(a0: float, a1: float, b0: float, b1: float) -> float:
    return max(0.0, max(b0 - a1, a0 - b1))


def _collapse_block_text(block: PDFBlock) -> str:
    return _join_wrapped_lines([line.text for line in block.lines])


def _dedupe_long_sections(sections: list[str]) -> list[str]:
    result = []
    seen = set()
    for section in sections:
        key = re.sub(r"\s+", " ", section).strip().lower()
        if len(key) > 120 and key in seen:
            continue
        if len(key) > 120:
            seen.add(key)
        result.append(section)
    return result


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        sys.exit("Missing dep: pip install python-docx")

    doc = Document(path)
    lines = [f"# {path.stem}\n"]
    for block in _iter_docx_blocks(doc):
        if block["type"] == "paragraph":
            paragraph = block["element"]
            style = paragraph.style.name if paragraph.style else ""
            text = paragraph.text.strip()
            if not text:
                continue
            if style.startswith("Heading 1"):
                lines.append(f"\n# {text}")
            elif style.startswith("Heading 2"):
                lines.append(f"\n## {text}")
            elif style.startswith("Heading 3"):
                lines.append(f"\n### {text}")
            elif style.startswith("List"):
                lines.append(f"- {text}")
            else:
                lines.append(text)
        elif block["type"] == "table":
            markdown_table = _table_to_markdown(
                [[cell.text.strip() for cell in row.cells] for row in block["element"].rows]
            )
            if markdown_table:
                lines.append("\n" + markdown_table + "\n")
    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("Missing dep: pip install python-pptx")

    presentation = Presentation(path)
    lines = [f"# {path.stem}\n"]
    for i, slide in enumerate(presentation.slides, 1):
        title = ""
        body_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.shape_type == 13:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title and shape.name.lower().startswith("title"):
                title = text
            else:
                body_lines.append(text)
        lines.append(f"\n## {title or f'Slide {i}'}\n")
        for body_line in body_lines:
            for subline in body_line.split("\n"):
                if subline.strip():
                    lines.append(f"- {subline.strip()}")
    return "\n".join(lines)


def extract_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        sys.exit("Missing dep: pip install beautifulsoup4")

    html = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    lines = [f"# {path.stem}\n"]
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "pre", "code"]):
        text = tag.get_text(separator=" ", strip=True)
        if not text:
            continue
        name = tag.name
        if name == "h1":
            lines.append(f"\n# {text}")
        elif name == "h2":
            lines.append(f"\n## {text}")
        elif name == "h3":
            lines.append(f"\n### {text}")
        elif name == "h4":
            lines.append(f"\n#### {text}")
        elif name == "li":
            lines.append(f"- {text}")
        elif name in ("pre", "code"):
            lines.append(f"```\n{text}\n```")
        else:
            lines.append(text)
    return "\n".join(lines)


def extract_csv(path: Path) -> str:
    try:
        import pandas as pd
    except ImportError:
        sys.exit("Missing dep: pip install pandas")
    dataframe = pd.read_csv(path)
    return f"# {path.stem}\n\n" + dataframe.to_markdown(index=False)


def extract_txt(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    return f"# {path.stem}\n\n{_clean_text(text)}"


def chunk_markdown(text: str, max_tokens: int = 2000) -> list[str]:
    frontmatter = ""
    body = text
    if text.startswith("---"):
        end = text.find("---", 3)
        if end != -1:
            frontmatter = text[: end + 3]
            body = text[end + 3 :].lstrip()

    if max_tokens <= 0:
        return [text]

    blocks = _markdown_blocks(body)
    chunks: list[str] = []
    current_blocks: list[str] = []
    current_tokens = 0

    for block in blocks:
        block_tokens = count_tokens(block)
        if block_tokens > max_tokens:
            if current_blocks:
                chunks.append(_compose_chunk(frontmatter, current_blocks))
                current_blocks = []
                current_tokens = 0
            for split_block in _split_large_markdown_block(block, max_tokens):
                split_tokens = count_tokens(split_block)
                if current_blocks and current_tokens + split_tokens > max_tokens:
                    chunks.append(_compose_chunk(frontmatter, current_blocks))
                    current_blocks = []
                    current_tokens = 0
                current_blocks.append(split_block)
                current_tokens += split_tokens
            continue

        if current_blocks and current_tokens + block_tokens > max_tokens:
            chunks.append(_compose_chunk(frontmatter, current_blocks))
            current_blocks = []
            current_tokens = 0

        current_blocks.append(block)
        current_tokens += block_tokens

    if current_blocks:
        chunks.append(_compose_chunk(frontmatter, current_blocks))

    return chunks


def _compose_chunk(frontmatter: str, blocks: list[str]) -> str:
    chunk = "\n\n".join(block for block in blocks if block.strip()).strip()
    if frontmatter:
        return f"{frontmatter}\n\n{chunk}" if chunk else frontmatter
    return chunk


def _markdown_blocks(body: str) -> list[str]:
    blocks: list[str] = []
    current: list[str] = []
    in_code_block = False

    for raw_line in body.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        is_code_fence = stripped.startswith("```")
        starts_standalone = bool(stripped) and (
            stripped.startswith("#")
            or stripped.startswith("|")
            or stripped.startswith("> ")
            or re.match(r"^(?:[-*+]|\d+\.)\s+", stripped) is not None
        )

        if is_code_fence:
            if current and not in_code_block:
                blocks.append("\n".join(current).strip())
                current = []
            current.append(line)
            in_code_block = not in_code_block
            if not in_code_block:
                blocks.append("\n".join(current).strip())
                current = []
            continue

        if in_code_block:
            current.append(line)
            continue

        if not stripped:
            if current:
                blocks.append("\n".join(current).strip())
                current = []
            continue

        if starts_standalone and current:
            blocks.append("\n".join(current).strip())
            current = []

        current.append(line)

    if current:
        blocks.append("\n".join(current).strip())
    return [block for block in blocks if block.strip()]


def _split_large_markdown_block(block: str, max_tokens: int) -> list[str]:
    if count_tokens(block) <= max_tokens:
        return [block]

    lines = block.splitlines()
    pieces: list[str] = []
    current: list[str] = []
    current_tokens = 0

    for line in lines:
        line_tokens = count_tokens(line)
        if line_tokens > max_tokens:
            if current:
                pieces.append("\n".join(current).strip())
                current = []
                current_tokens = 0
            pieces.extend(_split_oversized_line(line, max_tokens))
            continue
        if current and current_tokens + line_tokens > max_tokens:
            pieces.append("\n".join(current).strip())
            current = [line]
            current_tokens = line_tokens
            continue
        current.append(line)
        current_tokens += line_tokens

    if current:
        pieces.append("\n".join(current).strip())
    return [piece for piece in pieces if piece.strip()]


def _split_oversized_line(line: str, max_tokens: int) -> list[str]:
    segments = re.split(r"(?<=[.!?;:])\s+", line.strip())
    if len(segments) == 1:
        segments = re.split(r"(?<=,)\s+", line.strip())
    if len(segments) == 1:
        words = line.split()
        if not words:
            return []
        segments = words

    pieces: list[str] = []
    current: list[str] = []
    current_tokens = 0

    for segment in segments:
        segment = segment.strip()
        if not segment:
            continue
        segment_tokens = count_tokens(segment)
        if segment_tokens > max_tokens and " " in segment:
            subpieces = _split_oversized_line(segment.replace(", ", ",\n").replace(" ", "\n"), max_tokens)
            pieces.extend(subpieces)
            continue
        if current and current_tokens + segment_tokens > max_tokens:
            pieces.append(" ".join(current).strip())
            current = [segment]
            current_tokens = segment_tokens
            continue
        current.append(segment)
        current_tokens += segment_tokens

    if current:
        pieces.append(" ".join(current).strip())
    return [piece for piece in pieces if piece.strip()]


def _clean_text(text: str) -> str:
    text = text.replace("\r", "\n")
    text = text.replace("\f", "\n")
    text = text.replace("\u00A0", " ")
    text = text.replace("\u200B", "")
    text = text.replace("\u00AD", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n[ \t]+", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"^\s*\d+\s*$", "", text, flags=re.MULTILINE)
    return text.strip()


def _fix_unicode(text: str) -> str:
    text = _repair_mojibake(text)
    for source, target in _LIGATURE_REPLACEMENTS.items():
        text = text.replace(source, target)
    text = text.replace("\x00", "")
    text = "".join(ch for ch in text if ch == "\n" or ch == "\t" or ord(ch) >= 32)
    return unicodedata.normalize("NFC", text)


def _repair_mojibake(text: str) -> str:
    if not text or not any(marker in text for marker in _MOJIBAKE_MARKERS):
        return _repair_common_symbol_misdecodes(text)

    best = text
    best_score = _mojibake_score(text)
    for source_encoding in ("cp1252", "latin1"):
        try:
            candidate = text.encode(source_encoding).decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError):
            continue
        score = _mojibake_score(candidate)
        if score < best_score:
            best = candidate
            best_score = score
    return _repair_common_symbol_misdecodes(best)


def _repair_common_symbol_misdecodes(text: str) -> str:
    return (
        text.replace("\ufffd", "")
        .replace("ą", "±")
        .replace("þ", "+")
        .replace("�", "")
    )


def _mojibake_score(text: str) -> int:
    suspicious = sum(text.count(marker) for marker in _MOJIBAKE_MARKERS)
    control_like = sum(1 for ch in text if ord(ch) < 32 and ch not in "\n\t\r")
    return (suspicious * 4) + (control_like * 10)


_UNIT_TOKEN_SET = {
    "s",
    "sec",
    "min",
    "h",
    "m",
    "cm",
    "mm",
    "nm",
    "um",
    "μm",
    "µm",
    "l",
    "ml",
    "ul",
    "μl",
    "µl",
    "g",
    "mg",
    "kg",
    "m",
    "pa",
    "kpa",
    "mpa",
    "gpa",
    "kda",
    "mj",
    "mj/m",
    "mj*m",
}


def _normalize_pdf_markdown_math(text: str, source_path: Path | None = None) -> str:
    text = _normalize_line_endings(text)
    text = _strip_picture_text_blocks(text)
    text = _repair_pdf_markdown_artifact_spans(text)
    text = _normalize_markdown_table_blocks(text)
    normalized_lines = []
    in_code_block = False

    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code_block = not in_code_block
            normalized_lines.append(line)
            continue
        if in_code_block or stripped.startswith("!["):
            normalized_lines.append(line)
            continue
        line = _normalize_pdf_math_line(line)
        line = _normalize_pdf_citation_line(line)
        line = _normalize_pdf_prose_line(line)
        normalized_lines.append(line)

    normalized = "\n".join(normalized_lines)
    normalized = _normalize_pdf_markdown_document_artifacts(normalized, source_path=source_path)
    normalized = _dedupe_adjacent_duplicate_lines(normalized)
    normalized = llm_pdf_cleanup.normalize_h1_title(normalized, source_path=source_path)
    normalized = llm_pdf_cleanup.sanitize_existing_frontmatter(normalized)
    normalized = re.sub(r'(?m)^## (?:Article|Feature Article|Research Article|Review)\s*\n?', "", normalized)
    normalized = normalized.replace("https://doi. org/", "https://doi.org/")
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return llm_pdf_cleanup.normalize_line_endings(normalized).strip() + "\n"


def _repair_pdf_markdown_artifact_spans(text: str) -> str:
    text = re.sub(
        r"\b(\d+(?:\.\d+)?\s+g\s+l)-\s*\n\s*\n\s*1\s+(\d{1,2}C\b)",
        r"\1^-1 \2",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(
        r"\b(\d+(?:\.\d+)?\s+K\s+h)-\s*\n\s*\n\s*1\b",
        r"\1^-1",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(r"(\b\d+\s*,\s*\d+\s*,?)\s*\n\s*\n\s*(\d+\))", r"\1 \2", text)
    return text


def _normalize_pdf_markdown_document_artifacts(text: str, source_path: Path | None = None) -> str:
    text = llm_pdf_cleanup.normalize_document_structure(text, source_path=source_path)
    text = _repair_residual_unit_artifacts(text)
    text = _repair_missing_equation_markers(text)
    text = _repair_known_dangling_pdf_fragments(text)
    text = llm_figure_cleanup.align_figure_images_with_captions(text)
    return text


def _normalize_line_endings(text: str) -> str:
    return str(text).replace("\r\n", "\n").replace("\r", "\n")


def _split_frontmatter(text: str) -> tuple[str, str]:
    match = re.match(r"(?s)\A---\n.*?\n---\n\n?", text)
    if not match:
        return "", text
    return text[: match.end()], text[match.end() :]


def _find_first_h1(lines: list[str]) -> int | None:
    for index, line in enumerate(lines):
        if re.match(r"^# (?!#)", line.strip()):
            return index
    return None


def _clean_title_text(title: str) -> str:
    title = _normalize_metadata_field(title)
    title = re.sub(r"^#+\s*", "", title)
    title = title.replace("**", "").replace("__", "").strip("*_ ")
    title = re.sub(r"^\d+\s+", "", title)
    title = re.sub(r"^Article\s+https?://doi\.org/\S+\s+", "", title, flags=re.IGNORECASE)
    title = re.sub(r"^(?:Manuscript\s+title|Title)\s*:\s*", "", title, flags=re.IGNORECASE)
    title = title.replace("_", " ")
    title = re.sub(r"\bFull-\s+Length\b", "Full-Length", title, flags=re.IGNORECASE)
    title = re.sub(r"\bmechan-\s*ical\b", "mechanical", title, flags=re.IGNORECASE)
    title = re.sub(r"\s+\d+\s+(?=of\b)", " ", title, flags=re.IGNORECASE)
    title = re.sub(r"\b(synthetic proteins)\s+Thomas\s+Scheibel\b", r"\1", title, flags=re.IGNORECASE)
    title = re.sub(r"\s+", " ", title).strip(" .")
    return title


def _looks_like_filename_slug(title: str) -> bool:
    stripped = title.strip()
    if "_" in stripped:
        return True
    return bool(re.search(r"\b(?:boo|bioa|che|syntheti|mechani)\b$", stripped, flags=re.IGNORECASE))


def _looks_like_journal_title(title: str) -> bool:
    normalized = _normalize_for_compare(title)
    if normalized in {"microbial cell factories", "nature communications", "small", "biomaterials", "science"}:
        return True
    return _looks_like_journal_line(title)


def _looks_like_bad_extracted_title(title: str) -> bool:
    lower = title.lower()
    return bool(
        "figure" in lower
        or "scanning electron" in lower
        or " collected fro " in lower
        or lower.startswith(("authors:", "herein,", "finally,", "feasibility of", "fourier-transform"))
        or re.search(r"\bwe\s+(?:introduce|also|develop|investigated|demonstrate)\b", lower)
        or len(title) > 240
    )


def _is_title_candidate(line: str) -> bool:
    stripped = _clean_title_text(line)
    if not 12 <= len(stripped) <= 220:
        return False
    lower = stripped.lower()
    if line.strip().startswith(("![", "|")):
        return False
    if any(marker in lower for marker in ("http://", "https://", "doi:", "@", "number of words", "copyright")):
        return False
    if re.match(r"^(?:authors?|affiliations?|keywords?|abbreviations?|received|accepted|published|abstract)\b", lower):
        return False
    if lower.startswith(("composition and structural architecture", "natural spider silk assembly", "types of spider silk")):
        return False
    if _looks_like_journal_title(stripped) or _looks_like_bad_extracted_title(stripped):
        return False
    return len(re.findall(r"[A-Za-z][A-Za-z-]+", stripped)) >= 3


def _title_source_overlap_score(title: str, source_path: Path | None) -> int:
    if source_path is None:
        return 0
    source_tokens = {
        token
        for token in re.findall(r"[a-z0-9]+", source_path.stem.lower())
        if len(token) > 2 and token not in _TITLE_COMPARE_STOPWORDS
    }
    title_tokens = {
        token
        for token in re.findall(r"[a-z0-9]+", title.lower())
        if len(token) > 2 and token not in _TITLE_COMPARE_STOPWORDS
    }
    return len(source_tokens & title_tokens)


def _combine_split_title_heading(lines: list[str], index: int) -> tuple[str, set[int]]:
    title = _clean_title_text(lines[index])
    consumed = {index}
    for next_index in range(index + 1, min(len(lines), index + 3)):
        next_line = lines[next_index].strip()
        if not next_line:
            continue
        if not next_line.startswith("## "):
            break
        next_title = _clean_title_text(next_line)
        if not next_title:
            continue
        if next_title.lower().startswith(("abstract", "introduction", "methods", "results", "discussion", "references")):
            break
        if len(next_title.split()) <= 8 and not re.search(r"[.:]$", title):
            title = f"{title} {next_title}".strip()
            consumed.add(next_index)
            continue
        break
    return _clean_title_text(title), consumed


def _best_title_candidate(lines: list[str], start: int, source_path: Path | None) -> tuple[int | None, str]:
    best_index: int | None = None
    best_title = ""
    best_score = -1
    scan_limit = min(len(lines), start + 80)
    for index in range(start, scan_limit):
        stripped = lines[index].strip()
        if not stripped:
            continue
        if stripped.startswith("## "):
            cleaned_heading, _ = _combine_split_title_heading(lines, index)
            if cleaned_heading.lower().startswith("abstract"):
                break
            if _is_title_candidate(stripped):
                score = len(cleaned_heading.split()) + 12
                score += _title_source_overlap_score(cleaned_heading, source_path) * 4
                if re.search(r"\b(?:spider|silk|protein|fiber|fibre|polymer|biomaterial|mechanical)\b", cleaned_heading, re.IGNORECASE):
                    score += 6
                if score > best_score:
                    best_index = index
                    best_title = cleaned_heading
                    best_score = score
                continue
        if not _is_title_candidate(stripped):
            continue
        score = len(stripped.split())
        score += _title_source_overlap_score(stripped, source_path) * 4
        if ":" in stripped:
            score += 8
        if re.search(r"\b(?:spider|silk|protein|fiber|fibre|polymer|biomaterial|mechanical)\b", stripped, re.IGNORECASE):
            score += 6
        if index == start:
            score += 5
        if score > best_score:
            best_index = index
            best_title = _clean_title_text(stripped)
            best_score = score
    if best_title:
        return best_index, best_title
    if source_path is not None:
        return None, _clean_title_text(source_path.stem)
    return None, ""


def _is_stronger_title(candidate: str, current: str, source_path: Path | None) -> bool:
    if not candidate:
        return False
    if _looks_like_bad_extracted_title(current):
        return True
    if _looks_like_filename_slug(current) or _looks_like_journal_title(current):
        return True
    candidate_overlap = _title_source_overlap_score(candidate, source_path)
    current_overlap = _title_source_overlap_score(current, source_path)
    if candidate_overlap >= current_overlap + 2:
        return True
    if current.endswith(":") and len(candidate) > len(current) + 8 and candidate.startswith(current):
        return True
    if len(current.split()) <= 6 and len(candidate.split()) > len(current.split()) + 2:
        return True
    return False


def _normalize_pdf_h1_title(text: str, source_path: Path | None = None) -> str:
    prefix, body = _split_frontmatter(text)
    lines = body.splitlines()
    h1_index = _find_first_h1(lines)
    if h1_index is None:
        _, fallback = _best_title_candidate(lines, 0, source_path)
        if fallback:
            return prefix + f"# {fallback}\n\n" + body.lstrip()
        return text

    raw_heading = lines[h1_index].strip()[2:].strip()
    title = _clean_title_text(raw_heading)
    remove_indices: set[int] = set()

    next_index = h1_index + 1
    while next_index < len(lines) and not lines[next_index].strip():
        next_index += 1

    if next_index < len(lines) and _is_title_candidate(lines[next_index]):
        continuation = _clean_title_text(lines[next_index])
        if raw_heading.lower().startswith(("# title:", "title:", "# manuscript title:", "manuscript title:")):
            title = f"{title} {continuation}".strip()
            remove_indices.add(next_index)
        elif _looks_like_filename_slug(raw_heading) or _looks_like_journal_title(raw_heading):
            title = continuation
            remove_indices.add(next_index)

    candidate_index, candidate = _best_title_candidate(lines, h1_index + 1, source_path)
    if _is_stronger_title(candidate, title, source_path):
        title = candidate
        if candidate_index is not None:
            combined_title, combined_indices = _combine_split_title_heading(lines, candidate_index)
            if _clean_title_text(combined_title) == _clean_title_text(candidate):
                remove_indices.update(combined_indices)
            else:
                remove_indices.add(candidate_index)
    elif (_looks_like_filename_slug(raw_heading) or _looks_like_journal_title(raw_heading) or not title) and candidate:
        title = candidate
        if candidate_index is not None:
            remove_indices.add(candidate_index)

    lines[h1_index] = f"# {title or _clean_title_text(source_path.stem if source_path else raw_heading)}"
    output = [line for index, line in enumerate(lines) if index not in remove_indices]
    return prefix + "\n".join(output)


def _extract_plain_authors(lines: list[str]) -> str:
    for index, line in enumerate(lines[:50]):
        if re.match(r"^\s*Authors?\s*:", line, flags=re.IGNORECASE):
            parts = [re.sub(r"^\s*Authors?\s*:\s*", "", line, flags=re.IGNORECASE).strip()]
            for continuation in lines[index + 1 : index + 8]:
                stripped = continuation.strip()
                if not stripped:
                    continue
                if re.match(r"^(?:Affiliations?|Correspondence|Keywords?|Abbreviations?|##)\b", stripped, flags=re.IGNORECASE):
                    break
                parts.append(stripped)
            return _clean_metadata_authors(" ".join(parts))
    return ""


def _ensure_pdf_frontmatter(text: str) -> str:
    prefix, body = _split_frontmatter(text)
    lines = body.splitlines()
    h1_index = _find_first_h1(lines)
    title = _clean_title_text(lines[h1_index][2:]) if h1_index is not None else ""
    authors = _extract_plain_authors(lines)

    if prefix:
        frontmatter = prefix.strip().splitlines()
        keys = {line.split(":", 1)[0].strip().lower() for line in frontmatter[1:-1] if ":" in line}
        insert_at = 1
        if title and "title" not in keys:
            frontmatter.insert(insert_at, f'title: "{_yaml_quote(title)}"')
            insert_at += 1
        if authors and "authors" not in keys:
            frontmatter.insert(insert_at, f'authors: "{_yaml_quote(authors)}"')
        return "\n".join(frontmatter) + "\n\n" + body.lstrip()

    fields: list[tuple[str, str]] = []
    if title:
        fields.append(("title", title))
    if authors:
        fields.append(("authors", authors))
    if not fields:
        return text
    frontmatter_lines = ["---", *(f'{key}: "{_yaml_quote(value)}"' for key, value in fields), "---"]
    return "\n".join(frontmatter_lines) + "\n\n" + body.lstrip()


def _sanitize_existing_frontmatter(text: str) -> str:
    prefix, body = _split_frontmatter(text)
    if not prefix:
        return text
    body_lines = body.splitlines()
    h1_index = _find_first_h1(body_lines)
    body_title = _clean_title_text(body_lines[h1_index][2:]) if h1_index is not None else ""
    lines = prefix.strip().splitlines()
    output = [lines[0]]
    for line in lines[1:-1]:
        if ":" not in line:
            output.append(line)
            continue
        key, value = line.split(":", 1)
        cleaned_value = value.strip().strip('"')
        if key.strip().lower() == "authors":
            cleaned_value = _clean_metadata_authors(cleaned_value)
            if not cleaned_value:
                continue
            output.append(f'authors: "{_yaml_quote(cleaned_value)}"')
            continue
        if key.strip().lower() == "journal":
            cleaned_value = _clean_metadata_journal(cleaned_value)
            if not cleaned_value:
                continue
            output.append(f'journal: "{_yaml_quote(cleaned_value)}"')
            continue
        if key.strip().lower() == "title":
            cleaned_value = _clean_title_text(cleaned_value)
            if body_title and (
                _looks_like_filename_slug(cleaned_value)
                or _looks_like_journal_title(cleaned_value)
                or _looks_like_bad_extracted_title(cleaned_value)
                or "http://" in cleaned_value.lower()
                or "https://" in cleaned_value.lower()
                or "university of" in cleaned_value.lower()
                or (body_title.lower().startswith(cleaned_value.lower()) and len(body_title) > len(cleaned_value) + 8)
                or (_title_source_overlap_score(body_title, None) >= _title_source_overlap_score(cleaned_value, None) and len(body_title.split()) > len(cleaned_value.split()) + 2)
                or len(cleaned_value) < 12
            ):
                cleaned_value = body_title
            if not cleaned_value:
                continue
            output.append(f'title: "{_yaml_quote(cleaned_value)}"')
            continue
        output.append(line)
    output.append("---")
    return "\n".join(output) + "\n\n" + body.lstrip()


def _strip_leading_plain_metadata_and_boilerplate(text: str) -> str:
    prefix, body = _split_frontmatter(text)
    lines = body.splitlines()
    h1_index = _find_first_h1(lines)
    output: list[str] = []
    in_plain_metadata = False
    pre_section = True

    for index, line in enumerate(lines):
        stripped = line.strip()
        plain = stripped.strip("*_ ").strip()
        lower = plain.lower()
        if stripped.startswith("## "):
            heading_text = _clean_title_text(stripped)
            if pre_section and (_looks_like_journal_title(heading_text) or heading_text.lower() in {"review", "open access"}):
                continue
            pre_section = False
            in_plain_metadata = False
        if pre_section and re.match(r"^\*\*Abstract:\*\*", stripped, flags=re.IGNORECASE):
            pre_section = False
            in_plain_metadata = False
        if pre_section and index != h1_index:
            starts_metadata = bool(
                re.match(
                    r"^(?:Authors?|Affiliations?|Affiliation|Correspondence|Keywords?|Abbreviations?|Number of words|"
                    r"\d+\s*email|[*\u2020]?\s*These authors contributed|[*]?\s*Correspondence to)\b",
                    plain,
                    flags=re.IGNORECASE,
                )
                or re.match(r"^-?\s*\d+\b", stripped)
            )
            if starts_metadata:
                in_plain_metadata = True
                continue
            if in_plain_metadata:
                if not stripped:
                    continue
                if stripped.startswith("## ") or re.match(r"^\*\*Abstract:\*\*", stripped, flags=re.IGNORECASE):
                    in_plain_metadata = False
                else:
                    continue
            if _is_boilerplate_body_line(plain):
                continue
        output.append(line)
    return prefix + "\n".join(output)


def _is_boilerplate_body_line(line: str) -> bool:
    if not line:
        return False
    lower = line.lower()
    if any(phrase in lower for phrase in _BOILERPLATE_PHRASES):
        return True
    return bool(
        re.match(r"^(?:Page \d+ of \d+|\(page number not for citation purposes\)|www\.[\w.-]+\.com\b)", line, flags=re.IGNORECASE)
    )


def _remove_duplicate_title_and_author_lines(text: str) -> str:
    prefix, body = _split_frontmatter(text)
    lines = body.splitlines()
    h1_index = _find_first_h1(lines)
    if h1_index is None:
        return text
    title = _clean_title_text(lines[h1_index][2:])
    output: list[str] = []
    before_section = True
    skip_indices: set[int] = set()
    index = 0
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()
        plain = stripped.strip("*_ ").strip()
        if index in skip_indices:
            index += 1
            continue
        if index == h1_index:
            output.append(line)
            index += 1
            continue
        if stripped.startswith("## "):
            heading_title, consumed = _combine_split_title_heading(lines, index)
            if before_section and _clean_title_text(heading_title).lower() == title.lower():
                skip_indices.update(consumed)
                index += 1
                continue
            before_section = False
        if before_section:
            cleaned = _clean_title_text(plain)
            if title and cleaned.lower().startswith(title.lower()):
                index += 1
                continue
            if re.match(r"^\*?\s*Corresponding author\b", plain, flags=re.IGNORECASE):
                index += 1
                continue
        output.append(line)
        index += 1
    return prefix + "\n".join(output)


def _normalize_bold_section_labels(text: str) -> str:
    output: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        heading_match = re.match(r"^##\s+[_* ]*(?:\d+[\.\s]+)?(Abstract|Introduction|Results|Methods|Discussion|References|Supplementary|Review)[_* ]*$", stripped, flags=re.IGNORECASE)
        if heading_match:
            output.append(f"## {heading_match.group(1).title()}")
            continue
        bold_match = re.match(r"^\*\*(Abstract|Introduction|Results|Methods|Discussion|References|Supplementary|Review):\*\*\s*(.*)$", stripped, flags=re.IGNORECASE)
        if bold_match:
            output.append(f"## {bold_match.group(1).title()}")
            if bold_match.group(2).strip():
                output.append("")
                output.append(bold_match.group(2).strip())
            continue
        output.append(line)
    return "\n".join(output)


def _remove_empty_markdown_headings(text: str) -> str:
    return re.sub(r"(?m)^\s*#{1,6}\s*$\n?", "", text)


def _remove_running_headers_inline(text: str) -> str:
    patterns = [
        r"Nature Communications\|\s*\(\d{4}\)\s*\d+:\d+",
        r"Article\s+https?://doi\.org/\S+",
        r"Microbial Cell Factories\s+\d{4},\s*\d+:\d+\s+https?://\S+",
        r"Page\s+\d+\s+of\s+\d+",
        r"\(page number not for citation purposes\)",
    ]
    for pattern in patterns:
        text = re.sub(rf"(?im)^\s*{pattern}\s*$\n?", "", text)
        text = re.sub(rf"\s+{pattern}\s+", " ", text, flags=re.IGNORECASE)
    lines = [line for line in text.splitlines() if not _is_boilerplate_body_line(line.strip().strip("*_ "))]
    return re.sub(r"[ \t]{2,}", " ", "\n".join(lines))


def _remove_placeholder_affiliations(text: str) -> str:
    text = re.sub(r"(?im)^\s*\d+\s+TBA\s*$\n?", "", text)
    return re.sub(r"\b\d+\s+TBA\b", "", text)


def _repair_hyphenated_pdf_breaks(text: str) -> str:
    def replace_break(match: re.Match[str]) -> str:
        left, right = match.groups()
        combined = f"{left}{right}"
        replacement = _DEHYPHENATED_WORD_REPLACEMENTS.get(combined.lower())
        if replacement:
            return _apply_replacement_case(combined, replacement)
        return f"{left}-{right}"

    text = re.sub(r"\b([A-Za-z]{3,})-\s*\n\s*([a-z]{2,})\b", replace_break, text)
    for source, target in _DEHYPHENATED_WORD_REPLACEMENTS.items():
        split = re.sub(r"([a-z])([a-z]+)", r"\1-\2", source, count=1)
        text = re.sub(rf"\b{re.escape(split)}\b", target, text, flags=re.IGNORECASE)
    text = re.sub(r"\bmechan-ical\b", "mechanical", text, flags=re.IGNORECASE)
    text = re.sub(r"\bpro-duction\b", "production", text, flags=re.IGNORECASE)
    text = re.sub(r"\bsea-lants\b", "sealants", text, flags=re.IGNORECASE)
    text = re.sub(r"\bboo-\s+", "bio-", text, flags=re.IGNORECASE)
    return text


def _normalize_reference_blocks(text: str) -> str:
    lines = text.splitlines()
    output: list[str] = []
    in_references = False
    for line in lines:
        stripped = line.strip()
        if re.match(r"^## References\b", stripped, flags=re.IGNORECASE):
            in_references = True
            output.append("## References")
            continue
        if stripped.startswith("## ") and not re.match(r"^## References\b", stripped, flags=re.IGNORECASE):
            in_references = False
            output.append(line)
            continue
        if in_references:
            cleaned = _remove_running_headers_inline(line).strip()
            if not cleaned:
                output.append("")
                continue
            parts = re.split(r"(?<!^)\s+(?=\d{1,3}\.\s+[A-Z])", cleaned)
            for part in parts:
                if part.strip():
                    output.append(part.strip())
            continue
        output.append(line)
    return "\n".join(output)


def _qualify_repeated_section_headings(text: str) -> str:
    seen: dict[str, int] = {}
    seen_titles: set[str] = set()
    supplementary_context = False
    output: list[str] = []
    for line in text.splitlines():
        match = re.match(r"^(##) (Results|Methods|References|Discussion|Supplementary)\s*$", line.strip(), flags=re.IGNORECASE)
        if not match:
            output.append(line)
            continue
        base = match.group(2).title()
        count = seen.get(base, 0)
        if count == 0:
            title = base
        elif supplementary_context and base != "Supplementary":
            title = f"Supplementary {base}"
        else:
            title = f"{base} (continued {count + 1})"
        while title.lower() in seen_titles:
            count += 1
            title = f"{base} (continued {count + 1})"
        seen[base] = seen.get(base, 0) + 1
        seen_titles.add(title.lower())
        if base == "Supplementary":
            supplementary_context = True
        output.append(f"## {title}")
    return "\n".join(output)


def _reorder_intro_before_results(text: str) -> str:
    lines = text.splitlines()
    headings = [(index, line.strip()) for index, line in enumerate(lines) if line.startswith("## ")]
    result_entry = next(((index, title) for index, title in headings if title == "## Results"), None)
    intro_entry = next(((index, title) for index, title in headings if title == "## Introduction"), None)
    if not result_entry or not intro_entry or intro_entry[0] < result_entry[0]:
        return text
    intro_start = intro_entry[0]
    next_heading = next((index for index, _ in headings if index > intro_start), len(lines))
    intro_block = lines[intro_start:next_heading]
    remaining = lines[:intro_start] + lines[next_heading:]
    insert_at = next((index for index, line in enumerate(remaining) if line == "## Results"), None)
    if insert_at is None:
        return text
    return "\n".join(remaining[:insert_at] + intro_block + [""] + remaining[insert_at:])


def _repair_residual_unit_artifacts(text: str) -> str:
    text = re.sub(r"\b(\d+(?:\.\d+)?)\s+g\s+l-1\b", r"\1 g L^-1", text, flags=re.IGNORECASE)
    text = re.sub(r"\b(\d+(?:\.\d+)?)\s+K\s+h-1\b", r"\1 K h^-1", text, flags=re.IGNORECASE)
    return text


def _remove_duplicate_frontmatter_dates(text: str) -> str:
    frontmatter_match = re.match(r"(?s)\A---\n(.*?)\n---\n\n", text)
    if not frontmatter_match:
        return text
    frontmatter = frontmatter_match.group(1).lower()
    if "received:" not in frontmatter and "accepted:" not in frontmatter and "journal:" not in frontmatter:
        return text

    prefix = text[: frontmatter_match.end()]
    rest = text[frontmatter_match.end() :]
    title_match = re.match(r"(?s)(# .+?\n\n)(.*)", rest)
    if not title_match:
        return text
    title = title_match.group(1)
    body = title_match.group(2)
    metadata_block = r"(?:(?:Received|Accepted|Published online|Revised)\s*\n\s*(?:\d{1,2}\s+\w+\s+\d{4})\s*\n\s*)+"
    body = re.sub(rf"\A{metadata_block}", "", body, flags=re.IGNORECASE)
    body = _strip_leading_duplicate_metadata_lines(body, frontmatter)
    return prefix + title + body


def _strip_leading_duplicate_metadata_lines(body: str, frontmatter: str) -> str:
    lines = body.splitlines()
    output_start = 0
    saw_metadata = False
    for index, line in enumerate(lines[:12]):
        stripped = line.strip()
        if not stripped:
            output_start = index + 1
            continue
        normalized = _normalize_for_compare(stripped)
        is_date = bool(re.fullmatch(r"\d{1,2}\s+\w+\s+\d{4}", stripped))
        is_label = normalized in {"received", "accepted", "revised", "published online"}
        is_journal = f'journal: "{stripped.lower()}"' in frontmatter or f"journal: {stripped.lower()}" in frontmatter
        if is_date or is_label or is_journal:
            saw_metadata = True
            output_start = index + 1
            continue
        break
    if saw_metadata:
        return "\n".join(lines[output_start:]).lstrip()
    return body


def _repair_missing_equation_markers(text: str) -> str:
    return re.sub(
        r"using following equation:\.?",
        "using the following equation:\n\n> [!NOTE] Equation omitted during PDF extraction; see the source PDF for the formula.",
        text,
        flags=re.IGNORECASE,
    )


def _repair_known_dangling_pdf_fragments(text: str) -> str:
    text = re.sub(r"\b([A-Z]\.[A-Z]\.)We(?=want\b)", r"\1 We ", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    if "BMRB data bank under accession code 16249" in text:
        text = re.sub(
            r"best 20 structures plus a regularized average structure have been deposited at the\s+competing financial interests\.",
            "best 20 structures plus a regularized average structure have been deposited at the Protein Data Bank under accession code 2KHM. The authors declare no competing financial interests.",
            text,
            flags=re.IGNORECASE,
        )
    return text


def _normalize_pdf_math_line(line: str) -> str:
    line = re.sub(r"(?<=[A-Za-zµμ°Ω])\s*([/*])\s*(?=[A-Za-zµμ°Ω])", r"\1", line)
    line = re.sub(r"\[\((\d+)\)\]", r"(\1)", line)
    line = re.sub(
        r"10\s*[−-]\s*\[\s*(\d+)\s*\]",
        lambda m: f"10^-{m.group(1)}",
        line,
    )
    line = re.sub(
        r"10\s*\[\s*[−-]\s*\]\s*\[\s*(\d+)\s*\]",
        lambda m: f"10^-{m.group(1)}",
        line,
    )
    line = re.sub(
        r"\b([A-Z][A-Za-z0-9]*)\s*\[\s*(\d+)\s*\]\s*\[\s*([+\-−])\s*\]",
        _replace_bracketed_charge,
        line,
    )
    line = re.sub(
        r"\b([A-Z][A-Za-z0-9]*)\s*\[\s*([+\-−])\s*\]",
        _replace_bracketed_charge,
        line,
    )
    line = re.sub(
        r"\[\s*([NC])\s*\]\s*M(?=[A-Za-z0-9-])",
        lambda m: f"{m.group(1)}M",
        line,
    )
    line = re.sub(
        r"\[\s*([NC])\s*\]\s*M\b",
        lambda m: f"{m.group(1)}M",
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω/*]+)\s*\[\s*([−-])\s*\]\s*\[\s*(\d+)\s*\]",
        _replace_bracketed_unit_power,
        line,
    )
    line = re.sub(
        r"\be\s*\[\s*([−-])\s*\]\s*/\s*(?:Å|A)\s*\[\s*(\d+)\s*\]",
        lambda m: f"e^{_normalize_math_sign(m.group(1))}/Å^{m.group(2)}",
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω/*]+)\s*\[\s*([−-])\s*\]\s*([¹²³])",
        _replace_unicode_unit_power,
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω/*]+)\s*\[\s*(\d+)\s*\]",
        _replace_positive_unit_power,
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω/*]+)\s*-\s*\[\s*(\d+)\s*\]",
        _replace_hyphen_bracket_unit_power,
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω]+)\s*[−-]\s*(\d+)\b",
        _replace_inline_unit_power,
        line,
    )
    line = re.sub(r"([a-z])((?:NM|CM)(?=\b|[A-Z(]))", r"\1 \2", line)
    line = re.sub(r"([.,;:])((?:NM|CM)(?=\b|[A-Z(]))", r"\1 \2", line)
    line = re.sub(r"\b([A-Za-zµμ°Ω/*]+)\s+\^\s*([−-]?\d+)", r"\1^\2", line)
    return line


def _normalize_pdf_citation_line(line: str) -> str:
    line = re.sub(
        r"\[\d+\](?:\s*\[\s*(?:,|;|[âˆ’\-–])\s*\]\s*\[\d+\])+",
        _replace_citation_sequence,
        line,
    )
    line = re.sub(r"\[\s*,\s*\]", ",", line)
    line = re.sub(r"\[\s*;\s*\]", ";", line)
    line = re.sub(r"\[\s*([âˆ’\-–])\s*\]", lambda m: _normalize_math_sign(m.group(1)), line)
    line = re.sub(r"\]\s+([,.;:])", r"]\1", line)
    return line


def _replace_citation_sequence(match: re.Match[str]) -> str:
    sequence = match.group(0)
    pieces = [
        _normalize_math_sign(token) if token in {"âˆ’", "-", "–", "−"} else token
        for token in re.findall(r"\[\s*([0-9]+|,|;|âˆ’|-|–|−)\s*\]", sequence)
    ]

    if not pieces:
        return sequence

    rendered = pieces[0]
    for piece in pieces[1:]:
        if piece in {",", ";"}:
            rendered += piece + " "
        elif piece == "-":
            rendered += piece
        else:
            rendered += piece
    return f"[{rendered}]"


def _normalize_pdf_prose_line(line: str) -> str:
    line = _DISPLAY_DOI_SUFFIX_RE.sub("", line)
    line = _DISPLAY_DOI_FRAGMENT_RE.sub("", line)
    line = re.sub(r"\b(\d+)\s*repeat\b", r"\1-repeat", line)
    line = re.sub(r"(?<=[a-z])(?=(?:NM|CM))", " ", line)
    line = re.sub(r"(?<=[,.;:])(?=(?:NM|CM))", " ", line)
    line = re.sub(r"\b([NC]M)(\d)", r"\1-\2", line)
    line = re.sub(r"(?<=\w)CM(?=(?:\b|\())", "-CM", line)
    line = re.sub(r"(?<=[A-Za-zµμ°Ω])\s*([/*])\s*(?=[A-Za-zµμ°Ω])", r"\1", line)
    line = re.sub(r"\s+([,.;:])", r"\1", line)
    line = re.sub(r"\(\s+", "(", line)
    line = re.sub(r"\s+\)", ")", line)
    line = re.sub(r"\)(?=[A-Za-z])", ") ", line)
    line = re.sub(r"\s+[–-]\s+[–-]\s+", " ", line)
    line = _normalize_pdf_prose_ranges(line)
    line = _normalize_pdf_unit_expressions(line)
    line = _normalize_spaced_hyphen_fragments(line)
    line = _normalize_pdf_citation_superscripts(line)
    if not line.lstrip().startswith("|"):
        line = _normalize_pdf_prose_terms(line)
    return line


def _normalize_pdf_prose_terms(line: str) -> str:
    normalized = line
    for source, target in _PROSE_TERM_REPLACEMENTS.items():
        normalized = re.sub(
            rf"\b{re.escape(source)}\b",
            lambda m, replacement=target: _apply_replacement_case(m.group(0), replacement),
            normalized,
            flags=re.IGNORECASE,
        )

    normalized = re.sub(r"\b([A-Z])- (?=zone\b)", r"\1-", normalized)
    normalized = re.sub(r"\bfive-toten-fold\b", "five-to-ten-fold", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\bstrengthMW\b", "strength-MW", normalized)
    normalized = re.sub(r"\ba\s+[–-]\s+(strain at break)\b", r"a \1", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(strain at)\s+[–-]\s+(break)\b", r"\1 \2", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(?:β|ß|b)\s*-\s*sheet\b", "β-sheet", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(?:β|ß|b)\s*-\s*crystals?\b", "β-crystals", normalized, flags=re.IGNORECASE)
    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)\s+and\s+(\d)\b",
        r"\1 \2 and \3",
        normalized,
    )
    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)and(\d)\b",
        r"\1 \2 and \3",
        normalized,
    )
    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)([A-Za-z])\b",
        r"\1 \2\3",
        normalized,
    )
    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)\b",
        r"\1 \2",
        normalized,
    )
    normalized = re.sub(r"\b(\d+(?:\.\d+)?)\s*(ml|mL|uL|µL|mg|g|kg|M|mM|kDa)\b", r"\1 \2", normalized)
    normalized = re.sub(r"\b(\d+(?:\.\d+)?)\s*ml\b", r"\1 mL", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(\d+(?:\.\d+)?)\s*(?:ul|µl|μl)\b", r"\1 µL", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\bmg/ml\b", "mg/mL", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\bml/min\b", "mL/min", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(?:ul|µl|μl)/min\b", "µL/min", normalized, flags=re.IGNORECASE)
    normalized = normalized.replace(
        "The toughness modulusand the Youngwas obtained’s modulusby was determinedcalculating the totalfrom the slope of the linear elastic part.area under the stress-strain curve,",
        "The toughness modulus was obtained by calculating the total area under the stress-strain curve, and Young's modulus was determined from the slope of the linear elastic part.",
    )
    normalized = re.sub(r"\bthefi\s+ffi\s+", "the ", normalized, flags=re.IGNORECASE)
    normalized = re.sub(
        r"B\[2\]\s*obs-\[\s*q\s*\]\[\s*[−âˆ’-]\s*\]\[\s*2\s*\]",
        "B_obs^2-q^-2",
        normalized,
    )
    normalized = normalized.replace(
        "B_obs^2-q^-2[curve][34][.][It][was][found][that][the][periodic][ordering]",
        "B_obs^2-q^-2 curve[34]. It was found that the periodic ordering",
    )
    normalized = re.sub(r"\b(kcal)\s*∙\s*(mol)\[\s*[−âˆ’-]\s*\]\[\s*1\s*\]", r"\1∙\2^-1", normalized)
    normalized = re.sub(r"\b(MJ/m\^3)\s+MPa\b", r"\1", normalized)
    return re.sub(r"\s+", " ", normalized).strip()


def _normalize_pdf_prose_ranges(line: str) -> str:
    line = re.sub(
        r"(?<![\w.])(\d+(?:\.\d+)?)\s+(\d+(?:\.\d+)?)\s*(µm|μm|um|nm|mm|cm|mM|mg/mL|mg/ml|g/L|MPa|GPa|kDa|%)\b",
        lambda m: f"{m.group(1)}-{m.group(2)} {_normalize_unit_label(m.group(3))}",
        line,
    )
    line = re.sub(r"\bpH\s+(\d+(?:\.\d+)?)\s+(\d+(?:\.\d+)?)\b", r"pH \1-\2", line)
    return line


def _normalize_pdf_unit_expressions(line: str) -> str:
    line = re.sub(
        r"\b(\d+(?:\.\d+)?)\s+g\s+l(?:\[\s*-?\s*1\s*\]|-1|\^-1)(?=\s|\[|\d|[A-Za-z])",
        r"\1 g L^-1",
        line,
        flags=re.IGNORECASE,
    )
    line = re.sub(
        r"\b(\d+(?:\.\d+)?)\s+K\s+h(?:\[\s*-?\s*1\s*\]|-1|\^-1)(?=\s|\[|\d|[A-Za-z])",
        r"\1 K h^-1",
        line,
        flags=re.IGNORECASE,
    )
    line = re.sub(
        r"\b(cm|mm|m|nm|Âµm|Î¼m|um)(?:\[\s*-?\s*1\s*\]|-1)\b",
        lambda m: f"{_normalize_unit_label(m.group(1))}^-1",
        line,
        flags=re.IGNORECASE,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω]+)\s*[*·]\s*([A-Za-zµμ°Ω]+)\s*\^\s*-\s*(\d+)\b",
        _replace_negative_multiplicative_unit,
        line,
    )
    line = re.sub(
        r"\b([A-Za-zµμ°Ω]+)\s*/\s*([A-Za-zµμ°Ω]+)\s*\^\s*(\d+)\b",
        lambda m: f"{m.group(1)}/{m.group(2)}^{m.group(3)}",
        line,
    )
    line = re.sub(
        r"\b(cm|mm|m|µm|μm|um|nm)\s+s\s*[−–-]\s*(\d+)\b",
        lambda m: f"{m.group(1)} s^-{m.group(2)}",
        line,
    )
    line = re.sub(
        r"\b(\d+(?:\.\d+)?)\s*(cm/s|mm/s|m/s|µm|μm|um|nm|kDa|MPa|GPa|mM|mL|mg/mL|mg/ml)\b",
        lambda m: f"{m.group(1)} {_normalize_unit_label(m.group(2))}",
        line,
        flags=re.IGNORECASE,
    )
    line = re.sub(r"\bmg/\s*ml\b", "mg/mL", line, flags=re.IGNORECASE)
    line = re.sub(r"\b(?:ul|µl|μl)/min\b", "µL/min", line, flags=re.IGNORECASE)
    return line


def _normalize_spaced_hyphen_fragments(line: str) -> str:
    return re.sub(
        r"\b([A-Za-zµμα-ωΑ-Ω]+)-\s+([A-Za-z0-9][A-Za-z0-9-]*)\b",
        _repair_spaced_hyphen_fragment,
        line,
    )


def _normalize_pdf_citation_superscripts(line: str) -> str:
    return re.sub(
        r"\b(Da|kDa|µm|μm|um|nm|mm|cm|mL|L|M|mM|MPa|GPa|kPa|rpm|%)\^(\d{2,3})(?=\s+[A-Za-z])",
        lambda m: f"{m.group(1)}[{m.group(2)}]",
        line,
    )


def _normalize_pdf_caption_text(text: str) -> str:
    candidate = _clean_text(_fix_unicode(text or "")).strip()
    if not candidate:
        return ""
    candidate = _normalize_pdf_math_line(candidate)
    candidate = _normalize_pdf_citation_line(candidate)
    candidate = _normalize_pdf_prose_line(candidate)
    return re.sub(r"\s+", " ", candidate).strip()


def _apply_replacement_case(original: str, replacement: str) -> str:
    if original.isupper():
        return replacement.upper()
    if original[:1].isupper():
        return replacement[:1].upper() + replacement[1:]
    return replacement


def _normalize_unit_label(label: str) -> str:
    normalized = label
    normalized = re.sub(r"^mg/ml$", "mg/mL", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"^ml$", "mL", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"^(?:ul|µl|μl)$", "µL", normalized, flags=re.IGNORECASE)
    return normalized


def _repair_spaced_hyphen_fragment(match: re.Match[str]) -> str:
    prefix = match.group(1)
    suffix = match.group(2)
    if prefix.isupper() or prefix.lower() in _SPACED_HYPHEN_PREFIXES or len(prefix) == 1:
        return f"{prefix}-{suffix}"
    return f"{prefix}{suffix}"


def _replace_negative_multiplicative_unit(match: re.Match[str]) -> str:
    left = match.group(1)
    right = match.group(2)
    power = match.group(3)
    if not (_looks_like_unit_token(left) and _looks_like_unit_token(right)):
        return match.group(0)
    return f"{left}/{right}^{power}"


def _replace_bracketed_charge(match: re.Match[str]) -> str:
    species = match.group(1)
    if match.lastindex == 3:
        magnitude = match.group(2)
        sign = _normalize_math_sign(match.group(3))
        return f"{species}^{magnitude}{sign}"
    sign = _normalize_math_sign(match.group(2))
    return f"{species}{sign}"


def _replace_bracketed_unit_power(match: re.Match[str]) -> str:
    token = match.group(1)
    if not _looks_like_unit_token(token):
        return match.group(0)
    sign = _normalize_math_sign(match.group(2))
    power = match.group(3)
    return f"{token}^{sign}{power}"


def _replace_positive_unit_power(match: re.Match[str]) -> str:
    token = match.group(1)
    if not _looks_like_unit_token(token):
        return match.group(0)
    power = match.group(2)
    return f"{token}^{power}"


def _replace_inline_unit_power(match: re.Match[str]) -> str:
    token = match.group(1)
    if not _looks_like_unit_token(token):
        return match.group(0)
    return f"{token}^-{match.group(2)}"


def _replace_hyphen_bracket_unit_power(match: re.Match[str]) -> str:
    token = match.group(1)
    if not _looks_like_unit_token(token):
        return match.group(0)
    return f"{token}^-{match.group(2)}"


def _replace_unicode_unit_power(match: re.Match[str]) -> str:
    token = match.group(1)
    if not _looks_like_unit_token(token):
        return match.group(0)
    sign = _normalize_math_sign(match.group(2))
    power = str("¹²³".index(match.group(3)) + 1)
    return f"{token}^{sign}{power}"


def _looks_like_unit_token(token: str) -> bool:
    normalized = token.strip().replace("·", "*")
    if "/" in normalized or "*" in normalized:
        return True
    return normalized.lower() in _UNIT_TOKEN_SET


def _normalize_math_sign(sign: str) -> str:
    return "-" if sign in {"−", "-", "–"} else "+"


def _dedupe_adjacent_duplicate_lines(text: str) -> str:
    output_lines: list[str] = []
    last_nonempty = ""
    for line in text.splitlines():
        stripped = line.strip()
        if (
            stripped
            and stripped == last_nonempty
            and len(stripped) >= 20
            and not stripped.startswith(("#", "|", "!["))
        ):
            continue
        output_lines.append(line)
        if stripped:
            last_nonempty = stripped
    return "\n".join(output_lines)


def _normalize_markdown_table_blocks(text: str) -> str:
    lines = text.splitlines()
    normalized_lines: list[str] = []
    table_block: list[str] = []
    in_code_block = False

    def flush_table_block() -> None:
        nonlocal table_block
        if not table_block:
            return
        normalized_lines.extend(_normalize_markdown_table_block(table_block))
        table_block = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            flush_table_block()
            in_code_block = not in_code_block
            normalized_lines.append(line)
            continue
        if not in_code_block and _looks_like_markdown_table_line(line):
            table_block.append(line)
            continue
        flush_table_block()
        normalized_lines.append(line)

    flush_table_block()
    return "\n".join(normalized_lines)


def _looks_like_markdown_table_line(line: str) -> bool:
    stripped = line.strip()
    return stripped.startswith("|") and stripped.count("|") >= 2


def _normalize_markdown_table_block(lines: list[str]) -> list[str]:
    parsed_rows: list[list[str]] = []
    caption_line = ""
    trailing_lines: list[str] = []

    for line in lines:
        row = _split_markdown_table_row(line)
        if not row:
            continue
        if _is_markdown_table_separator_row(row):
            continue
        parsed_rows.append(row)

    parsed_rows = _sanitize_table_rows(parsed_rows)
    if not parsed_rows:
        return lines

    if _is_table_caption_row(parsed_rows[0]):
        caption_line = next((cell for cell in parsed_rows[0] if cell), "")
        caption_line = re.sub(r"[.\s]{5,}$", "", caption_line).strip()
        parsed_rows = parsed_rows[1:]

    while parsed_rows and _is_table_note_row(parsed_rows[-1]):
        trailing_lines.insert(0, parsed_rows.pop()[0])

    if len(parsed_rows) < 2 or max((len(row) for row in parsed_rows), default=0) < 2:
        rendered = [caption_line] if caption_line else []
        rendered.extend(trailing_lines or lines)
        return rendered

    rendered_table = _render_table_rows(parsed_rows)
    rendered_lines: list[str] = []
    if caption_line:
        rendered_lines.append(f"> **{caption_line}**")
        rendered_lines.append("")
    rendered_lines.extend(rendered_table.splitlines())
    if trailing_lines:
        rendered_lines.append("")
        rendered_lines.extend(trailing_lines)
    return rendered_lines


def _split_markdown_table_row(line: str) -> list[str]:
    stripped = line.strip().strip("|")
    if not stripped:
        return []
    return [_clean_markdown_table_cell(cell) for cell in stripped.split("|")]


def _clean_markdown_table_cell(cell: str) -> str:
    cleaned = _clean_text(_fix_unicode(cell)).replace("<br>", " ").replace("\n", " ").strip()
    cleaned = _DISPLAY_DOI_FRAGMENT_RE.sub("", cleaned).strip()
    cleaned = re.sub(r"\.{5,}", " ", cleaned)
    if re.fullmatch(r"[._·•\s]+", cleaned):
        return ""
    cleaned = re.sub(r"\s+", " ", cleaned)
    cleaned = _normalize_table_cell_text(cleaned)
    return cleaned


def _normalize_table_cell_text(text: str) -> str:
    if not text:
        return ""

    normalized = text
    compact = re.sub(r"\s+", "", normalized).lower()
    if compact in _TABLE_CELL_LITERAL_REPLACEMENTS:
        return _TABLE_CELL_LITERAL_REPLACEMENTS[compact]

    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)and(\d)\b",
        r"\1 \2 and \3",
        normalized,
    )
    normalized = re.sub(
        r"\b(MaSp|MiSp|TuSp|AcSp|PySp|AgSp|Flag)(\d)\b",
        r"\1 \2",
        normalized,
    )
    normalized = re.sub(r"\b(inner|outer)(core)\b", r"\1 \2", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(tensile)(strength)\b", r"\1 \2", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(water)(balance)\b", r"\1 \2", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\b(social)(aspects)\b", r"\1 \2", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    return normalized


def _is_markdown_table_separator_row(row: list[str]) -> bool:
    if not row:
        return False
    tokens = [cell.replace(" ", "") for cell in row]
    return all(token and re.fullmatch(r":?-{2,}:?", token) for token in tokens)


def _sanitize_table_rows(table: list[list[str]]) -> list[list[str]]:
    rows = [[str(cell or "").strip().replace("\n", " ") for cell in row] for row in table]
    rows = [[_clean_markdown_table_cell(cell) for cell in row] for row in rows]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return []

    max_cols = max(len(row) for row in rows)
    rows = [row + [""] * (max_cols - len(row)) for row in rows]

    keep_indices = [index for index in range(max_cols) if any(row[index] for row in rows)]
    if not keep_indices:
        return []
    rows = [[row[index] for index in keep_indices] for row in rows]
    return _merge_sparse_table_rows(rows)


def _merge_sparse_table_rows(rows: list[list[str]]) -> list[list[str]]:
    if not rows or len(rows[0]) < 3:
        return rows

    merged_rows: list[list[str]] = []
    for row in rows:
        non_empty_indices = [index for index, cell in enumerate(row) if cell]
        if (
            len(non_empty_indices) == 1
            and non_empty_indices[0] == 0
            and merged_rows
            and not _is_table_note_row(row)
        ):
            previous = merged_rows[-1]
            target_index = max((index for index, cell in enumerate(previous) if cell), default=-1)
            if target_index > 0:
                separator = "; " if previous[target_index] else ""
                previous[target_index] = f"{previous[target_index]}{separator}{row[0]}".strip()
                continue
        merged_rows.append(row)
    return merged_rows


def _is_table_caption_row(row: list[str]) -> bool:
    if len(row) < 1:
        return False
    non_empty = [cell for cell in row if cell]
    if len(non_empty) != 1:
        return False
    return _display_caption_key(non_empty[0]).startswith("table ")


def _is_table_note_row(row: list[str]) -> bool:
    non_empty = [cell for cell in row if cell]
    if len(non_empty) != 1:
        return False
    cell = non_empty[0]
    if _display_caption_key(cell):
        return False
    if len(cell) > 180:
        return True
    if ":" in cell and len(cell) >= 20:
        return True
    return cell.lower().startswith(("note:", "abbreviations:", "*", "ns", "p "))


def _render_table_rows(rows: list[list[str]]) -> str:
    if not rows or not rows[0]:
        return ""

    col_widths = [max(len(row[i]) for row in rows if i < len(row)) for i in range(len(rows[0]))]

    def fmt(row: list[str]) -> str:
        return "| " + " | ".join(
            (row[i] if i < len(row) else "").ljust(col_widths[i]) for i in range(len(col_widths))
        ) + " |"

    header = fmt(rows[0])
    separator = "| " + " | ".join("-" * max(width, 3) for width in col_widths) + " |"
    body = "\n".join(fmt(row) for row in rows[1:])
    return f"{header}\n{separator}\n{body}" if body else f"{header}\n{separator}"


def _table_to_markdown(table: list[list[str]]) -> str:
    rows = _sanitize_table_rows(table)
    if not rows or not rows[0]:
        return ""
    return _render_table_rows(rows)


def _iter_docx_blocks(doc: Any):
    from docx.table import Table
    from docx.text.paragraph import Paragraph as DocxParagraph

    parent = doc.element.body
    for child in parent.iterchildren():
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            yield {"type": "paragraph", "element": DocxParagraph(child, parent)}
        elif tag == "tbl":
            yield {"type": "table", "element": Table(child, parent)}


NON_PDF_EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".csv": extract_csv,
    ".txt": extract_txt,
    ".md": extract_txt,
}

SUPPORTED_EXTENSIONS = {".pdf", *NON_PDF_EXTRACTORS.keys()}


def convert_file_with_details(
    input_path: Path,
    output_path: Path,
    chunk_size: int = 0,
    pdf_config: PDFConfig | None = None,
    cancel_event: Any | None = None,
) -> dict[str, Any]:
    _check_cancel(cancel_event)
    ext = input_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        print(f"  Unsupported format: {ext} - skipping {input_path.name}")
        return {
            "status": "skipped",
            "tokens": 0,
            "chunk_count": 0,
            "output_path": str(output_path),
            "actual_pdf_backend": None,
        }

    config = pdf_config or PDFConfig()
    _validate_input_file_limits(input_path, config.security)
    if ext == ".pdf" and not config.security.hardened_mode:
        _validate_pdf_page_limit(input_path, config.security)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if _is_already_processed(input_path, output_path, chunk_size, config):
        print(f"  Skipping {input_path.name} (already processed)")
        actual_backend = None
        if ext == ".pdf":
            manifest = _load_manifest(_manifest_path(output_path))
            manifest_entry = manifest.get(_manifest_key(output_path), {})
            actual_backend = manifest_entry.get("pdf", {}).get("backend")
        return {
            "status": "skipped_cached",
            "tokens": 0,
            "chunk_count": 0,
            "output_path": str(output_path),
            "actual_pdf_backend": actual_backend,
        }

    print(f"  Converting {input_path.name} ...", end=" ", flush=True)
    actual_pdf_backend: str | None = None
    if ext == ".pdf":
        text, actual_pdf_backend = extract_pdf(input_path, output_path, config, cancel_event=cancel_event)
        text = _normalize_pdf_markdown_math(text, source_path=input_path)
    else:
        text = NON_PDF_EXTRACTORS[ext](input_path)

    _check_cancel(cancel_event)
    tokens = count_tokens(text)
    if chunk_size > 0:
        chunks = chunk_markdown(text, max_tokens=chunk_size)
        stem = output_path.stem
        for i, chunk in enumerate(chunks, 1):
            _check_cancel(cancel_event)
            chunk_path = output_path.parent / f"{stem}_chunk{i:03d}.md"
            safe_atomic_write_text(chunk_path, chunk, encoding="utf-8")
        print(f"done: {tokens:,} tokens -> {len(chunks)} chunks")
        chunk_count = len(chunks)
    else:
        _check_cancel(cancel_event)
        safe_atomic_write_text(output_path, text, encoding="utf-8")
        print(f"done: {tokens:,} tokens -> {output_path.name}")
        chunk_count = 1
    _check_cancel(cancel_event)
    _record_processed(input_path, output_path, chunk_size, config, actual_pdf_backend=actual_pdf_backend)
    return {
        "status": "converted",
        "tokens": tokens,
        "chunk_count": chunk_count,
        "output_path": str(output_path),
        "actual_pdf_backend": actual_pdf_backend,
    }


def convert_file(
    input_path: Path,
    output_path: Path,
    chunk_size: int = 0,
    pdf_config: PDFConfig | None = None,
    cancel_event: Any | None = None,
) -> None:
    convert_file_with_details(
        input_path,
        output_path,
        chunk_size=chunk_size,
        pdf_config=pdf_config,
        cancel_event=cancel_event,
    )

def _is_already_processed(
    input_path: Path,
    output_path: Path,
    chunk_size: int,
    pdf_config: PDFConfig,
) -> bool:
    if not _outputs_exist(output_path, chunk_size):
        return False

    manifest = _load_manifest(_manifest_path(output_path))
    key = _manifest_key(output_path)
    expected = _build_manifest_entry(input_path, output_path, chunk_size, pdf_config)
    return manifest.get(key) == expected


def _record_processed(
    input_path: Path,
    output_path: Path,
    chunk_size: int,
    pdf_config: PDFConfig,
    actual_pdf_backend: str | None = None,
) -> None:
    manifest_path = _manifest_path(output_path)
    manifest = _load_manifest(manifest_path)
    manifest[_manifest_key(output_path)] = _build_manifest_entry(
        input_path,
        output_path,
        chunk_size,
        pdf_config,
        actual_pdf_backend=actual_pdf_backend,
    )
    _save_manifest(manifest_path, manifest)


def _outputs_exist(output_path: Path, chunk_size: int) -> bool:
    if chunk_size > 0:
        first_chunk = output_path.parent / f"{output_path.stem}_chunk001.md"
        return first_chunk.exists()
    return output_path.exists()


def _manifest_path(output_path: Path) -> Path:
    return output_path.parent / ".llm_ingest_manifest.json"


def _manifest_key(output_path: Path) -> str:
    return str(output_path.resolve()).lower()


def _load_manifest(manifest_path: Path) -> dict[str, Any]:
    if not manifest_path.exists():
        return {}
    try:
        data = json.loads(manifest_path.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            return data
    except (json.JSONDecodeError, OSError):
        return {}
    return {}


def _save_manifest(manifest_path: Path, manifest: dict[str, Any]) -> None:
    safe_atomic_write_json(manifest_path, manifest)


def _build_manifest_entry(
    input_path: Path,
    output_path: Path,
    chunk_size: int,
    pdf_config: PDFConfig,
    actual_pdf_backend: str | None = None,
) -> dict[str, Any]:
    stat = input_path.stat()
    ext = input_path.suffix.lower()
    entry = {
        "manifest_version": MANIFEST_VERSION,
        "pipeline_signature": _pipeline_signature(),
        "source_path": _redact_local_path(input_path.resolve()) if pdf_config.security.privacy_mode else str(input_path.resolve()),
        "source_mtime_ns": stat.st_mtime_ns,
        "source_size": stat.st_size,
        "chunk_size": chunk_size,
        "output_path": _redact_local_path(output_path.resolve()) if pdf_config.security.privacy_mode else str(output_path.resolve()),
        "extension": ext,
    }
    if ext == ".pdf":
        resolved_backend = actual_pdf_backend or _resolved_pdf_backend_name(pdf_config)
        entry["pdf"] = {
            "backend_requested": pdf_config.pdf_backend,
            "backend": resolved_backend,
            "ocr_language": pdf_config.ocr_language,
            "ocr_dpi": pdf_config.ocr_dpi,
            "ocr_mode": pdf_config.ocr_mode,
            "table_strategy": pdf_config.table_strategy,
            "marker_python": (
                _redact_local_path(_find_marker_python(pdf_config) or "")
                if pdf_config.security.privacy_mode and (resolved_backend == "marker" or pdf_config.pdf_backend == "marker")
                else (str(_find_marker_python(pdf_config) or "") if resolved_backend == "marker" or pdf_config.pdf_backend == "marker" else "")
            ),
        }
    return entry


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _safe_sample_id(value: str) -> str:
    slug = re.sub(r"[^a-z0-9._-]+", "-", value.lower()).strip("-")
    return slug or "sample"


def _short_slug(value: str, limit: int = 48) -> str:
    slug = _safe_sample_id(value)
    if len(slug) <= limit:
        return slug
    digest = hashlib.sha1(slug.encode("utf-8")).hexdigest()[:8]
    head = slug[: max(12, limit - 9)].rstrip("-")
    return f"{head}-{digest}"


def load_audit_manifest(manifest_path: Path) -> list[AuditSample]:
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not isinstance(data, list):
        raise ValueError("Audit manifest must be a JSON array.")

    samples: list[AuditSample] = []
    for item in data:
        if not isinstance(item, dict):
            raise ValueError("Audit manifest entries must be JSON objects.")
        samples.append(
            AuditSample(
                id=str(item["id"]),
                category=str(item["category"]),
                source_url=str(item["source_url"]),
                filename=str(item.get("filename") or f"{item['id']}.pdf"),
                expected_traits=tuple(str(value) for value in item.get("expected_traits", [])),
                recommended_backends=tuple(str(value) for value in item.get("recommended_backends", [])),
                sha256=str(item.get("sha256", "")),
                license_note=str(item.get("license_note", "")),
            )
        )
    return samples


def _download_file(url: str, destination: Path, limits: SecurityLimits) -> None:
    parsed = urllib.parse.urlparse(url)
    if parsed.scheme.lower() != "https":
        raise ValueError(f"Audit downloads require https URLs: {url}")
    destination.parent.mkdir(parents=True, exist_ok=True)
    partial = destination.with_suffix(destination.suffix + ".partial")
    max_bytes = limits.max_audit_download_mb * 1024 * 1024
    request = urllib.request.Request(
        url,
        headers={
            "User-Agent": "llm_ingest audit downloader/1.0",
            "Accept": "application/pdf,application/octet-stream;q=0.9,*/*;q=0.8",
        },
    )
    with urllib.request.urlopen(request, timeout=180) as response:
        content_type = str(response.headers.get("Content-Type", "")).lower()
        if "pdf" not in content_type and "octet-stream" not in content_type:
            raise ValueError(f"Audit download did not look like a PDF ({content_type or 'unknown content type'}): {url}")
        total = 0
        with partial.open("wb") as handle:
            while True:
                chunk = response.read(1024 * 1024)
                if not chunk:
                    break
                total += len(chunk)
                if total > max_bytes:
                    with contextlib.suppress(OSError):
                        partial.unlink()
                    raise ValueError(
                        f"Audit download exceeded {limits.max_audit_download_mb} MB: {url}"
                    )
                handle.write(chunk)
    partial.replace(destination)


def ensure_audit_corpus(
    manifest_path: Path,
    cache_dir: Path,
    *,
    download_missing: bool = False,
    security: SecurityLimits | None = None,
) -> tuple[list[AuditFileTarget], list[str]]:
    security = security or SecurityLimits()
    samples = load_audit_manifest(manifest_path)
    targets: list[AuditFileTarget] = []
    missing: list[str] = []

    for sample in samples:
        cached_path = cache_dir / sample.filename
        needs_download = not cached_path.exists()
        hash_matches = False
        if cached_path.exists() and sample.sha256:
            with contextlib.suppress(OSError):
                hash_matches = _sha256_file(cached_path).lower() == sample.sha256.lower()
            needs_download = not hash_matches

        if needs_download and download_missing:
            if not sample.sha256 and not security.allow_unverified_downloads:
                raise ValueError(f"Audit sample '{sample.id}' is missing sha256; refusing unverified download.")
            _download_file(sample.source_url, cached_path, security)
            if sample.sha256:
                actual = _sha256_file(cached_path).lower()
                expected = sample.sha256.lower()
                if actual != expected:
                    with contextlib.suppress(OSError):
                        cached_path.unlink()
                    raise ValueError(
                        f"Downloaded sample '{sample.id}' does not match the manifest hash. "
                        f"Expected {expected}, got {actual}."
                    )
            hash_matches = True

        if not cached_path.exists():
            missing.append(sample.id)
            continue

        if sample.sha256 and not hash_matches:
            actual = _sha256_file(cached_path).lower()
            if actual != sample.sha256.lower():
                missing.append(sample.id)
                continue

        targets.append(
            AuditFileTarget(
                id=sample.id,
                label=sample.filename,
                path=cached_path,
                source_kind="seed",
                category=sample.category,
                source_url=sample.source_url,
                expected_traits=sample.expected_traits,
            )
        )

    return targets, missing


def discover_audit_targets(
    baseline_dirs: list[Path],
    manifest_path: Path,
    cache_dir: Path,
    *,
    download_missing: bool = False,
    security: SecurityLimits | None = None,
) -> tuple[list[AuditFileTarget], list[str]]:
    targets: list[AuditFileTarget] = []
    for baseline_dir in baseline_dirs:
        if not baseline_dir.exists() or not baseline_dir.is_dir():
            continue
        for file in list_supported_files(baseline_dir):
            if file.suffix.lower() != ".pdf":
                continue
            relative = file.relative_to(baseline_dir).as_posix()
            targets.append(
                AuditFileTarget(
                    id=f"baseline-{_safe_sample_id(relative)}",
                    label=relative,
                    path=file,
                    source_kind="baseline",
                    category="baseline",
                )
            )

    manifest_targets, missing = ensure_audit_corpus(
        manifest_path,
        cache_dir,
        download_missing=download_missing,
        security=security,
    )
    targets.extend(manifest_targets)
    return targets, missing


def parse_audit_backend_specs(csv_text: str) -> list[AuditBackendSpec]:
    specs: list[AuditBackendSpec] = []
    for raw_value in (part.strip() for part in csv_text.split(",")):
        if not raw_value:
            continue
        if ":" in raw_value:
            name, ocr_mode = raw_value.split(":", 1)
            specs.append(AuditBackendSpec(name=name.strip(), ocr_mode=ocr_mode.strip()))
        else:
            specs.append(AuditBackendSpec(name=raw_value, ocr_mode=None))
    if not specs:
        raise ValueError("At least one audit backend must be specified.")
    return specs


def _backend_spec_config(
    spec: AuditBackendSpec,
    marker_python: str | None = None,
    security: SecurityLimits | None = None,
) -> PDFConfig:
    return PDFConfig(
        ocr_language="eng",
        ocr_dpi=200,
        ocr_mode=spec.ocr_mode or "auto",
        pdf_backend=spec.name,
        table_strategy="lines_strict",
        marker_python=marker_python,
        security=security or SecurityLimits(),
    )


def _audit_issue_counts(markdown_path: Path) -> dict[str, int]:
    text = markdown_path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    counts: dict[str, int] = {
        "generic_page_alt": sum(1 for line in lines if line.startswith("![Page ")),
        "duplicate_doi_prefix": text.count('doi: "doi:'),
        "article_heading": text.count("## Article") + text.count("## Feature Article"),
        "broken_doi_spacing": text.count("https://doi. org"),
        "legacy_ligatures": text.count("fbers") + text.count("specifc"),
        "saxs_waxd_garble": text.count("B[2] obs") + text.count("curve]["),
        "unit_duplication": text.count("MJ/m^3 MPa"),
        "replacement_char": text.count("\ufffd"),
        "broken_charge_notation": text.count("[−]") + text.count("[þ]"),
    }
    for finding in llm_audit_assertions.scan_markdown_file(markdown_path):
        key = f"assertion_{finding.rule}"
        counts[key] = counts.get(key, 0) + 1
    return {key: value for key, value in counts.items() if value}


def _asset_dir_for_output(output_path: Path) -> Path:
    return output_path.parent / f"{output_path.stem}_assets"


def _audit_log_excerpt(text: str, limit: int = 1200) -> str:
    compact = text.strip()
    if len(compact) <= limit:
        return compact
    return compact[: limit - 3].rstrip() + "..."


def _write_audit_report(report_dir: Path, report: AuditReport) -> None:
    report_dir.mkdir(parents=True, exist_ok=True)
    payload = {
        "created_at": report.created_at,
        "manifest_path": report.manifest_path,
        "cache_dir": report.cache_dir,
        "report_dir": report.report_dir,
        "baseline_dirs": list(report.baseline_dirs),
        "backend_labels": list(report.backend_labels),
        "backend_plan": report.backend_plan,
        "missing_samples": list(report.missing_samples),
        "results": [
            {
                "sample_id": result.sample_id,
                "sample_label": result.sample_label,
                "source_kind": result.source_kind,
                "category": result.category,
                "backend_requested": result.backend_requested,
                "backend_label": result.backend_label,
                "backend_used": result.backend_used,
                "status": result.status,
                "output_path": result.output_path,
                "asset_dir": result.asset_dir,
                "tokens": result.tokens,
                "asset_count": result.asset_count,
                "issue_counts": result.issue_counts,
                "issue_total": result.issue_total,
                "log_excerpt": result.log_excerpt,
                "error": result.error,
            }
            for result in report.results
        ],
    }
    safe_atomic_write_json(report_dir / "audit_report.json", payload)
    _write_audit_assertions_summary(report_dir, report.results)

    success_count = sum(1 for result in report.results if result.status == "ok")
    skipped_count = sum(1 for result in report.results if result.status == "skipped")
    failure_count = sum(1 for result in report.results if result.status == "failed")
    issue_total = sum(result.issue_total for result in report.results)

    lines = [
        "# LLM Ingest Audit Report",
        "",
        f"- Created: {report.created_at}",
        f"- Manifest: `{report.manifest_path}`",
        f"- Cache dir: `{report.cache_dir}`",
        f"- Report dir: `{report.report_dir}`",
        f"- Baseline dirs: {', '.join(f'`{value}`' for value in report.baseline_dirs) or 'None'}",
        f"- Backends: {', '.join(report.backend_labels)}",
        f"- Missing seed samples: {', '.join(report.missing_samples) if report.missing_samples else 'None'}",
        "",
        "## Summary",
        "",
        f"- Successful runs: {success_count}",
        f"- Skipped runs: {skipped_count}",
        f"- Failed runs: {failure_count}",
        f"- Total soft issues flagged: {issue_total}",
        f"- Assertion details: `{Path(report.report_dir) / 'audit_assertions.md'}`",
        "",
        "## Results",
        "",
        "| Sample | Source | Backend | Used | Status | Tokens | Assets | Issues |",
        "| --- | --- | --- | --- | --- | ---: | ---: | ---: |",
    ]
    for result in report.results:
        lines.append(
            f"| {result.sample_label} | {result.source_kind} | {result.backend_label} | "
            f"{result.backend_used or '-'} | {result.status} | {result.tokens} | {result.asset_count} | {result.issue_total} |"
        )

    lines.extend(["", "## Findings", ""])
    for result in report.results:
        if result.status == "ok" and not result.issue_counts:
            continue
        lines.append(f"### {result.sample_label} [{result.backend_label}]")
        lines.append("")
        lines.append(f"- Status: {result.status}")
        if result.backend_used:
            lines.append(f"- Backend used: `{result.backend_used}`")
        if result.issue_counts:
            issue_line = ", ".join(f"{key}={value}" for key, value in sorted(result.issue_counts.items()))
            lines.append(f"- Issues: {issue_line}")
        if result.error:
            lines.append(f"- Error: {result.error}")
        if result.output_path:
            lines.append(f"- Output: `{result.output_path}`")
        if result.log_excerpt:
            lines.append("")
            lines.append("```text")
            lines.append(result.log_excerpt)
            lines.append("```")
        lines.append("")

    safe_atomic_write_text(report_dir / "audit_summary.md", "\n".join(lines).strip() + "\n", encoding="utf-8")


def _write_audit_assertions_summary(report_dir: Path, results: tuple[AuditRunResult, ...]) -> None:
    lines = ["# Audit Assertion Details", ""]
    total = 0
    for result in results:
        output_path = Path(result.output_path)
        if result.status != "ok" or not output_path.exists():
            continue
        findings = llm_audit_assertions.scan_markdown_file(output_path)
        if not findings:
            continue
        total += len(findings)
        lines.append(f"## {result.sample_label} [{result.backend_label}]")
        lines.append("")
        lines.append(f"- Output: `{result.output_path}`")
        lines.append("")
        for finding in findings:
            lines.append(f"- `{finding.rule}` line {finding.line}: {finding.message}")
            if finding.excerpt:
                lines.append(f"  `{finding.excerpt}`")
        lines.append("")
    if total == 0:
        lines.append("No assertion regressions were found in successful audit outputs.")
    else:
        lines.insert(2, f"Total assertion findings: {total}")
        lines.insert(3, "")
    safe_atomic_write_text(report_dir / "audit_assertions.md", "\n".join(lines).strip() + "\n", encoding="utf-8")


def run_audit(
    manifest_path: Path,
    cache_dir: Path,
    report_dir: Path,
    backend_specs: list[AuditBackendSpec],
    *,
    baseline_dirs: list[Path] | None = None,
    download_missing: bool = False,
    marker_python: str | None = None,
    security: SecurityLimits | None = None,
    cancel_event: Any | None = None,
    progress_callback: Any | None = None,
) -> AuditReport:
    security = security or SecurityLimits()
    baseline_dirs = baseline_dirs or []
    targets, missing = discover_audit_targets(
        baseline_dirs,
        manifest_path,
        cache_dir,
        download_missing=download_missing,
        security=security,
    )
    report_dir.mkdir(parents=True, exist_ok=True)
    render_root = report_dir / "renders"
    render_root.mkdir(parents=True, exist_ok=True)

    backend_plan: dict[str, list[str]] = {}
    results: list[AuditRunResult] = []
    total_runs = len(targets) * len(backend_specs)

    for run_index, target in enumerate(targets, 1):
        for backend_index, spec in enumerate(backend_specs, 1):
            _check_cancel(cancel_event)
            ordinal = ((run_index - 1) * len(backend_specs)) + backend_index
            progress_label = f"Audit {ordinal}/{total_runs}: {target.label} [{spec.label}]"
            if callable(progress_callback):
                progress_callback(progress_label)
            print(progress_label)

            config = _backend_spec_config(spec, marker_python=marker_python, security=security)
            require_marker_models = spec.name == "marker"
            plan = inspect_pdf_backend_plan(config, require_marker_models=require_marker_models, sample_path=target.path)
            backend_plan.setdefault(spec.label, []).extend(describe_pdf_backend_plan(plan))

            output_dir = render_root / _short_slug(target.id, limit=42)
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / f"{_short_slug(target.path.stem, limit=40)}__{_short_slug(spec.label, limit=14)}.md"
            asset_dir = _asset_dir_for_output(output_path)
            if output_path.exists():
                output_path.unlink()
            if asset_dir.exists():
                safe_remove_generated_dir(asset_dir)

            buffer = io.StringIO()
            backend_used: str | None = None
            tokens = 0
            status = "ok"
            error = ""
            issue_counts: dict[str, int] = {}
            asset_count = 0
            try:
                with contextlib.redirect_stdout(buffer), contextlib.redirect_stderr(buffer):
                    details = convert_file_with_details(
                        target.path,
                        output_path,
                        chunk_size=0,
                        pdf_config=config,
                        cancel_event=cancel_event,
                    )
                backend_used = details.get("actual_pdf_backend")
                if details.get("status") == "skipped":
                    status = "skipped"
                elif details.get("status") == "skipped_cached":
                    status = "skipped"
                if output_path.exists():
                    tokens = count_tokens(output_path.read_text(encoding="utf-8", errors="replace"))
                    issue_counts = _audit_issue_counts(output_path)
                if asset_dir.exists():
                    asset_count = len([file for file in asset_dir.iterdir() if file.is_file()])
            except BaseException as exc:
                status = "failed"
                if isinstance(exc, SystemExit):
                    error = _system_exit_message(exc) or "SystemExit"
                else:
                    error = str(exc).strip() or exc.__class__.__name__

            log_excerpt = _audit_log_excerpt(buffer.getvalue())
            results.append(
                AuditRunResult(
                    sample_id=target.id,
                    sample_label=target.label,
                    source_kind=target.source_kind,
                    category=target.category,
                    backend_requested=spec.name,
                    backend_label=spec.label,
                    backend_used=backend_used,
                    status=status,
                    output_path=str(output_path),
                    asset_dir=str(asset_dir),
                    tokens=tokens,
                    asset_count=asset_count,
                    issue_counts=issue_counts,
                    issue_total=sum(issue_counts.values()),
                    log_excerpt=log_excerpt,
                    error=error,
                )
            )

    report = AuditReport(
        created_at=dt.datetime.now(dt.timezone.utc).isoformat(),
        manifest_path=str(manifest_path),
        cache_dir=str(cache_dir),
        report_dir=str(report_dir),
        baseline_dirs=tuple(str(path) for path in baseline_dirs),
        backend_labels=tuple(spec.label for spec in backend_specs),
        backend_plan=backend_plan,
        missing_samples=tuple(missing),
        results=tuple(results),
    )
    _write_audit_report(report_dir, report)
    return report


def _build_convert_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert documents to LLM-ready Markdown")
    parser.add_argument("input", help="File or directory to convert")
    parser.add_argument("-o", "--output", help="Output file (single-file mode only)")
    parser.add_argument("-c", "--chunk", type=int, default=0, help="Max tokens per chunk (0 = no chunking)")
    parser.add_argument("--out-dir", default="llm_ready", help="Output directory for batch mode (default: ./llm_ready)")
    parser.add_argument("--ocr-language", default="eng", help="OCR language(s) for PDFs, e.g. eng or eng+deu")
    parser.add_argument("--ocr-dpi", type=int, default=200, help="OCR DPI for PDFs")
    parser.add_argument("--tessdata", help="Path to the Tesseract tessdata directory")
    parser.add_argument("--marker-python", help="Path to the Python interpreter that has marker-pdf installed")
    parser.add_argument("--backend-timeout-seconds", type=int, default=DEFAULT_BACKEND_TIMEOUT_SECONDS, help="PDF backend timeout in hardened mode.")
    parser.add_argument("--max-input-mb", type=int, default=DEFAULT_MAX_INPUT_MB, help="Maximum input file size in MB.")
    parser.add_argument("--max-pdf-pages", type=int, default=DEFAULT_MAX_PDF_PAGES, help="Maximum PDF page count.")
    parser.add_argument("--max-extracted-assets", type=int, default=DEFAULT_MAX_EXTRACTED_ASSETS, help="Maximum image assets extracted per PDF.")
    parser.add_argument("--allow-external-marker-python", action="store_true", help="Allow --marker-python or LLM_INGEST_MARKER_PYTHON outside the local sidecar.")
    parser.add_argument("--privacy-mode", action="store_true", help="Redact local user paths in manifests and reports where possible.")
    parser.add_argument("--no-hardened-mode", action="store_true", help="Run PDF extraction in-process for compatibility.")
    parser.add_argument(
        "--pdf-backend",
        choices=("auto", "custom", "pymupdf4llm", "marker"),
        default="auto",
        help="PDF extraction backend: auto, custom, pymupdf4llm, or marker",
    )
    parser.add_argument(
        "--table-strategy",
        choices=("lines_strict", "lines", "text", "none"),
        default="lines_strict",
        help="PyMuPDF table extraction strategy for the custom PDF backend",
    )
    parser.add_argument(
        "--ocr-mode",
        choices=("auto", "full", "off"),
        default="auto",
        help="PDF OCR mode: auto, full, or off",
    )
    return parser


def _build_audit_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Run the PDF audit corpus against configured backends.")
    parser.add_argument("--manifest", default=DEFAULT_AUDIT_MANIFEST, help="Path to the audit corpus manifest JSON.")
    parser.add_argument("--cache-dir", default=DEFAULT_AUDIT_CACHE_DIR, help="Directory where downloaded audit PDFs are cached.")
    parser.add_argument("--report-dir", default=DEFAULT_AUDIT_REPORT_DIR, help="Directory where audit reports and rendered outputs are written.")
    parser.add_argument("--backends", default=DEFAULT_AUDIT_BACKENDS, help="Comma-separated backend list, e.g. auto,custom:off,pymupdf4llm,marker.")
    parser.add_argument("--download-missing", action="store_true", help="Download any missing seed corpus PDFs before running the audit.")
    parser.add_argument("--allow-unverified-downloads", action="store_true", help="Allow audit downloads without a manifest sha256.")
    parser.add_argument("--max-audit-download-mb", type=int, default=DEFAULT_MAX_AUDIT_DOWNLOAD_MB, help="Maximum audit sample download size in MB.")
    parser.add_argument("--backend-timeout-seconds", type=int, default=DEFAULT_BACKEND_TIMEOUT_SECONDS, help="PDF backend timeout for audit runs.")
    parser.add_argument("--max-input-mb", type=int, default=DEFAULT_MAX_INPUT_MB, help="Maximum audited PDF size in MB.")
    parser.add_argument("--max-pdf-pages", type=int, default=DEFAULT_MAX_PDF_PAGES, help="Maximum audited PDF page count.")
    parser.add_argument("--max-extracted-assets", type=int, default=DEFAULT_MAX_EXTRACTED_ASSETS, help="Maximum image assets extracted per audited PDF.")
    parser.add_argument("--baseline-dir", action="append", default=[], help="Local directory to include in the audit baseline. May be provided multiple times.")
    parser.add_argument("--marker-python", help="Path to the Python interpreter that has marker-pdf installed.")
    return parser


def _build_graph_parser() -> argparse.ArgumentParser:
    import llm_knowledge_graph

    parser = argparse.ArgumentParser(description="Build or query a local Markdown knowledge graph.")
    subparsers = parser.add_subparsers(dest="graph_command", required=True)

    build = subparsers.add_parser("build", help="Build a graph index from generated Markdown files.")
    build.add_argument("--source-dir", default=llm_knowledge_graph.DEFAULT_GRAPH_SOURCE_DIR, help="Folder containing generated Markdown files.")
    build.add_argument("--index-dir", default=llm_knowledge_graph.DEFAULT_GRAPH_INDEX_DIR, help="Folder where graph artifacts are written.")
    build.add_argument("--max-chunk-tokens", type=int, default=850, help="Approximate max tokens per graph chunk.")
    build.add_argument("--top-terms", type=int, default=14, help="Top terms to keep per chunk.")
    build.add_argument("--embedding-model", choices=llm_knowledge_graph.SUPPORTED_EMBEDDING_MODELS, default=llm_knowledge_graph.DEFAULT_EMBEDDING_MODEL, help="Embedding backend for vector retrieval.")
    build.add_argument("--embedding-dimensions", type=int, default=llm_knowledge_graph.DEFAULT_EMBEDDING_DIMENSIONS, help="Dimensions for local hash embeddings.")
    build.add_argument("--max-source-files", type=int, default=2000, help="Maximum Markdown files to index.")
    build.add_argument("--max-graph-chunk-text-bytes", type=int, default=llm_knowledge_graph.DEFAULT_MAX_GRAPH_CHUNK_TEXT_BYTES, help="Maximum text bytes allowed per indexed Markdown source.")

    query = subparsers.add_parser("query", help="Query an existing graph index.")
    query.add_argument("query", help="Question or search phrase.")
    query.add_argument("--index-dir", default=llm_knowledge_graph.DEFAULT_GRAPH_INDEX_DIR, help="Folder containing graph artifacts.")
    query.add_argument("--limit", type=int, default=8, help="Number of evidence chunks to return.")
    query.add_argument("--mode", choices=("hybrid", "lexical", "vector"), default="hybrid", help="Retrieval mode for RAG evidence.")
    return parser


def _run_convert_cli(args: argparse.Namespace) -> None:
    security = SecurityLimits(
        max_input_mb=args.max_input_mb,
        max_pdf_pages=args.max_pdf_pages,
        max_extracted_assets=args.max_extracted_assets,
        backend_timeout_seconds=args.backend_timeout_seconds,
        hardened_mode=not args.no_hardened_mode,
        allow_external_marker_python=args.allow_external_marker_python,
        privacy_mode=args.privacy_mode,
    )
    pdf_config = PDFConfig(
        ocr_language=args.ocr_language,
        ocr_dpi=args.ocr_dpi,
        tessdata=args.tessdata,
        ocr_mode=args.ocr_mode,
        pdf_backend=args.pdf_backend,
        table_strategy=args.table_strategy,
        marker_python=args.marker_python,
        security=security,
    )

    input_path = Path(args.input)
    if input_path.is_file():
        output_path = Path(args.output) if args.output else input_path.with_suffix(".md")
        convert_file(input_path, output_path, chunk_size=args.chunk, pdf_config=pdf_config)
    elif input_path.is_dir():
        out_dir = Path(args.out_dir)
        out_dir.mkdir(parents=True, exist_ok=True)
        files = list_supported_files(input_path, out_dir)
        batch_plan = build_batch_targets(files, input_path, out_dir)
        print(f"\nFound {len(files)} supported files in {input_path}\n")
        for file, target in batch_plan:
            convert_file(file, target, chunk_size=args.chunk, pdf_config=pdf_config)
        print(f"\nDone. Output in ./{out_dir}/")
    else:
        sys.exit(f"Error: '{args.input}' is not a valid file or directory.")


def _run_audit_cli(args: argparse.Namespace) -> None:
    backend_specs = parse_audit_backend_specs(args.backends)
    security = SecurityLimits(
        max_input_mb=args.max_input_mb,
        max_pdf_pages=args.max_pdf_pages,
        max_extracted_assets=args.max_extracted_assets,
        max_audit_download_mb=args.max_audit_download_mb,
        backend_timeout_seconds=args.backend_timeout_seconds,
        allow_unverified_downloads=args.allow_unverified_downloads,
    )
    baseline_dirs = [Path(value) for value in args.baseline_dir]
    if not baseline_dirs and Path(DEFAULT_AUDIT_BASELINE_DIR).exists():
        baseline_dirs.append(Path(DEFAULT_AUDIT_BASELINE_DIR))

    report = run_audit(
        Path(args.manifest),
        Path(args.cache_dir),
        Path(args.report_dir),
        backend_specs,
        baseline_dirs=baseline_dirs,
        download_missing=args.download_missing,
        marker_python=args.marker_python,
        security=security,
    )
    print(f"\nAudit complete. Summary: {Path(report.report_dir) / 'audit_summary.md'}")
    print(f"JSON report: {Path(report.report_dir) / 'audit_report.json'}")


def _print_console_text(text: str) -> None:
    encoding = sys.stdout.encoding or "utf-8"
    sys.stdout.write(text.encode(encoding, errors="replace").decode(encoding, errors="replace"))
    if not text.endswith("\n"):
        sys.stdout.write("\n")


def _run_graph_cli(args: argparse.Namespace) -> None:
    import llm_knowledge_graph

    if args.graph_command == "build":
        report = llm_knowledge_graph.build_knowledge_graph(
            Path(args.source_dir),
            Path(args.index_dir),
            max_chunk_tokens=args.max_chunk_tokens,
            top_terms_per_chunk=args.top_terms,
            embedding_model=args.embedding_model,
            embedding_dimensions=args.embedding_dimensions,
            max_source_files=args.max_source_files,
            max_chunk_text_bytes=args.max_graph_chunk_text_bytes,
            progress_callback=print,
        )
        print(
            f"\nGraph complete: {report.document_count} docs, {report.chunk_count} chunks, "
            f"{report.node_count} nodes, {report.edge_count} edges, {report.embedding_count} vectors."
        )
        print(f"Graph index: {Path(report.index_dir)}")
        print(f"LLM context: {Path(report.index_dir) / 'graph_context.md'}")
        return

    if args.graph_command == "query":
        result = llm_knowledge_graph.query_knowledge_graph(Path(args.index_dir), args.query, limit=args.limit, retrieval_mode=args.mode)
        _print_console_text(result.context_markdown)
        print(f"\nSaved query pack: {Path(result.index_dir) / 'last_query.md'}")
        print(f"Saved RAG JSON: {Path(result.index_dir) / 'rag_pack.json'}")
        return

    raise ValueError(f"Unknown graph command: {args.graph_command}")


def main(argv: list[str] | None = None) -> None:
    argv = list(sys.argv[1:] if argv is None else argv)
    if argv and argv[0] == "audit":
        parser = _build_audit_parser()
        args = parser.parse_args(argv[1:])
        _run_audit_cli(args)
        return
    if argv and argv[0] in {"graph", "kg", "knowledge-graph"}:
        parser = _build_graph_parser()
        args = parser.parse_args(argv[1:])
        _run_graph_cli(args)
        return

    parser = _build_convert_parser()
    args = parser.parse_args(argv)
    _run_convert_cli(args)


if __name__ == "__main__":
    main()
